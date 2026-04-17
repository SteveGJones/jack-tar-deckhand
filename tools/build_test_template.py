#!/usr/bin/env python3
"""Generate a Metamirror-branded test template .pptx for template-driven layout tests.

Creates a minimal template with 1 slide master and layouts from python-pptx defaults.
Includes 2 example slides to verify stripping logic.

Usage:
    python3 tools/build_test_template.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), '..', 'tests', 'fixtures', 'templates', 'metamirror-template.pptx')

# Metamirror brand colours
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)    # dark navy
SECONDARY = RGBColor(0x00, 0xD4, 0xAA)  # teal accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def build_template():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add 2 example slides that the assembler must strip
    title_layout = prs.slide_layouts[0]  # Title Slide
    blank_layout = prs.slide_layouts[6]  # Blank

    slide1 = prs.slides.add_slide(title_layout)
    slide1.placeholders[0].text = 'Example Title Slide — DELETE ME'
    slide1.placeholders[1].text = 'This is an example subtitle'

    slide2 = prs.slides.add_slide(blank_layout)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f'Test template saved to {OUTPUT_PATH}')
    print(f'  Slide masters: {len(prs.slide_masters)}')
    print(f'  Layouts: {len(prs.slide_layouts)}')
    print(f'  Example slides: {len(prs.slides)}')

    # Print layout inventory
    for i, layout in enumerate(prs.slide_layouts):
        ph_types = [str(ph.placeholder_format.type).split('(')[0].strip()
                    for ph in layout.placeholders]
        print(f'  Layout {i}: "{layout.name}" — placeholders: {ph_types}')

    return OUTPUT_PATH


if __name__ == '__main__':
    build_template()
