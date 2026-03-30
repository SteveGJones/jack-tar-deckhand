"""Tests for keynote-specific QA checks."""

import io
import os
import shutil
import tempfile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def tmp_dir():
    d = tempfile.mkdtemp()
    yield d
    shutil.rmtree(d)


def _make_pptx_with_image(tmp_dir, image_color, filename='test.pptx'):
    """Create a minimal pptx with a single full-bleed image of the given RGB colour."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Create image in memory
    img = Image.new('RGB', (1920, 1080), color=image_color)
    img_path = os.path.join(tmp_dir, 'test_img.png')
    img.save(img_path)
    slide.shapes.add_picture(img_path, 0, 0, Inches(13.333), Inches(7.5))

    pptx_path = os.path.join(tmp_dir, filename)
    prs.save(pptx_path)
    return pptx_path


def test_palette_drift_within_tolerance(tmp_dir):
    from src.qa.checks.keynote_checks import check_palette_drift
    # Image is exact brand teal — no drift
    pptx_path = _make_pptx_with_image(tmp_dir, (0, 107, 94))  # #006B5E
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    brand_palette = ['006B5E', '5CDBC0', '4B635B', 'F5FBF7']
    issues = check_palette_drift(slide, 1, brand_palette=brand_palette)
    assert len(issues) == 0


def test_palette_drift_detected(tmp_dir):
    from src.qa.checks.keynote_checks import check_palette_drift
    # Image is bright red — completely off-brand
    pptx_path = _make_pptx_with_image(tmp_dir, (255, 0, 0))
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    brand_palette = ['006B5E', '5CDBC0', '4B635B', 'F5FBF7']
    issues = check_palette_drift(slide, 1, brand_palette=brand_palette)
    assert len(issues) >= 1
    assert issues[0]['category'] == 'palette_drift'
    assert issues[0]['severity'] == 'warning'


def test_palette_drift_no_images_no_issues(tmp_dir):
    from src.qa.checks.keynote_checks import check_palette_drift
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand_palette = ['006B5E']
    issues = check_palette_drift(slide, 1, brand_palette=brand_palette)
    assert len(issues) == 0


def test_run_qa_skips_text_checks_on_full_render(tmp_dir):
    """full_render slides should not get wall-of-text or font-size errors."""
    import json
    from src.qa.run_qa import run_qa

    # Create a deck dir with a strategy map marking slide 1 as full_render
    deck_dir = os.path.join(tmp_dir, 'deck')
    os.makedirs(deck_dir, exist_ok=True)
    strategy_map = {
        'approval_mode': 'review',
        'slides': [{'slide_number': 1, 'strategy': 'full_render', 'rationale': 'test', 'render_funnel': ['ollama']}],
    }
    with open(os.path.join(deck_dir, 'strategy-map.json'), 'w') as f:
        json.dump(strategy_map, f)

    # Create a minimal pptx with one image-only slide
    pptx_path = _make_pptx_with_image(tmp_dir, (0, 107, 94), 'qa_test.pptx')
    report = run_qa(pptx_path, deck_dir)

    # Should not have wall-of-text or font-size errors for slide 1
    text_errors = [f for f in report['findings']
                   if f['slide_number'] == 1 and f['category'] in ('text_overflow', 'consistency')]
    assert len(text_errors) == 0


def test_run_qa_applies_palette_drift_on_full_render(tmp_dir):
    """full_render slides should get palette drift checks."""
    import json
    from src.qa.run_qa import run_qa

    deck_dir = os.path.join(tmp_dir, 'deck')
    os.makedirs(deck_dir, exist_ok=True)
    strategy_map = {
        'approval_mode': 'review',
        'slides': [{'slide_number': 1, 'strategy': 'full_render', 'rationale': 'test', 'render_funnel': ['ollama']}],
    }
    with open(os.path.join(deck_dir, 'strategy-map.json'), 'w') as f:
        json.dump(strategy_map, f)

    # Brand profile with palette
    brand_profile = {'palette': {'primary': '006B5E', 'secondary': '4B635B'}}
    with open(os.path.join(deck_dir, 'brand-profile.json'), 'w') as f:
        json.dump(brand_profile, f)

    # Create pptx with an off-brand red image
    pptx_path = _make_pptx_with_image(tmp_dir, (255, 0, 0), 'drift_test.pptx')
    report = run_qa(pptx_path, deck_dir)

    drift_findings = [f for f in report['findings'] if f['category'] == 'palette_drift']
    assert len(drift_findings) >= 1
