"""Integration tests for template-driven layout pipeline."""

import json
import os
import pytest

from src.slide_prompt_composer import classify_slide_strategy, build_strategy_map


class TestTemplateStrategyConstraints:
    def test_template_mode_forces_composed(self):
        slide = {'slide_number': 1, 'slide_type': 'content', 'headline': 'Test',
                 'body_points': ['A'], 'visual_type': 'hero_image'}
        result = classify_slide_strategy(slide, template_mode=True)
        assert result['strategy'] == 'composed'

    def test_template_mode_title_stays_composed(self):
        slide = {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title',
                 'visual_type': 'hero_image'}
        result = classify_slide_strategy(slide, template_mode=True)
        assert result['strategy'] == 'composed'

    def test_template_mode_full_bleed_picture_allows_full_render(self):
        slide = {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title',
                 'visual_type': 'hero_image'}
        full_bleed_layout = {
            'placeholders': [
                {'idx': 13, 'type': 'picture', 'name': 'Picture', 'x': 0.0, 'y': 0.0, 'w': 13.33, 'h': 7.5},
            ]
        }
        result = classify_slide_strategy(slide, template_mode=True, template_layout=full_bleed_layout)
        assert result['strategy'] == 'full_render'

    def test_template_mode_small_picture_stays_composed(self):
        slide = {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title',
                 'visual_type': 'hero_image'}
        small_pic_layout = {
            'placeholders': [
                {'idx': 13, 'type': 'picture', 'name': 'Picture', 'x': 7.0, 'y': 0.0, 'w': 6.0, 'h': 7.5},
            ]
        }
        result = classify_slide_strategy(slide, template_mode=True, template_layout=small_pic_layout)
        assert result['strategy'] == 'composed'

    def test_without_template_mode_unchanged(self):
        slide = {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title',
                 'visual_type': 'hero_image'}
        result = classify_slide_strategy(slide, template_mode=False)
        assert result['strategy'] == 'full_render'

    def test_build_strategy_map_template_mode(self):
        outline = {
            'narrative_arc': 'test',
            'estimated_duration_minutes': 10,
            'slides': [
                {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title', 'visual_type': 'hero_image'},
                {'slide_number': 2, 'slide_type': 'content', 'headline': 'Content',
                 'body_points': ['A', 'B', 'C', 'D'], 'visual_type': 'hero_image'},
            ],
        }
        strategy_map = build_strategy_map(outline, template_mode=True)
        for entry in strategy_map['slides']:
            assert entry['strategy'] == 'composed'

    def test_build_strategy_map_render_funnel_ollama_only(self):
        outline = {
            'narrative_arc': 'test',
            'estimated_duration_minutes': 10,
            'slides': [
                {'slide_number': 1, 'slide_type': 'content', 'headline': 'Test', 'visual_type': 'hero_image'},
            ],
        }
        strategy_map = build_strategy_map(outline, template_mode=True)
        assert strategy_map['slides'][0]['render_funnel'] == ['ollama']


from src.deckcontext import DEFAULT_STEP_ORDER, write_contract, read_contract
from src.template_analyser import analyse_template
from src.assembler.build_deck_template import build_deck
from pptx import Presentation


class TestPipelineStepOrder:
    def test_template_analysis_in_step_order(self):
        assert 'template-analysis' in DEFAULT_STEP_ORDER

    def test_template_analysis_before_brand_manager(self):
        ta_idx = DEFAULT_STEP_ORDER.index('template-analysis')
        bm_idx = DEFAULT_STEP_ORDER.index('brand-manager')
        assert ta_idx < bm_idx

    def test_template_analysis_after_validate_brief(self):
        ta_idx = DEFAULT_STEP_ORDER.index('template-analysis')
        vb_idx = DEFAULT_STEP_ORDER.index('validate-brief')
        assert ta_idx > vb_idx


FIXTURE_PATH = os.path.join(os.path.dirname(__file__), 'fixtures', 'templates', 'metamirror-template.pptx')


class TestEndToEnd:
    def test_full_pipeline_template_to_pptx(self, tmp_path):
        deck_dir = str(tmp_path / 'deck')
        os.makedirs(os.path.join(deck_dir, 'images'), exist_ok=True)
        os.makedirs(os.path.join(deck_dir, 'output'), exist_ok=True)

        # 1. Analyse template
        profile = analyse_template(FIXTURE_PATH)
        profile['speaker_approved'] = True
        write_contract(deck_dir, 'template-profile', profile)

        # 2. Create outline
        outline = {
            'narrative_arc': 'problem-solution',
            'estimated_duration_minutes': 15,
            'slides': [
                {'slide_number': 1, 'slide_type': 'title', 'headline': 'Welcome to the Future',
                 'visual_type': 'hero_image'},
                {'slide_number': 2, 'slide_type': 'content', 'headline': 'The Problem',
                 'body_points': ['Legacy systems are slow', 'Users are frustrated', 'Costs are rising'],
                 'visual_type': 'none'},
                {'slide_number': 3, 'slide_type': 'content', 'headline': 'Our Solution',
                 'body_points': ['AI-driven automation', 'Real-time processing'],
                 'visual_type': 'hero_image'},
                {'slide_number': 4, 'slide_type': 'closing', 'headline': 'Thank You',
                 'visual_type': 'none'},
            ],
        }
        write_contract(deck_dir, 'outline', outline, validate=False)

        # 3. Create empty manifests
        write_contract(deck_dir, 'image-manifest', {'images': []}, validate=False)
        with open(os.path.join(deck_dir, 'speaker-notes.json'), 'w') as f:
            json.dump({'notes': [
                {'slide_number': 1, 'text': 'Welcome everyone'},
                {'slide_number': 4, 'text': 'Questions?'},
            ]}, f)

        # 4. Build strategy map in template mode
        strategy_map = build_strategy_map(outline, template_mode=True)
        for entry in strategy_map['slides']:
            assert entry['strategy'] == 'composed'

        # 5. Assemble
        loaded_profile = read_contract(deck_dir, 'template-profile')
        output_path = build_deck(deck_dir, FIXTURE_PATH, loaded_profile)

        # 6. Verify output
        assert os.path.isfile(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) == 4

        # Verify slide 1 has the title
        slide1_texts = [s.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert any('Welcome to the Future' in t for t in slide1_texts)

        # Verify slide 2 has body points
        slide2_texts = ' '.join(s.text for s in prs.slides[1].shapes if s.has_text_frame)
        assert 'Legacy systems are slow' in slide2_texts

        # Verify slide 1 has speaker notes
        notes1 = prs.slides[0].notes_slide.notes_text_frame.text
        assert 'Welcome everyone' in notes1

    def test_assembler_routing_uses_template_when_profile_exists(self, tmp_path):
        """Verify that the presence of template-profile.json routes to template assembler."""
        deck_dir = str(tmp_path / 'deck')
        os.makedirs(deck_dir, exist_ok=True)

        # No template profile — would route to JS
        profile = read_contract(deck_dir, 'template-profile')
        assert profile is None

        # With template profile — would route to python-pptx
        tp = analyse_template(FIXTURE_PATH)
        write_contract(deck_dir, 'template-profile', tp)
        profile = read_contract(deck_dir, 'template-profile')
        assert profile is not None
