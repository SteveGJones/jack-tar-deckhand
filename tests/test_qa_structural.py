"""Tests for structural QA checks using python-pptx to create test presentations."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from src.qa.checks.structural import (
    check_wall_of_text,
    check_font_size,
    check_orphan_widow,
    check_safe_margins,
    check_speaker_notes,
    check_slide_count_ratio,
    check_placeholder_residue,
    check_bullet_count,
    check_title_slide,
    check_closing_slide,
    check_text_overflow,
    check_dead_slides,
    check_element_overlap,
)


def make_slide_with_text(text, font_size=None):
    """Create a minimal presentation with one slide containing text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size is not None:
        run.font.size = Pt(font_size)
    return prs, slide


class TestWallOfText:
    def test_passes_short_text(self):
        _, slide = make_slide_with_text("Hello world")
        issues = check_wall_of_text(slide, 1)
        assert len(issues) == 0

    def test_fails_long_textbox(self):
        long_text = " ".join(["word"] * 50)
        _, slide = make_slide_with_text(long_text)
        issues = check_wall_of_text(slide, 1)
        assert any(i['category'] == 'text_overflow' for i in issues)

    def test_custom_threshold(self):
        text = " ".join(["word"] * 10)
        _, slide = make_slide_with_text(text)
        issues = check_wall_of_text(slide, 1, config={'max_words_per_textbox': 5, 'max_words_per_slide': 100})
        assert len(issues) > 0


class TestFontSize:
    def test_passes_large_font(self):
        _, slide = make_slide_with_text("Big text", font_size=24)
        issues = check_font_size(slide, 1)
        assert len(issues) == 0

    def test_fails_small_font(self):
        _, slide = make_slide_with_text("Tiny text", font_size=10)
        issues = check_font_size(slide, 1)
        assert len(issues) > 0
        assert issues[0]['severity'] == 'error'


class TestPlaceholderResidue:
    def test_passes_normal_text(self):
        _, slide = make_slide_with_text("Regular content here")
        issues = check_placeholder_residue(slide, 1)
        assert len(issues) == 0

    def test_detects_todo(self):
        _, slide = make_slide_with_text("This section is TODO")
        issues = check_placeholder_residue(slide, 1)
        assert len(issues) > 0
        assert issues[0]['category'] == 'placeholder_residue'

    def test_detects_lorem_ipsum(self):
        _, slide = make_slide_with_text("Lorem ipsum dolor sit amet")
        issues = check_placeholder_residue(slide, 1)
        assert len(issues) > 0


class TestDeadSlides:
    def test_passes_slide_with_content(self):
        _, slide = make_slide_with_text("Content")
        issues = check_dead_slides(slide, 1)
        assert len(issues) == 0

    def test_detects_empty_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        issues = check_dead_slides(slide, 1)
        assert len(issues) > 0
        assert issues[0]['category'] == 'missing_content'


class TestSpeakerNotes:
    def test_passes_with_notes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.paragraphs[0].add_run().text = "Content"
        slide.notes_slide.notes_text_frame.text = "Speaker notes here"
        issues = check_speaker_notes(prs)
        assert len(issues) == 0

    def test_warns_without_notes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.paragraphs[0].add_run().text = "Content"
        issues = check_speaker_notes(prs)
        assert len(issues) > 0


class TestSlideCountRatio:
    def test_passes_reasonable_ratio(self):
        prs = Presentation()
        for _ in range(15):
            prs.slides.add_slide(prs.slide_layouts[5])
        issues = check_slide_count_ratio(prs, 30)
        assert len(issues) == 0

    def test_warns_too_many_slides(self):
        prs = Presentation()
        for _ in range(100):
            prs.slides.add_slide(prs.slide_layouts[5])
        issues = check_slide_count_ratio(prs, 10)
        assert len(issues) > 0

    def test_skips_without_duration(self):
        prs = Presentation()
        issues = check_slide_count_ratio(prs, None)
        assert len(issues) == 0


class TestTitleSlide:
    def test_passes_title_layout(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
        issues = check_title_slide(prs)
        assert len(issues) == 0

    def test_warns_no_slides(self):
        prs = Presentation()
        issues = check_title_slide(prs)
        assert len(issues) > 0


class TestClosingSlide:
    def test_passes_with_cta_keywords(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # first slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.paragraphs[0].add_run().text = "Thank you! Questions?"
        issues = check_closing_slide(prs)
        assert len(issues) == 0

    def test_warns_without_cta(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.paragraphs[0].add_run().text = "Random content"
        issues = check_closing_slide(prs)
        assert len(issues) > 0


class TestSafeMargins:
    def test_passes_within_margins(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.paragraphs[0].add_run().text = "Within margins"
        issues = check_safe_margins(slide, 1, prs)
        assert len(issues) == 0


class TestElementOverlap:
    def test_passes_no_overlap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
        box1.text_frame.paragraphs[0].add_run().text = "Box 1"
        box2 = slide.shapes.add_textbox(Inches(5), Inches(0.5), Inches(2), Inches(1))
        box2.text_frame.paragraphs[0].add_run().text = "Box 2"
        issues = check_element_overlap(slide, 1)
        assert len(issues) == 0

    def test_detects_overlap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        box1.text_frame.paragraphs[0].add_run().text = "Box 1"
        box2 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        box2.text_frame.paragraphs[0].add_run().text = "Box 2"
        issues = check_element_overlap(slide, 1)
        assert len(issues) > 0
