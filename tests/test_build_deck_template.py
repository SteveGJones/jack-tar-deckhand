"""Tests for template-driven deck assembly via python-pptx."""

import json
import os
import pytest
from PIL import Image
from pptx import Presentation

from src.template_analyser import analyse_template
from src.assembler.build_deck_template import build_deck

FIXTURE_PATH = os.path.join(os.path.dirname(__file__), 'fixtures', 'templates', 'metamirror-template.pptx')


def _setup_deck_dir(tmp_path, outline_slides, image_manifest_images=None, speaker_notes=None):
    """Create a minimal DeckContext for testing."""
    deck_dir = str(tmp_path / 'deck')
    os.makedirs(os.path.join(deck_dir, 'images'), exist_ok=True)
    os.makedirs(os.path.join(deck_dir, 'output'), exist_ok=True)

    # Template profile
    profile = analyse_template(FIXTURE_PATH)
    with open(os.path.join(deck_dir, 'template-profile.json'), 'w') as f:
        json.dump(profile, f)

    # Outline
    outline = {
        'narrative_arc': 'test',
        'estimated_duration_minutes': 10,
        'slides': outline_slides,
    }
    with open(os.path.join(deck_dir, 'outline.json'), 'w') as f:
        json.dump(outline, f)

    # Image manifest
    manifest = {'images': image_manifest_images or []}
    with open(os.path.join(deck_dir, 'image-manifest.json'), 'w') as f:
        json.dump(manifest, f)

    # Speaker notes
    notes = speaker_notes or {'notes': []}
    with open(os.path.join(deck_dir, 'speaker-notes.json'), 'w') as f:
        json.dump(notes, f)

    return deck_dir, profile


class TestBuildDeckTemplate:
    def test_strips_example_slides(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'Real Slide'},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        # Should have exactly 1 slide (the one from our outline), not the template examples
        assert len(prs.slides) == 1

    def test_populates_title_placeholder(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'My Headline'},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        slide = prs.slides[0]
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert any('My Headline' in t for t in texts)

    def test_populates_body_points(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'Title',
             'body_points': ['Point one', 'Point two', 'Point three']},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        slide = prs.slides[0]
        all_text = ' '.join(shape.text for shape in slide.shapes if shape.has_text_frame)
        assert 'Point one' in all_text
        assert 'Point two' in all_text
        assert 'Point three' in all_text

    def test_output_path_correct(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'Test'},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        assert output_path.endswith('output/presentation.pptx')
        assert os.path.isfile(output_path)

    def test_multiple_slides(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'title', 'headline': 'Title Slide'},
            {'slide_number': 2, 'slide_type': 'content', 'headline': 'Content Slide'},
            {'slide_number': 3, 'slide_type': 'content', 'headline': 'Another Content'},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        assert len(prs.slides) == 3


class TestTemplateAssemblyImages:
    def test_image_inserted_into_picture_placeholder(self, tmp_path):
        img = Image.new('RGB', (400, 300), color='red')
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'With Image'},
        ])
        img_path = os.path.join(deck_dir, 'images', 'slide1.png')
        img.save(img_path)

        manifest = {
            'images': [{
                'image_id': 'img-1',
                'slide_number': 1,
                'file_path': 'images/slide1.png',
                'status': 'generated',
            }]
        }
        with open(os.path.join(deck_dir, 'image-manifest.json'), 'w') as f:
            json.dump(manifest, f)

        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        assert os.path.isfile(output_path)

    def test_missing_image_does_not_crash(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'No Image'},
        ], image_manifest_images=[{
            'image_id': 'img-1',
            'slide_number': 1,
            'file_path': 'images/nonexistent.png',
            'status': 'generated',
        }])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        assert os.path.isfile(output_path)


class TestTemplateAssemblyNotes:
    def test_speaker_notes_populated(self, tmp_path):
        notes = {
            'notes': [
                {'slide_number': 1, 'text': 'Remember to introduce yourself'},
            ]
        }
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'Intro'},
        ], speaker_notes=notes)
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert 'Remember to introduce yourself' in notes_text

    def test_no_notes_does_not_crash(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'content', 'headline': 'No Notes'},
        ])
        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        assert os.path.isfile(output_path)


class TestTemplateSmartArtPlaceholder:
    def test_smartart_placeholder_emitted_for_smartart_slide(self, tmp_path):
        deck_dir, profile = _setup_deck_dir(tmp_path, [
            {'slide_number': 1, 'slide_type': 'diagram', 'headline': 'Architecture'},
        ])
        sa_manifest = {
            'graphics': [{
                'smartart_id': 'sa-1',
                'slide_number': 1,
                'graphic_type': 'flowchart',
                'engine_used': 'pptx_native',
                'enrichment_tier': 'pure_programmatic',
                'file_path': 'smartart/carrier_1.pptx',
                'status': 'rendered',
            }]
        }
        with open(os.path.join(deck_dir, 'smartart-manifest.json'), 'w') as f:
            json.dump(sa_manifest, f)

        output_path = build_deck(deck_dir, FIXTURE_PATH, profile)
        prs = Presentation(output_path)
        slide = prs.slides[0]
        shape_names = [s.name for s in slide.shapes]
        assert any('pptx_native_placeholder_1' in n for n in shape_names)
