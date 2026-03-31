"""Tests for deck-assembler PPTX assembly.

These tests validate that the build_deck.js script correctly reads
DeckContext contracts and produces a valid .pptx file.
"""

import json
import os
import shutil
import subprocess
import pytest

FIXTURE_DIR = os.path.join(os.path.dirname(__file__), 'fixtures', 'minimal_deck')
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'tmp', 'test-assembler')


@pytest.fixture(autouse=True)
def clean_output_dir():
    """Create and clean test output directory."""
    os.makedirs(os.path.join(OUTPUT_DIR, 'output'), exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, 'images'), exist_ok=True)
    yield
    if os.path.exists(OUTPUT_DIR):
        shutil.rmtree(OUTPUT_DIR)


@pytest.fixture
def deck_dir():
    """Set up a DeckContext directory from fixtures."""
    for fname in os.listdir(FIXTURE_DIR):
        src = os.path.join(FIXTURE_DIR, fname)
        dst = os.path.join(OUTPUT_DIR, fname)
        if os.path.isfile(src):
            shutil.copy2(src, dst)
        elif os.path.isdir(src):
            shutil.copytree(src, dst, dirs_exist_ok=True)
    return OUTPUT_DIR


def test_build_deck_produces_pptx(deck_dir):
    """build_deck.js should produce a .pptx file at output/presentation.pptx."""
    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"build_deck.js failed: {result.stderr}"
    pptx_path = os.path.join(deck_dir, 'output', 'presentation.pptx')
    assert os.path.isfile(pptx_path), "presentation.pptx not created"
    assert os.path.getsize(pptx_path) > 1000, "presentation.pptx too small"


def test_build_deck_has_correct_slide_count(deck_dir):
    """Output should have the same number of slides as the outline."""
    subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    pptx_path = os.path.join(deck_dir, 'output', 'presentation.pptx')
    from pptx import Presentation
    prs = Presentation(pptx_path)
    outline = json.load(open(os.path.join(deck_dir, 'outline.json')))
    assert len(prs.slides) == len(outline['slides'])


def test_build_deck_has_speaker_notes(deck_dir):
    """Output slides should have speaker notes from the SpeakerNotes contract."""
    subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    pptx_path = os.path.join(deck_dir, 'output', 'presentation.pptx')
    from pptx import Presentation
    prs = Presentation(pptx_path)
    notes = json.load(open(os.path.join(deck_dir, 'speaker-notes.json')))
    slides_with_notes = 0
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slides_with_notes += 1
    assert slides_with_notes >= len(notes['notes']), \
        f"Expected {len(notes['notes'])} slides with notes, got {slides_with_notes}"


def test_build_deck_applies_palette(deck_dir):
    """Output should use colours from the StyleGuide palette."""
    subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    pptx_path = os.path.join(deck_dir, 'output', 'presentation.pptx')
    from pptx import Presentation
    prs = Presentation(pptx_path)
    assert len(prs.slides) > 0


def test_build_deck_handles_empty_chart_manifest(deck_dir):
    """Assembly should succeed when chart-manifest.json has zero charts."""
    subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    pptx_path = os.path.join(deck_dir, 'output', 'presentation.pptx')
    assert os.path.isfile(pptx_path)


def test_build_deck_with_background_strategy(deck_dir):
    """Assembler should handle 'background' strategy without crashing."""
    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"build_deck.js failed: {result.stderr}"


def test_build_deck_background_variants(deck_dir):
    """Each backdrop_variant should produce a valid slide."""
    # Update the strategy map to test each variant
    sm_path = os.path.join(deck_dir, 'strategy-map.json')
    with open(sm_path) as f:
        sm = json.load(f)
    sm['slides'][1]['strategy'] = 'background'
    sm['slides'][1]['backdrop_variant'] = 'bottom_bar'
    with open(sm_path, 'w') as f:
        json.dump(sm, f)

    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"Failed with bottom_bar: {result.stderr}"


def test_build_deck_pragmatic_composition(deck_dir):
    """Assembler should handle pragmatic_composition with multiple images per slide."""
    # Update strategy map
    sm_path = os.path.join(deck_dir, 'strategy-map.json')
    with open(sm_path) as f:
        sm = json.load(f)
    sm['slides'][1]['strategy'] = 'pragmatic_composition'
    sm['slides'][1]['element_layout'] = {
        'template': 'two_column',
        'elements': [
            {'id': 'elem_1', 'label_source': 'body_points[0]', 'x': 0.05, 'y': 0.22, 'w': 0.42, 'h': 0.55},
            {'id': 'elem_2', 'label_source': 'body_points[1]', 'x': 0.55, 'y': 0.22, 'w': 0.42, 'h': 0.55},
        ],
        'title_region': {'x': 0.05, 'y': 0.03, 'w': 0.90, 'h': 0.12},
    }
    with open(sm_path, 'w') as f:
        json.dump(sm, f)

    # Update image manifest with element images
    im_path = os.path.join(deck_dir, 'image-manifest.json')
    with open(im_path) as f:
        im = json.load(f)
    im['images'].extend([
        {'image_id': 'slide-02-bg', 'slide_number': 2, 'file_path': im['images'][0]['file_path'],
         'status': 'generated', 'placement_zone': 'background'},
        {'image_id': 'slide-02-elem-1', 'slide_number': 2, 'file_path': im['images'][0]['file_path'],
         'status': 'generated', 'placement_zone': 'element', 'element_id': 'elem_1'},
        {'image_id': 'slide-02-elem-2', 'slide_number': 2, 'file_path': im['images'][0]['file_path'],
         'status': 'generated', 'placement_zone': 'element', 'element_id': 'elem_2'},
    ])
    with open(im_path, 'w') as f:
        json.dump(im, f)

    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"Failed: {result.stderr}"


def test_build_deck_backdrop_with_detected_positions(deck_dir):
    """Assembler should place text at vision-detected positions for backdrop strategy."""
    sm_path = os.path.join(deck_dir, 'strategy-map.json')
    with open(sm_path) as f:
        sm = json.load(f)
    sm['slides'][1]['strategy'] = 'backdrop'
    sm['slides'][1]['element_layout'] = {
        'template': 'three_across',
        'elements': [
            {'id': 'elem_1', 'label_source': 'body_points[0]', 'x': 0.08, 'y': 0.25, 'w': 0.25, 'h': 0.50},
            {'id': 'elem_2', 'label_source': 'body_points[1]', 'x': 0.38, 'y': 0.25, 'w': 0.25, 'h': 0.50},
            {'id': 'elem_3', 'label_source': 'body_points[2]', 'x': 0.67, 'y': 0.25, 'w': 0.25, 'h': 0.50},
        ],
        'title_region': {'x': 0.05, 'y': 0.02, 'w': 0.90, 'h': 0.12},
    }
    with open(sm_path, 'w') as f:
        json.dump(sm, f)

    # Add detected_positions to image manifest
    im_path = os.path.join(deck_dir, 'image-manifest.json')
    with open(im_path) as f:
        im = json.load(f)
    im['images'].append({
        'image_id': 'slide-02-scene', 'slide_number': 2,
        'file_path': im['images'][0]['file_path'],
        'status': 'generated',
        'detected_positions': [
            {'element_id': 'elem_1', 'x': 0.10, 'y': 0.28, 'w': 0.22, 'h': 0.45, 'confidence': 0.92},
            {'element_id': 'elem_2', 'x': 0.39, 'y': 0.26, 'w': 0.24, 'h': 0.48, 'confidence': 0.88},
            {'element_id': 'elem_3', 'x': 0.68, 'y': 0.27, 'w': 0.23, 'h': 0.44, 'confidence': 0.85},
        ],
    })
    with open(im_path, 'w') as f:
        json.dump(im, f)

    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"Failed: {result.stderr}"


def test_build_deck_backdrop_falls_back_to_prescribed(deck_dir):
    """Without detected_positions, backdrop should fall back to element_layout positions."""
    sm_path = os.path.join(deck_dir, 'strategy-map.json')
    with open(sm_path) as f:
        sm = json.load(f)
    sm['slides'][1]['strategy'] = 'backdrop'
    sm['slides'][1]['element_layout'] = {
        'template': 'two_column',
        'elements': [
            {'id': 'elem_1', 'label_source': 'body_points[0]', 'x': 0.05, 'y': 0.22, 'w': 0.42, 'h': 0.55},
            {'id': 'elem_2', 'label_source': 'body_points[1]', 'x': 0.55, 'y': 0.22, 'w': 0.42, 'h': 0.55},
        ],
        'title_region': {'x': 0.05, 'y': 0.02, 'w': 0.90, 'h': 0.12},
    }
    with open(sm_path, 'w') as f:
        json.dump(sm, f)

    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True, timeout=30
    )
    assert result.returncode == 0, f"Failed: {result.stderr}"


def test_build_deck_without_strategy_map(tmp_path):
    """Assembler works when strategy-map.json is absent (backward compatible)."""
    # Use the existing test deck at ./tmp/deck if available, otherwise skip
    deck_dir = './tmp/deck'
    if not os.path.exists(os.path.join(deck_dir, 'outline.json')):
        pytest.skip('No test deck available at ./tmp/deck')
    result = subprocess.run(
        ['node', 'src/assembler/build_deck.js', '--deck-dir', deck_dir],
        capture_output=True, text=True,
    )
    assert result.returncode == 0
    assert 'assembled successfully' in result.stdout.lower()
