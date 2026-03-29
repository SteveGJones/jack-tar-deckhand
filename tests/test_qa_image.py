"""Tests for image quality QA checks."""

import io
import pytest
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from src.qa.checks.image_quality import (
    check_image_resolution,
    check_aspect_ratio_distortion,
)


def _make_image_bytes(width, height, fmt='PNG'):
    """Create a test image as bytes."""
    img = Image.new('RGB', (width, height), color=(100, 150, 200))
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    buf.seek(0)
    return buf


class TestImageResolution:
    def test_passes_high_res_image(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_buf = _make_image_bytes(1920, 1080)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(0.5), Inches(5), Inches(3))
        issues = check_image_resolution(slide, 1, prs)
        assert len(issues) == 0

    def test_fails_low_res_image(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_buf = _make_image_bytes(50, 50)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(0.5), Inches(5), Inches(3))
        issues = check_image_resolution(slide, 1, prs)
        assert len(issues) > 0
        assert issues[0]['severity'] == 'error'

    def test_passes_no_images(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        issues = check_image_resolution(slide, 1, prs)
        assert len(issues) == 0


class TestAspectRatioDistortion:
    def test_passes_correct_ratio(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_buf = _make_image_bytes(400, 300)
        # Display at 4:3 ratio (matching native)
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1), Inches(4), Inches(3))
        issues = check_aspect_ratio_distortion(slide, 1, prs)
        assert len(issues) == 0

    def test_warns_distorted_ratio(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_buf = _make_image_bytes(400, 400)  # Square
        # Display at 4:1 ratio (severely distorted)
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1), Inches(8), Inches(2))
        issues = check_aspect_ratio_distortion(slide, 1, prs)
        assert len(issues) > 0
