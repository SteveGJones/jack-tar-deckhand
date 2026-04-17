"""Tests for template analyser and template profile schema."""

import json
import os
import pytest
import jsonschema

from src.deckcontext import read_contract, write_contract


SCHEMA_DIR = os.path.join(os.path.dirname(__file__), '..', 'src', 'schemas')


def _load_schema():
    with open(os.path.join(SCHEMA_DIR, 'template_profile.schema.json')) as f:
        return json.load(f)


def _minimal_profile():
    return {
        'template_path': '/tmp/template.pptx',
        'template_hash': 'a' * 64,
        'slide_width_inches': 13.333,
        'slide_height_inches': 7.5,
        'master_index': 0,
        'layouts': [
            {
                'name': 'Content 1',
                'index': 0,
                'placeholder_count': 2,
                'placeholders': [
                    {'idx': 0, 'type': 'title', 'name': 'Title 1', 'x': 0.6, 'y': 0.6, 'w': 12.13, 'h': 1.02},
                    {'idx': 31, 'type': 'content', 'name': 'Content Placeholder 5', 'x': 0.6, 'y': 2.33, 'w': 12.13, 'h': 4.57},
                ],
                'decorative_shape_count': 0,
            }
        ],
        'layout_mapping': {
            'content': {'layout_name': 'Content 1', 'layout_index': 0},
        },
        'unmapped_fallback': {'layout_name': 'Content 1', 'layout_index': 0},
        'constrains_strategies': True,
        'speaker_approved': False,
    }


class TestTemplateProfileSchema:
    def test_valid_profile_passes(self):
        schema = _load_schema()
        jsonschema.validate(instance=_minimal_profile(), schema=schema)

    def test_missing_required_field_fails(self):
        schema = _load_schema()
        profile = _minimal_profile()
        del profile['template_hash']
        with pytest.raises(jsonschema.ValidationError):
            jsonschema.validate(instance=profile, schema=schema)

    def test_invalid_hash_pattern_fails(self):
        schema = _load_schema()
        profile = _minimal_profile()
        profile['template_hash'] = 'not-a-hash'
        with pytest.raises(jsonschema.ValidationError):
            jsonschema.validate(instance=profile, schema=schema)

    def test_invalid_placeholder_type_fails(self):
        schema = _load_schema()
        profile = _minimal_profile()
        profile['layouts'][0]['placeholders'][0]['type'] = 'invalid_type'
        with pytest.raises(jsonschema.ValidationError):
            jsonschema.validate(instance=profile, schema=schema)

    def test_write_and_read_via_deckcontext(self, tmp_path):
        profile = _minimal_profile()
        deck_dir = str(tmp_path)
        write_contract(deck_dir, 'template-profile', profile)
        loaded = read_contract(deck_dir, 'template-profile')
        assert loaded['template_hash'] == 'a' * 64
        assert loaded['layouts'][0]['name'] == 'Content 1'


from src.template_analyser import classify_placeholder


class TestClassifyPlaceholder:
    def test_title_type(self):
        assert classify_placeholder('TITLE (1)', 'Title 1') == 'title'

    def test_subtitle_type(self):
        assert classify_placeholder('SUBTITLE (2)', 'Subtitle 1') == 'subtitle'

    def test_body_type(self):
        assert classify_placeholder('BODY (2)', 'Text Placeholder 3') == 'body'

    def test_object_type_maps_to_content(self):
        assert classify_placeholder('OBJECT (7)', 'Content Placeholder 5') == 'content'

    def test_picture_type(self):
        assert classify_placeholder('PICTURE (18)', 'Picture Placeholder 10') == 'picture'

    def test_chart_type(self):
        assert classify_placeholder('CHART (12)', 'Chart Placeholder 1') == 'chart'

    def test_table_type(self):
        assert classify_placeholder('TABLE (11)', 'Table Placeholder 1') == 'table'

    def test_date_type(self):
        assert classify_placeholder('DATE (16)', 'Date Placeholder 1') == 'date'

    def test_footer_type(self):
        assert classify_placeholder('FOOTER (15)', 'Footer Placeholder 1') == 'footer'

    def test_slide_number_type(self):
        assert classify_placeholder('SLIDE_NUMBER (13)', 'Slide Number Placeholder 1') == 'slide_number'

    def test_chapter_name_heuristic(self):
        assert classify_placeholder('BODY (2)', 'Chapter Placeholder 2') == 'chapter_box'

    def test_logo_name_heuristic(self):
        assert classify_placeholder('BODY (2)', 'Logo Placeholder 7') == 'other'

    def test_unknown_type(self):
        assert classify_placeholder('UNKNOWN (99)', 'Mystery Shape') == 'other'


from src.template_analyser import extract_layouts

FIXTURE_PATH = os.path.join(os.path.dirname(__file__), 'fixtures', 'templates', 'metamirror-template.pptx')


class TestExtractLayouts:
    def test_returns_list_of_layouts(self):
        layouts = extract_layouts(FIXTURE_PATH)
        assert isinstance(layouts, list)
        assert len(layouts) > 0

    def test_layout_has_required_keys(self):
        layouts = extract_layouts(FIXTURE_PATH)
        layout = layouts[0]
        assert 'name' in layout
        assert 'index' in layout
        assert 'placeholder_count' in layout
        assert 'placeholders' in layout
        assert 'decorative_shape_count' in layout

    def test_placeholder_has_geometry(self):
        layouts = extract_layouts(FIXTURE_PATH)
        # Find a layout with placeholders
        layout_with_phs = next(l for l in layouts if l['placeholder_count'] > 0)
        ph = layout_with_phs['placeholders'][0]
        assert 'idx' in ph
        assert 'type' in ph
        assert 'name' in ph
        assert isinstance(ph['x'], float)
        assert isinstance(ph['y'], float)
        assert isinstance(ph['w'], float)
        assert isinstance(ph['h'], float)

    def test_title_layout_has_title_placeholder(self):
        layouts = extract_layouts(FIXTURE_PATH)
        title_layout = next(l for l in layouts if 'title' in l['name'].lower())
        ph_types = [p['type'] for p in title_layout['placeholders']]
        assert 'title' in ph_types

    def test_master_index_zero_by_default(self):
        layouts = extract_layouts(FIXTURE_PATH, master_index=0)
        assert len(layouts) > 0


from src.template_analyser import auto_map_layouts


class TestAutoMapLayouts:
    def test_maps_title_slide(self):
        layouts = [
            {'name': 'Title Slide', 'index': 0, 'placeholder_count': 2,
             'placeholders': [{'idx': 0, 'type': 'title', 'name': 'Title 1', 'x': 0.5, 'y': 1.0, 'w': 8.0, 'h': 4.5}],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'title' in mapping
        assert mapping['title']['layout_name'] == 'Title Slide'

    def test_maps_content_layout(self):
        layouts = [
            {'name': 'Content 1', 'index': 5, 'placeholder_count': 2,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 12.0, 'h': 1.0},
                 {'idx': 31, 'type': 'content', 'name': 'Content Placeholder', 'x': 0.6, 'y': 2.3, 'w': 12.0, 'h': 4.5},
             ],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'content' in mapping
        assert mapping['content']['layout_name'] == 'Content 1'

    def test_maps_divider_to_section_divider(self):
        layouts = [
            {'name': 'Divider 1', 'index': 12, 'placeholder_count': 2,
             'placeholders': [{'idx': 0, 'type': 'title', 'name': 'Title', 'x': 5.6, 'y': 5.0, 'w': 7.0, 'h': 1.8}],
             'decorative_shape_count': 1},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'section_divider' in mapping
        assert mapping['section_divider']['layout_name'] == 'Divider 1'

    def test_maps_end_slide_to_closing(self):
        layouts = [
            {'name': 'End Slide 1', 'index': 56, 'placeholder_count': 0,
             'placeholders': [],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'closing' in mapping

    def test_maps_comparison_layout(self):
        layouts = [
            {'name': 'Comparison', 'index': 33, 'placeholder_count': 5,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 12.0, 'h': 1.0},
                 {'idx': 1, 'type': 'content', 'name': 'Content 1', 'x': 0.6, 'y': 2.3, 'w': 5.5, 'h': 4.5},
                 {'idx': 2, 'type': 'content', 'name': 'Content 2', 'x': 7.0, 'y': 2.3, 'w': 5.5, 'h': 4.5},
             ],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'comparison' in mapping

    def test_content_with_picture_maps_to_content_with_image(self):
        layouts = [
            {'name': 'Content Photo 1', 'index': 35, 'placeholder_count': 3,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 6.0, 'h': 1.0},
                 {'idx': 1, 'type': 'content', 'name': 'Content', 'x': 0.6, 'y': 2.3, 'w': 6.0, 'h': 4.5},
                 {'idx': 13, 'type': 'picture', 'name': 'Picture', 'x': 7.0, 'y': 0.0, 'w': 6.3, 'h': 7.5},
             ],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert 'content_with_image' in mapping

    def test_prefers_simpler_layout_when_multiple_match(self):
        layouts = [
            {'name': 'Content 1', 'index': 5, 'placeholder_count': 2,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 12.0, 'h': 1.0},
                 {'idx': 31, 'type': 'content', 'name': 'Content', 'x': 0.6, 'y': 2.3, 'w': 12.0, 'h': 4.5},
             ],
             'decorative_shape_count': 0},
            {'name': 'Content 1 Chapterbox', 'index': 6, 'placeholder_count': 3,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 12.0, 'h': 1.0},
                 {'idx': 31, 'type': 'content', 'name': 'Content', 'x': 0.6, 'y': 2.3, 'w': 12.0, 'h': 4.5},
                 {'idx': 13, 'type': 'chapter_box', 'name': 'Chapter', 'x': 0.6, 'y': 0.29, 'w': 12.0, 'h': 0.17},
             ],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        assert mapping['content']['layout_name'] == 'Content 1'

    def test_fallback_is_largest_content_layout(self):
        layouts = [
            {'name': 'Content 1', 'index': 5, 'placeholder_count': 2,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 12.0, 'h': 1.0},
                 {'idx': 31, 'type': 'content', 'name': 'Content', 'x': 0.6, 'y': 2.3, 'w': 12.0, 'h': 4.57},
             ],
             'decorative_shape_count': 0},
            {'name': 'Content 2', 'index': 6, 'placeholder_count': 3,
             'placeholders': [
                 {'idx': 0, 'type': 'title', 'name': 'Title', 'x': 0.6, 'y': 0.6, 'w': 3.6, 'h': 1.0},
                 {'idx': 31, 'type': 'content', 'name': 'Content', 'x': 0.6, 'y': 2.3, 'w': 3.6, 'h': 4.57},
                 {'idx': 28, 'type': 'content', 'name': 'Content 2', 'x': 5.0, 'y': 2.3, 'w': 7.6, 'h': 4.6},
             ],
             'decorative_shape_count': 0},
        ]
        mapping, fallback = auto_map_layouts(layouts)
        # Content 1 has the single largest content placeholder (12.0 * 4.57)
        assert fallback['layout_name'] == 'Content 1'

    def test_empty_layouts_returns_empty_mapping(self):
        mapping, fallback = auto_map_layouts([])
        assert mapping == {}
        assert fallback is None


from src.template_analyser import analyse_template


class TestAnalyseTemplate:
    def test_returns_complete_profile(self):
        profile = analyse_template(FIXTURE_PATH)
        assert profile['template_path'] == FIXTURE_PATH
        assert len(profile['template_hash']) == 64
        assert profile['slide_width_inches'] == 13.33
        assert profile['slide_height_inches'] == 7.5
        assert profile['master_index'] == 0
        assert isinstance(profile['layouts'], list)
        assert isinstance(profile['layout_mapping'], dict)
        assert profile['constrains_strategies'] is True
        assert profile['speaker_approved'] is False

    def test_hash_changes_with_different_file(self, tmp_path):
        from pptx import Presentation as Prs
        prs = Prs()
        path1 = str(tmp_path / 'a.pptx')
        prs.save(path1)
        profile1 = analyse_template(path1)

        # Modify and save again
        prs2 = Prs()
        prs2.slides.add_slide(prs2.slide_layouts[0])
        path2 = str(tmp_path / 'b.pptx')
        prs2.save(path2)
        profile2 = analyse_template(path2)

        assert profile1['template_hash'] != profile2['template_hash']

    def test_validates_against_schema(self):
        profile = analyse_template(FIXTURE_PATH)
        schema = _load_schema()
        jsonschema.validate(instance=profile, schema=schema)

    def test_unmapped_fallback_present(self):
        profile = analyse_template(FIXTURE_PATH)
        assert profile['unmapped_fallback'] is not None or len(profile['layout_mapping']) > 0
