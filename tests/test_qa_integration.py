"""Integration test for deck-qa: creates a .pptx and runs QA against it."""

import json
import os
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.qa.run_qa import run_qa


@pytest.fixture
def minimal_pptx(tmp_path):
    """Create a minimal .pptx for QA testing."""
    pptx_path = str(tmp_path / 'test.pptx')
    prs = Presentation()

    # Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Test Presentation"
    slide1.placeholders[1].text = "A test deck"

    # Content slide with notes
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Content Slide"
    slide2.placeholders[1].text = "Some content here"
    slide2.notes_slide.notes_text_frame.text = "Speaker notes for slide 2"

    # Closing slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    txBox.text_frame.paragraphs[0].add_run().text = "Thank you! Questions?"

    prs.save(pptx_path)
    return pptx_path


def test_qa_runs_without_error(minimal_pptx, tmp_path):
    """QA pipeline should run and produce a valid report."""
    report = run_qa(minimal_pptx, str(tmp_path))
    assert report['verdict'] in ('pass', 'pass_with_warnings', 'fail')
    assert 'summary' in report
    assert 'findings' in report
    assert isinstance(report['findings'], list)


def test_qa_report_has_correct_slide_count(minimal_pptx, tmp_path):
    """Report should reflect actual slide count."""
    report = run_qa(minimal_pptx, str(tmp_path))
    assert report['summary']['total_slides'] == 3


def test_qa_with_duration(minimal_pptx, tmp_path):
    """QA should accept duration parameter for AP-10 check."""
    report = run_qa(minimal_pptx, str(tmp_path), duration_minutes=30)
    assert report['verdict'] in ('pass', 'pass_with_warnings', 'fail')


def test_qa_schema_validation(minimal_pptx, tmp_path):
    """QAReport should validate against the JSON schema."""
    report = run_qa(minimal_pptx, str(tmp_path))

    schema_path = os.path.join(os.path.dirname(__file__), '..', 'src', 'schemas', 'qa_report.schema.json')
    if os.path.exists(schema_path):
        from jsonschema import validate
        with open(schema_path) as f:
            schema = json.load(f)
        validate(instance=report, schema=schema)
