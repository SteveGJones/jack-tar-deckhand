"""Tests for cross-slide consistency QA checks."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.qa.checks.consistency import (
    check_font_families,
    check_consecutive_bullet_slides,
    check_heading_consistency,
    check_branding_consistency,
)


class TestFontFamilies:
    def test_passes_single_font(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "Hello"
        run.font.name = "Arial"
        issues = check_font_families(prs)
        assert len(issues) == 0

    def test_passes_two_fonts(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        run1 = box1.text_frame.paragraphs[0].add_run()
        run1.text = "Heading"
        run1.font.name = "Arial"
        box2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        run2 = box2.text_frame.paragraphs[0].add_run()
        run2.text = "Body text"
        run2.font.name = "Times New Roman"
        issues = check_font_families(prs)
        assert len(issues) == 0

    def test_warns_three_fonts(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        for i, font_name in enumerate(["Arial", "Times New Roman", "Courier New"]):
            box = slide.shapes.add_textbox(Inches(1), Inches(i + 1), Inches(2), Inches(0.5))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = f"Text in {font_name}"
            run.font.name = font_name
        issues = check_font_families(prs)
        assert len(issues) > 0
        assert issues[0]['category'] == 'consistency'


class TestConsecutiveBulletSlides:
    def test_passes_few_bullet_slides(self):
        prs = Presentation()
        for _ in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
            tf = box.text_frame
            for j in range(5):
                p = tf.add_paragraph()
                p.text = f"Point {j}"
        issues = check_consecutive_bullet_slides(prs)
        assert len(issues) == 0

    def test_passes_empty_deck(self):
        prs = Presentation()
        issues = check_consecutive_bullet_slides(prs)
        assert len(issues) == 0


class TestHeadingConsistency:
    def test_passes_empty_deck(self):
        prs = Presentation()
        issues = check_heading_consistency(prs)
        assert len(issues) == 0


class TestBrandingConsistency:
    def test_passes_no_logos(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        issues = check_branding_consistency(prs)
        assert len(issues) == 0
