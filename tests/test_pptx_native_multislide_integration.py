"""Phase 3.4 — multi-slide deck integration test.

Verifies the pptx_native injection coexists with other Jack-Tar
rendering strategies in a realistic multi-slide deck. The Phase 3
unit/integration tests previously used single-slide hosts; this
test fills that gap.

Setup:
    - Copies the minimal_deck fixture (3 slides: full_render title,
      background content, composed closing) and extends it to 6 slides
      by adding:
        Slide 4: composed (intentional non-target between two pptx_native)
        Slide 5: smartart pptx_native flowchart
        Slide 6: smartart pptx_native cycle

    Wait — for the test to actually exercise pptx_native injection,
    we need the smartart slides to be IN the deck. Final layout:
        Slide 1: full_render title (existing)
        Slide 2: background content (existing — has image)
        Slide 3: composed closing (existing)
        Slide 4: composed content (new, no image)
        Slide 5: smartart pptx_native flowchart (new) — INJECT TARGET
        Slide 6: smartart pptx_native cycle (new)        — INJECT TARGET

Verification:
    - Slides 1-4 byte-identical between pre-injection and post-injection
      .pptx (proves zero side effects on non-target slides)
    - Shared resources (theme, slide masters, slide layouts,
      presentation.xml) byte-identical
    - Slides 5 and 6 have <p:graphicFrame> elements wired to
      independently-numbered diagrams (data1+data2 etc.)
    - rId allocation per slide is independent (slide 5's rels and
      slide 6's rels don't collide)
    - [Content_Types].xml has 8 new diagram overrides
    - python-pptx can still open the deck (zip + structure validity)
    - SA-06/07/08 QA checks pass for both injected slides
"""
from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

import pytest

REPO_ROOT = Path(__file__).resolve().parent.parent
MINIMAL_DECK_FIXTURE = REPO_ROOT / "tests" / "fixtures" / "minimal_deck"


# ---------------------------------------------------------------------------
# Fixture builder
# ---------------------------------------------------------------------------

def _build_extended_deck_dir(tmp: Path) -> Path:
    """Copy minimal_deck and extend its contracts to 6 slides with mixed
    strategies including 2 pptx_native slides.

    Returns the path to the prepared deck directory.
    """
    deck_dir = tmp / "deck"
    shutil.copytree(MINIMAL_DECK_FIXTURE, deck_dir)
    (deck_dir / "output").mkdir(exist_ok=True)
    (deck_dir / "images").mkdir(exist_ok=True)

    # ---- Extend outline.json -------------------------------------------
    outline = json.loads((deck_dir / "outline.json").read_text())
    # Existing slides 1-3 stay as they are.
    # Add slides 4, 5, 6.
    outline["total_slides"] = 6
    outline["slides"].extend([
        {
            "slide_number": 4,
            "slide_type": "content",
            "headline": "Why This Matters",
            "body_points": [
                "Faster iteration cycles",
                "Lower error rates",
                "Better maintainability",
            ],
            "narrative_beat": "evidence-2",
            "visual_type": "none",
            "layout_template": "content",
        },
        {
            "slide_number": 5,
            "slide_type": "content",
            "headline": "Our Development Process",
            "body_points": [
                "Research",
                "Design",
                "Build",
                "Test",
                "Ship",
            ],
            "narrative_beat": "evidence-3",
            "visual_type": "diagram",
            "layout_template": "content",
        },
        {
            "slide_number": 6,
            "slide_type": "content",
            "headline": "Continuous Improvement",
            "body_points": [
                "Plan",
                "Build",
                "Test",
                "Deploy",
                "Monitor",
                "Learn",
            ],
            "narrative_beat": "evidence-4",
            "visual_type": "diagram",
            "layout_template": "content",
        },
    ])
    (deck_dir / "outline.json").write_text(json.dumps(outline, indent=2))

    # ---- Extend strategy-map.json --------------------------------------
    strategy = json.loads((deck_dir / "strategy-map.json").read_text())
    # Existing slides 1-3 keep their strategies (full_render, background, composed)
    strategy["slides"].extend([
        {
            "slide_number": 4,
            "strategy": "composed",
            "rationale": "Plain content slide between two smartart slides",
            "render_funnel": ["ollama"],
            "speaker_override": None,
        },
        {
            "slide_number": 5,
            "strategy": "smartart",
            "rationale": "Process flowchart via pptx_native",
            "render_funnel": ["ollama"],
            "speaker_override": None,
        },
        {
            "slide_number": 6,
            "strategy": "smartart",
            "rationale": "Cycle diagram via pptx_native",
            "render_funnel": ["ollama"],
            "speaker_override": None,
        },
    ])
    (deck_dir / "strategy-map.json").write_text(json.dumps(strategy, indent=2))

    # ---- Extend speaker-notes.json -------------------------------------
    notes = json.loads((deck_dir / "speaker-notes.json").read_text())
    notes["notes"].extend([
        {
            "slide_number": 4,
            "text": "Three concrete benefits speakers should remember.",
            "estimated_seconds": 30,
            "timing_marker": "~5:00",
            "cues": [],
        },
        {
            "slide_number": 5,
            "text": "Walk through the five-step development process. The graphic is editable in PowerPoint.",
            "estimated_seconds": 45,
            "timing_marker": "~5:30",
            "cues": [],
        },
        {
            "slide_number": 6,
            "text": "Cycle diagram showing the continuous improvement loop. Six stages.",
            "estimated_seconds": 45,
            "timing_marker": "~6:15",
            "cues": [],
        },
    ])
    (deck_dir / "speaker-notes.json").write_text(json.dumps(notes, indent=2))

    return deck_dir


def _generate_carriers(tmp: Path) -> tuple[Path, Path]:
    """Generate carrier .pptx files for the two pptx_native slides.

    Returns (slide5_carrier_path, slide6_carrier_path).
    """
    from src.smartart_pptx_native import engine

    carriers_dir = tmp / "carriers"
    carriers_dir.mkdir()

    slide5_result = engine.render(
        {
            "slide_number": 5,
            "graphic_type": "flowchart",
            "data": {"steps": ["Research", "Design", "Build", "Test", "Ship"]},
        },
        carriers_dir,
        output_name="slide5_flowchart.pptx",
    )
    slide6_result = engine.render(
        {
            "slide_number": 6,
            "graphic_type": "cycle",
            "data": {"stages": ["Plan", "Build", "Test", "Deploy", "Monitor", "Learn"]},
        },
        carriers_dir,
        output_name="slide6_cycle.pptx",
    )
    return slide5_result.output_path, slide6_result.output_path


def _write_smartart_manifest(deck_dir: Path, slide5_carrier: Path, slide6_carrier: Path):
    """Write smartart-manifest.json with the two pptx_native entries."""
    manifest = {
        "generated_at": "2026-04-08T00:00:00Z",
        "graphics": [
            {
                "smartart_id": "sa-slide-5-flowchart",
                "slide_number": 5,
                "graphic_type": "flowchart",
                "engine_used": "pptx_native",
                "enrichment_tier": "pure_programmatic",
                "file_path": str(slide5_carrier),
                "svg_source_path": "",
                "status": "rendered",
                "dimensions": {"width": 1920, "height": 1080},
                "node_count": 5,
                "alt_text": "Five-step development process",
                "content_hash": "test0001",
            },
            {
                "smartart_id": "sa-slide-6-cycle",
                "slide_number": 6,
                "graphic_type": "cycle",
                "engine_used": "pptx_native",
                "enrichment_tier": "pure_programmatic",
                "file_path": str(slide6_carrier),
                "svg_source_path": "",
                "status": "rendered",
                "dimensions": {"width": 1920, "height": 1080},
                "node_count": 6,
                "alt_text": "Six-stage continuous improvement cycle",
                "content_hash": "test0002",
            },
        ],
    }
    (deck_dir / "smartart-manifest.json").write_text(json.dumps(manifest, indent=2))


def _run_assembler(deck_dir: Path) -> subprocess.CompletedProcess:
    return subprocess.run(
        ["node", "src/assembler/build_deck.js", "--deck-dir", str(deck_dir)],
        capture_output=True,
        text=True,
        timeout=120,
        cwd=str(REPO_ROOT),
    )


def _zip_contents(pptx_path: Path) -> dict[str, bytes]:
    """Read all entries from a .pptx zip into a dict for byte comparison."""
    contents: dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as z:
        for name in z.namelist():
            contents[name] = z.read(name)
    return contents


# ---------------------------------------------------------------------------
# The integration test
# ---------------------------------------------------------------------------

def test_multislide_pptx_native_injection_coexists_with_other_strategies():
    """End-to-end: 6-slide deck with mixed strategies, two pptx_native
    slides, full pipeline (carriers → JS assembler → Python injection
    → QA), with byte-identity guarantees on non-target slides.
    """
    from src.smartart_pptx_native.pipeline import run_injection_step
    from src.qa.checks.smartart_checks import check_all_pptx_native_graphics

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)

        # ---- 1. Build the extended deck fixture -----------------------
        deck_dir = _build_extended_deck_dir(tmp)

        # ---- 2. Generate carriers via the engine ----------------------
        slide5_carrier, slide6_carrier = _generate_carriers(tmp)
        assert slide5_carrier.exists()
        assert slide6_carrier.exists()

        # ---- 3. Write the smartart manifest ---------------------------
        _write_smartart_manifest(deck_dir, slide5_carrier, slide6_carrier)

        # ---- 4. Run the JS assembler ----------------------------------
        result = _run_assembler(deck_dir)
        assert result.returncode == 0, (
            f"build_deck.js failed:\nstdout={result.stdout}\nstderr={result.stderr}"
        )
        assembled = deck_dir / "output" / "presentation.pptx"
        assert assembled.exists(), "build_deck.js did not produce presentation.pptx"

        # ---- 5. Verify it's a 6-slide deck before injection -----------
        from pptx import Presentation
        prs_before = Presentation(assembled)
        assert len(prs_before.slides) == 6, (
            f"expected 6 slides before injection, got {len(prs_before.slides)}"
        )

        # ---- 6. Capture pre-injection state ---------------------------
        pre_state = _zip_contents(assembled)

        # Verify placeholders are present in slides 5 and 6 pre-injection
        slide5_xml_pre = pre_state["ppt/slides/slide5.xml"].decode("utf-8")
        slide6_xml_pre = pre_state["ppt/slides/slide6.xml"].decode("utf-8")
        assert "pptx_native_placeholder_5" in slide5_xml_pre, (
            "JS assembler did not emit placeholder for slide 5"
        )
        assert "pptx_native_placeholder_6" in slide6_xml_pre, (
            "JS assembler did not emit placeholder for slide 6"
        )

        # ---- 7. Run injection -----------------------------------------
        injection_result = run_injection_step(deck_dir)
        assert injection_result.injected_count == 2, (
            f"expected 2 injections, got {injection_result.injected_count}"
        )

        # ---- 8. Capture post-injection state --------------------------
        post_state = _zip_contents(assembled)

        # ---- 9. NON-TARGET SLIDES MUST BE BYTE-IDENTICAL --------------
        # Slides 1, 2, 3, 4 — and their per-slide rels — must not have
        # been touched by the injection process.
        non_target_files = [
            "ppt/slides/slide1.xml",
            "ppt/slides/slide2.xml",
            "ppt/slides/slide3.xml",
            "ppt/slides/slide4.xml",
            "ppt/slides/_rels/slide1.xml.rels",
            "ppt/slides/_rels/slide2.xml.rels",
            "ppt/slides/_rels/slide3.xml.rels",
            "ppt/slides/_rels/slide4.xml.rels",
        ]
        for path in non_target_files:
            assert pre_state[path] == post_state[path], (
                f"non-target file {path!r} was modified by injection — "
                f"this is the side-effect-on-other-slides bug we built "
                f"this test to catch"
            )

        # ---- 10. SHARED RESOURCES MUST BE BYTE-IDENTICAL --------------
        shared_resources = [
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/theme/theme1.xml",
        ]
        # Slide masters and slide layouts (variable count)
        for name in pre_state:
            if name.startswith("ppt/slideMasters/") or name.startswith("ppt/slideLayouts/"):
                shared_resources.append(name)

        for path in shared_resources:
            assert pre_state[path] == post_state[path], (
                f"shared resource {path!r} was modified by injection — "
                f"injection must only touch the target slide files and "
                f"package-level [Content_Types].xml"
            )

        # ---- 11. TARGET SLIDES MUST HAVE GRAPHIC FRAMES --------------
        slide5_xml_post = post_state["ppt/slides/slide5.xml"].decode("utf-8")
        slide6_xml_post = post_state["ppt/slides/slide6.xml"].decode("utf-8")

        # Placeholders removed
        assert "pptx_native_placeholder_5" not in slide5_xml_post
        assert "pptx_native_placeholder_6" not in slide6_xml_post

        # Graphic frames present
        assert "<p:graphicFrame>" in slide5_xml_post
        assert "<p:graphicFrame>" in slide6_xml_post

        # Diagram references present
        assert "dgm:relIds" in slide5_xml_post
        assert "dgm:relIds" in slide6_xml_post

        # ---- 12. INDEPENDENT DIAGRAM NUMBERING ------------------------
        # Slide 5 should reference data1.xml etc., slide 6 should
        # reference data2.xml etc., proving _next_free_diagram_number
        # works across multiple injections in one call.
        assert "ppt/diagrams/data1.xml" in post_state
        assert "ppt/diagrams/layout1.xml" in post_state
        assert "ppt/diagrams/quickStyle1.xml" in post_state
        assert "ppt/diagrams/colors1.xml" in post_state
        assert "ppt/diagrams/data2.xml" in post_state
        assert "ppt/diagrams/layout2.xml" in post_state
        assert "ppt/diagrams/quickStyle2.xml" in post_state
        assert "ppt/diagrams/colors2.xml" in post_state

        # ---- 13. PER-SLIDE RID ALLOCATION IS INDEPENDENT --------------
        slide5_rels = post_state[
            "ppt/slides/_rels/slide5.xml.rels"
        ].decode("utf-8")
        slide6_rels = post_state[
            "ppt/slides/_rels/slide6.xml.rels"
        ].decode("utf-8")

        # Each slide's rels should have its own diagram relationships,
        # pointing at the correct diagram_number
        assert "diagramData" in slide5_rels
        assert "data1.xml" in slide5_rels  # slide 5 → diagram 1
        assert "diagramData" in slide6_rels
        assert "data2.xml" in slide6_rels  # slide 6 → diagram 2

        # ---- 14. CONTENT TYPES HAS 8 NEW DIAGRAM OVERRIDES ------------
        ct_post = post_state["[Content_Types].xml"].decode("utf-8")
        for n in (1, 2):
            for part in ("data", "layout", "quickStyle", "colors"):
                expected = f'/ppt/diagrams/{part}{n}.xml'
                assert expected in ct_post, (
                    f"missing content-type override for {expected}"
                )

        # ---- 15. python-pptx STILL OPENS THE DECK CLEANLY ------------
        # If injection corrupted any structural part, python-pptx would
        # raise here.
        prs_after = Presentation(assembled)
        assert len(prs_after.slides) == 6
        # Speaker notes should still be present on the slides that had them
        for slide in prs_after.slides:
            # Just touching the slide_layout traversal to surface any
            # broken refs
            _ = slide.slide_layout

        # ---- 16. QA CHECKS PASS FOR BOTH PPTX_NATIVE SLIDES -----------
        manifest = json.loads(
            (deck_dir / "smartart-manifest.json").read_text()
        )
        findings = check_all_pptx_native_graphics(assembled, manifest)
        assert findings == [], f"QA findings: {findings}"

        # ---- 17. STEP LABELS PRESENT IN BOTH DIAGRAMS -----------------
        # Sanity check: the carrier content actually made it through
        data1 = post_state["ppt/diagrams/data1.xml"].decode("utf-8")
        data2 = post_state["ppt/diagrams/data2.xml"].decode("utf-8")
        for step in ["Research", "Design", "Build", "Test", "Ship"]:
            assert f"<a:t>{step}</a:t>" in data1, (
                f"slide 5 diagram missing label {step!r}"
            )
        for stage in ["Plan", "Test", "Deploy", "Monitor", "Learn"]:
            assert f"<a:t>{stage}</a:t>" in data2, (
                f"slide 6 diagram missing label {stage!r}"
            )

        # ---- 18. STRUCTURAL INVARIANT: NO STALE DRAWING CACHES --------
        for n in (1, 2):
            assert f"ppt/diagrams/drawing{n}.xml" not in post_state, (
                f"drawing{n}.xml should not exist — PowerPoint regenerates "
                f"the cache from layout{n}.xml on first open"
            )
