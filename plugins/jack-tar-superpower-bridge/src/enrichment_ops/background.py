"""Op1 — apply an AI background image to a slide and remove its BG marker.

In-memory variant — the orchestrator's transactional gate (src/enrichment.py)
saves once after all ops succeed.
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.parts.image import Image, ImagePart

from src.security import resolve_within_allowlist


def _remove_existing_bg(cSld_elem) -> None:
    existing = cSld_elem.find(qn("p:bg"))
    if existing is not None:
        cSld_elem.remove(existing)


def _build_bg_element(rid: str):
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    xml = (
        f'<p:bg xmlns:p="{nsmap["p"]}" xmlns:a="{nsmap["a"]}" xmlns:r="{nsmap["r"]}">'
        "<p:bgPr>"
        '<a:blipFill dpi="0" rotWithShape="1">'
        f'<a:blip r:embed="{rid}"/>'
        "<a:srcRect/>"
        "<a:stretch><a:fillRect/></a:stretch>"
        "</a:blipFill>"
        "<a:effectLst/>"
        "</p:bgPr>"
        "</p:bg>"
    )
    return etree.fromstring(xml)


def _shape_text_content(sp_elem) -> str:
    """Concatenate all <a:t> text in a shape's <p:txBody> for matching.

    Returns the empty string when the shape has no text body. Used to detect
    the Run 5 Finding #17 anti-pattern where /pptx emits a separate addText
    shape whose visible text reproduces the BG marker name as a label.
    """
    a_t_qn = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    parts: list[str] = []
    for t_elem in sp_elem.iter(a_t_qn):
        if t_elem.text:
            parts.append(t_elem.text)
    return "".join(parts).strip()


def apply_background_in_memory(
    *,
    prs,
    slide_index_1based: int,
    image_path: Path,
    marker_name: str | None,
    allowed_image_roots: list[Path],
) -> None:
    """Apply `image_path` as the background of slide `slide_index_1based`.

    `image_path` is resolved through the image-path allowlist.
    The BG marker shape (if `marker_name` provided) is removed from the slide.
    The Presentation `prs` is mutated in place; the caller saves.

    Run 5 Finding #17 — when ``marker_name`` is supplied, this function also
    removes any shape on the slide whose visible text content matches the
    marker name, even if the shape's ``cNvPr@name`` differs. /pptx commonly
    emits a separate ``addText`` rendering "BG:slug" as a visible label
    inside the marker rect; that label survives the rect swap as residual
    cosmetic text on top of the new background image. Defence-in-depth on
    top of the narrative-brief-architect's authoring guidance.
    """
    safe_image = resolve_within_allowlist(image_path, allowed_roots=allowed_image_roots)
    if slide_index_1based < 1 or slide_index_1based > len(prs.slides):
        raise IndexError(
            f"slide {slide_index_1based} out of range (1..{len(prs.slides)})"
        )
    slide = prs.slides[slide_index_1based - 1]

    img = Image.from_file(str(safe_image))
    image_part = ImagePart.new(slide.part.package, img)
    rid = slide.part.relate_to(image_part, RT.IMAGE)

    cSld = slide.element.find(qn("p:cSld"))
    _remove_existing_bg(cSld)
    bg = _build_bg_element(rid)
    cSld.insert(0, bg)

    if marker_name:
        spTree = cSld.find(qn("p:spTree"))
        for sp in list(spTree):
            nvSpPr = sp.find(qn("p:nvSpPr"))
            if nvSpPr is None:
                continue
            cNvPr = nvSpPr.find(qn("p:cNvPr"))
            shape_name = cNvPr.get("name") if cNvPr is not None else None
            if shape_name == marker_name:
                spTree.remove(sp)
                continue
            # Finding #17: remove residual label whose visible text
            # exactly equals the marker name. Strict equality keeps
            # legitimate slide text that merely contains the marker
            # string as a substring intact.
            if _shape_text_content(sp) == marker_name:
                spTree.remove(sp)
