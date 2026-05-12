"""Op — replace a named CHART:* marker shape with a native python-pptx chart.

Issue #55 — CHART:slug marker kind for the superpower bridge.

Design decisions (locked at plan time):
  1. python-pptx ``slide.shapes.add_chart()`` — all-Python, no Node round-trip.
  2. ``chart_spec`` supplied by operator on EnrichmentItem at /enrich-deck time.
  3. Seven v1 chart types: bar, column, line, area, pie, doughnut, scatter.
  4. Series colours derived via ``derive_chart_palette()`` from BrandPalette.
  5. Marker rect + sibling addText label removed via existing v0.2 #23 helpers.
"""
from __future__ import annotations

import logging
from typing import Any

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

from src.chart_palette import derive_chart_palette
from src.colors_xml_builder import BrandPalette

log = logging.getLogger(__name__)

# Map from chart_spec["type"] string to python-pptx XL_CHART_TYPE member.
# v1 scope: bar, column, line, area, pie, doughnut, scatter.
_CHART_TYPE_MAP: dict[str, Any] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

# Default palette role when chart_spec omits "palette_role".
_DEFAULT_PALETTE_ROLE = "default"


def _validate_chart_spec(spec: dict[str, Any], marker_name: str) -> None:
    """Raise ValueError if required fields are missing or malformed.

    Uses ValueError (not EnrichmentApplyError) to avoid a circular import:
    enrichment_ops.chart is imported by enrichment, so importing
    EnrichmentApplyError back from enrichment would create a cycle.
    The transactional orchestrator in enrichment.py wraps all exceptions
    as EnrichmentApplyError automatically.
    """
    missing = []
    for field in ("type", "categories", "series"):
        if field not in spec:
            missing.append(field)
    if missing:
        raise ValueError(
            f"chart_spec for {marker_name!r} is missing required "
            f"field(s): {', '.join(missing)}. Required fields: "
            f"type (str), categories (list), series (list of {{name, values}})."
        )

    chart_type = spec["type"]
    if chart_type not in _CHART_TYPE_MAP:
        valid = ", ".join(sorted(_CHART_TYPE_MAP))
        raise ValueError(
            f"chart_spec for {marker_name!r} has unsupported type {chart_type!r}. "
            f"Valid v1 types: {valid}."
        )

    if not isinstance(spec["categories"], list) or len(spec["categories"]) == 0:
        raise ValueError(
            f"chart_spec for {marker_name!r}: 'categories' must be a non-empty list."
        )

    series = spec["series"]
    if not isinstance(series, list) or len(series) == 0:
        raise ValueError(
            f"chart_spec for {marker_name!r}: 'series' must be a non-empty list."
        )
    for i, s in enumerate(series):
        if not isinstance(s, dict) or "name" not in s or "values" not in s:
            raise ValueError(
                f"chart_spec for {marker_name!r}: series[{i}] must be a dict "
                f"with 'name' (str) and 'values' (list) keys."
            )
        if not isinstance(s["values"], list) or len(s["values"]) == 0:
            raise ValueError(
                f"chart_spec for {marker_name!r}: series[{i}]['values'] "
                f"must be a non-empty list."
            )


def _find_marker_shape(prs, slide_index_1based: int, marker_name: str):
    """Return the shape whose name matches marker_name on the given slide.

    Returns None if the marker is not found — caller raises EnrichmentApplyError.
    """
    slide = prs.slides[slide_index_1based - 1]
    for shape in slide.shapes:
        if shape.name == marker_name:
            return shape, slide
    return None, slide


def _apply_series_colours(chart, hex_colours: list[str]) -> None:
    """Apply series fill colours to a python-pptx Chart object.

    Sets a solid fill (RGB) on each series's ``format.fill``. Series that
    don't have a corresponding colour in the list are left at their default.
    Colours are 6-char uppercase hex strings without leading ``#``.
    """
    from pptx.dml.color import RGBColor

    for i, series in enumerate(chart.series):
        if i >= len(hex_colours):
            break
        hex_str = hex_colours[i]
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(r, g, b)


def _apply_chart_formatting(chart, spec: dict[str, Any]) -> None:
    """Apply optional formatting fields from chart_spec to the chart object."""
    # Chart title
    title_text = spec.get("title")
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(title_text)
    else:
        chart.has_title = False

    # Legend
    show_legend = spec.get("show_legend")
    if show_legend is None:
        # Default: show for multi-series, hide for single-series
        show_legend = len(spec.get("series", [])) > 1
    chart.has_legend = bool(show_legend)

    # Value format on data labels — apply if specified
    value_format = spec.get("value_format")
    if value_format:
        try:
            chart.plots[0].has_data_labels = True
            chart.plots[0].data_labels.number_format = str(value_format)
            chart.plots[0].data_labels.number_format_is_linked = False
        except Exception as exc:
            log.warning(
                "Could not apply value_format %r to chart: %s", value_format, exc
            )

    # Axis titles (category / value axes — not available on pie/doughnut)
    x_title = spec.get("x_axis_title")
    y_title = spec.get("y_axis_title")
    chart_type = spec.get("type", "")
    if chart_type not in ("pie", "doughnut"):
        if x_title:
            try:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = str(x_title)
            except Exception as exc:
                log.warning("Could not set x_axis_title: %s", exc)
        if y_title:
            try:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = str(y_title)
            except Exception as exc:
                log.warning("Could not set y_axis_title: %s", exc)


def apply_chart_enrichment(
    prs,
    slide_index_1based: int,
    marker_name: str,
    chart_spec: dict[str, Any],
    brand_palette: BrandPalette | None = None,
) -> None:
    """Replace the named CHART marker shape with a native python-pptx chart.

    Steps:
      1. Validate chart_spec — raise EnrichmentApplyError on missing required fields.
      2. Find the marker shape and capture its coordinates.
      3. Drop the marker shape element and any sibling addText label.
      4. Build CategoryChartData (or XyChartData for scatter).
      5. Derive series colours from brand_palette (or default).
      6. Insert chart via slide.shapes.add_chart().
      7. Apply series colours and optional formatting.

    Parameters
    ----------
    prs:
        python-pptx Presentation object (in-memory, caller owns save).
    slide_index_1based:
        1-based slide index.
    marker_name:
        Full marker string, e.g. ``"CHART:revenue-vs-costs"``.
    chart_spec:
        Dict with required fields: type, categories, series. Optional:
        title, x_axis_title, y_axis_title, show_legend, value_format,
        palette_role, series_colours.
    brand_palette:
        BrandPalette for series colour derivation. When None, falls back to
        a neutral grey-blue palette derived from a default primary colour.
    """
    _validate_chart_spec(chart_spec, marker_name)

    shape, slide = _find_marker_shape(prs, slide_index_1based, marker_name)
    if shape is None:
        raise LookupError(
            f"CHART marker shape {marker_name!r} not found on slide "
            f"{slide_index_1based}."
        )

    # Capture marker geometry before removing the shape.
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Remove the marker placeholder rect.
    sp_elem = shape._element
    sp_elem.getparent().remove(sp_elem)

    # Remove any sibling addText label whose visible text equals marker_name
    # (the v0.2 #23 cleanup pattern, generalised from BG/SMARTART to CHART).
    _drop_residual_marker_label_shapes_local(prs, slide_index_1based, marker_name)

    # Build chart data.
    chart_type_key = chart_spec["type"]
    chart_type = _CHART_TYPE_MAP[chart_type_key]
    categories = chart_spec["categories"]
    series = chart_spec["series"]
    n_series = len(series)

    if chart_type_key == "scatter":
        chart_data = XyChartData()
        for s in series:
            series_data = chart_data.add_series(str(s["name"]))
            vals = s["values"]
            # Scatter expects (x, y) pairs. If values is a flat list of
            # scalars, pair them with categories (converted to float if possible).
            if vals and not isinstance(vals[0], (list, tuple)):
                for cat, val in zip(categories, vals):
                    try:
                        x_val = float(cat)
                    except (TypeError, ValueError):
                        x_val = float(categories.index(cat))
                    series_data.add_data_point(x_val, float(val) if val is not None else 0.0)
            else:
                for pair in vals:
                    series_data.add_data_point(float(pair[0]), float(pair[1]))
    else:
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]
        for s in series:
            chart_data.add_series(
                str(s["name"]),
                [v if v is not None else 0 for v in s["values"]],
            )

    # Derive series colours.
    palette_role = chart_spec.get("palette_role", _DEFAULT_PALETTE_ROLE)
    explicit_colours = chart_spec.get("series_colours")

    if brand_palette is not None:
        hex_colours = derive_chart_palette(
            brand=brand_palette,
            role=palette_role,
            n_series=n_series,
            explicit=explicit_colours,
        )
    elif explicit_colours:
        from src.chart_palette import _normalise_hex
        normalised = [_normalise_hex(c) for c in explicit_colours]
        if len(normalised) < n_series:
            normalised = normalised + [normalised[-1]] * (n_series - len(normalised))
        hex_colours = normalised[:n_series]
    else:
        # Fallback: neutral slate-blue family when no palette is available.
        hex_colours = ["4472C4", "5B9BD5", "70AD47", "ED7D31", "FFC000", "FF0000"][:n_series]
        if n_series > len(hex_colours):
            hex_colours = (hex_colours * (n_series // len(hex_colours) + 1))[:n_series]

    # Insert native chart at marker coordinates.
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart

    # Apply series colours.
    _apply_series_colours(chart, hex_colours)

    # Apply optional formatting (title, legend, axis titles, value format).
    _apply_chart_formatting(chart, chart_spec)

    log.info(
        "apply_chart_enrichment: inserted %s chart at slide %d (%s)",
        chart_type_key, slide_index_1based, marker_name,
    )


def _drop_residual_marker_label_shapes_local(
    prs, slide_index_1based: int, marker_name: str,
) -> None:
    """Drop addText label shapes whose visible text equals marker_name.

    This is the same logic as the top-level
    ``_drop_residual_marker_label_shapes`` in enrichment.py, duplicated here
    to avoid a circular import (enrichment_ops.chart → enrichment would
    cycle). The two implementations must stay in sync.

    The marker shape itself has already been removed by the time this is
    called, so the cNvPr-name-mismatch guard is a belt-and-braces check
    for any edge-case where name re-use occurs.
    """
    a_t_qn = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    slide = prs.slides[slide_index_1based - 1]
    cSld = slide.element.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return
    for sp in list(spTree):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        shape_name = cNvPr.get("name") if cNvPr is not None else None
        if shape_name == marker_name:
            continue
        parts: list[str] = []
        for t_elem in sp.iter(a_t_qn):
            if t_elem.text:
                parts.append(t_elem.text)
        visible_text = "".join(parts).strip()
        if visible_text == marker_name:
            spTree.remove(sp)
