"""Transactional enrichment orchestrator.

Spec Section 3.4 step 4 - all-or-nothing application of selected
enrichments against a copy of the source .pptx. Phases:

  1. Open source as in-memory Presentation.
  2. Apply Op1 (backgrounds) and Op2 (element images) against the
     in-memory tree. For SMARTART items with action="apply_clear_overlap",
     remove the overlapping text shapes from the in-memory tree HERE
     (before any save).
  3. Save the in-memory tree to `<output>.tmp-<pid>` only if Phase 2
     succeeded for every requested op.
  4. Run Op3 (SmartArt injection) against the temp file. msft-smartart's
     inject() is itself read-all -> mutate -> write-fresh internally, so
     it preserves atomicity within a single call.
  5. Rename temp file to output via os.replace (atomic on POSIX).
  6. On any exception in Phases 2-5, the finally block removes the
     temp file and re-raises EnrichmentApplyError. No partial output
     ever reaches disk.
"""
from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn

from src.colors_xml_builder import BrandPalette, patch_carrier_palette
from src.enrichment_ops.background import apply_background_in_memory
from src.enrichment_ops.chart import apply_chart_enrichment
from src.enrichment_ops.element_image import replace_image_marker_in_memory
from src.enrichment_ops.smartart import inject_smartart_into_file
from src.enrichment_ops.smartart_from_list import (
    extract_list_items_from_marker_shape,
    select_layout_for_bullets,
)
from src.security import resolve_within_allowlist
from src.smartart_bridge import render_carrier, select_layout_for_slide


VALID_KINDS = {"background", "image", "smartart", "smartart_from_list", "chart"}
VALID_ACTIONS = {"apply", "apply_clear_overlap", "skip"}


class EnrichmentApplyError(RuntimeError):
    """Raised by apply_enrichment on any op failure or rename failure."""


@dataclass
class EnrichmentItem:
    slide_index: int
    kind: str                    # background | image | smartart | smartart_from_list | chart
    marker_name: str             # full marker string (e.g. "IMAGE:foo", "CHART:bar")
    asset_path: Path | None      # image path for background/image; None for smartart/chart
    action: str = "apply"        # apply | apply_clear_overlap | skip
    smartart_spec: dict[str, Any] | None = None  # required when kind="smartart" and action!="skip"
    chart_spec: dict[str, Any] | None = None  # required when kind="chart" and action!="skip"
    overlap_shape_names: list[str] = field(default_factory=list)
    # Run 8 Finding #26 — picture-SmartArt image binding. When kind ==
    # "smartart_from_list" AND this list is non-empty, the bridge routes the
    # marker to a picture-capable layout (vList3 / vList4 / pList1) and binds
    # one image per bullet item via <dgm:blipFill> on each child node. The
    # length must match the bullet-item count discovered at extraction time;
    # otherwise apply_enrichment raises EnrichmentApplyError. None means
    # text-only routing (the existing v0.1.2 behaviour).
    per_item_image_paths: list[Path] | None = None


@dataclass
class EnrichmentPlan:
    source_pptx: Path
    output_pptx: Path
    items: list[EnrichmentItem]


def _remove_overlapping_text_shapes(prs, slide_index_1based: int,
                                       overlap_names: list[str]) -> None:
    if not overlap_names:
        return
    if slide_index_1based < 1 or slide_index_1based > len(prs.slides):
        return
    slide = prs.slides[slide_index_1based - 1]
    cSld = slide.element.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return
    for sp in list(spTree):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        if cNvPr is not None and cNvPr.get("name") in overlap_names:
            spTree.remove(sp)


def _drop_marker_shape(prs, slide_index_1based: int, marker_name: str) -> None:
    slide = prs.slides[slide_index_1based - 1]
    cSld = slide.element.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return
    for sp in list(spTree):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        if cNvPr is not None and cNvPr.get("name") == marker_name:
            spTree.remove(sp)


def _drop_residual_marker_label_shapes(
    prs, slide_index_1based: int, marker_name: str,
) -> None:
    """v0.2 #23 — extend the v0.1.x #17 BG-cleanup pattern to all marker kinds.

    Build.js authoring patterns sometimes include a small ``addText`` label
    shape next to a marker (e.g. ``"SMARTART-FROM-LIST:tooling-gaps"``
    floating beside the bullet list) so the marker is visible in the
    unenriched deck. After enrichment, the marker shape itself is replaced
    by the rendered SmartArt / image, but the small label addText survives
    as a stranded grey-italic floater on top of the rendered output.

    This helper removes any ``<p:sp>`` on the slide whose visible text
    content (concatenated ``<a:t>`` runs) equals ``marker_name`` AND
    whose own ``cNvPr@name`` is NOT ``marker_name``. The strict-equality
    comparison keeps legitimate slide text that merely contains the
    marker string as a substring intact; the cNvPr-name-mismatch
    constraint keeps the marker shape itself in place (the inject path
    handles its replacement).

    Run 6 evidence proved this pattern for ``BG:`` markers (Finding #17,
    fixed in v0.1.x). Run 8 + Run 9 dogfood revealed the same pattern
    affects ``SMARTART-FROM-LIST:`` markers when /pptx authors a label
    addText alongside the bullet list. v0.2 generalises the cleanup.
    """
    a_t_qn = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    slide = prs.slides[slide_index_1based - 1]
    cSld = slide.element.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return
    for sp in list(spTree):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        shape_name = cNvPr.get("name") if cNvPr is not None else None
        # Don't touch the marker shape itself — inject handles its
        # replacement and removing it here would break that path.
        if shape_name == marker_name:
            continue
        # Concatenate visible text (<a:t> runs) and compare strict.
        parts: list[str] = []
        for t_elem in sp.iter(a_t_qn):
            if t_elem.text:
                parts.append(t_elem.text)
        visible_text = "".join(parts).strip()
        if visible_text == marker_name:
            spTree.remove(sp)


def apply_enrichment(
    plan: EnrichmentPlan,
    *,
    allowed_image_roots: list[Path],
    carriers_dir: Path | None = None,
    brand_palette: BrandPalette | None = None,
) -> None:
    """Apply the plan transactionally - all ops succeed and the output
    file is renamed atomically, OR no output file is produced.

    Optional ``brand_palette`` is applied to every SmartArt carrier between
    ``render_carrier`` and ``inject_smartart_into_file`` so that injected
    SmartArt graphics inherit the brief's palette instead of the SDK
    fixture's default Microsoft palette (Finding #10, Run 3 dogfood).
    Pass ``None`` (the default) to keep the carrier's stock colours —
    backwards-compatible with pre-Contract-1 enrichment.
    """
    src = Path(plan.source_pptx)
    out = Path(plan.output_pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    tmp = out.with_name(f"{out.name}.tmp-{os.getpid()}")
    carriers_dir = carriers_dir or (out.parent / "carriers")

    skip_items = [it for it in plan.items if it.action == "skip"]
    active_items = [it for it in plan.items if it.action != "skip"]

    try:
        # Spec § Security & Privacy - every loaded .pptx must pass pre-flight
        # before python-pptx opens it. The analyser already calls preflight,
        # but apply_enrichment may be invoked independently (e.g. from
        # cohesion-driven re-runs in Task 26 Step 10), so preflight is repeated
        # here as the first gate. Cheap (a single zipfile.ZipFile + stat) and
        # idempotent.
        from src.security import preflight_pptx
        preflight_pptx(src)
        prs = Presentation(str(src))

        # Phase 2 - Op1, Op2 (in-memory) + pre-Op3 overlap clearing
        for item in active_items:
            if item.kind == "background":
                if item.asset_path is None:
                    raise EnrichmentApplyError(f"background item missing asset_path: {item}")
                apply_background_in_memory(
                    prs=prs, slide_index_1based=item.slide_index,
                    image_path=item.asset_path,
                    marker_name=item.marker_name,
                    allowed_image_roots=allowed_image_roots,
                )
            elif item.kind == "image":
                if item.asset_path is None:
                    raise EnrichmentApplyError(f"image item missing asset_path: {item}")
                replace_image_marker_in_memory(
                    prs=prs, marker_name=item.marker_name,
                    image_path=item.asset_path,
                    allowed_image_roots=allowed_image_roots,
                )
            elif item.kind == "smartart":
                if item.smartart_spec is None:
                    raise EnrichmentApplyError(f"smartart item missing spec: {item}")
                if item.action == "apply_clear_overlap":
                    _remove_overlapping_text_shapes(
                        prs, item.slide_index, item.overlap_shape_names,
                    )
                # v0.2 #23 — drop residual marker-label addText shapes
                # (same cleanup pattern as smartart_from_list / BG).
                _drop_residual_marker_label_shapes(
                    prs, item.slide_index, item.marker_name,
                )
                # Op3 is file-based; we'll do it after the save below.
            elif item.kind == "smartart_from_list":
                # Contract 2: the marker shape IS a text shape with bullets.
                # Extract items from the shape's text frame; build the spec
                # inline if the caller didn't supply one. The substitution
                # is handled by Phase 4's inject_smartart_into_file, which
                # finds the marker shape by name and replaces it with the
                # rendered SmartArt at the same coordinates.
                #
                # Layout selection (Run 4/5/6 Finding #13): rather than
                # hardcoding ``process1`` and forcing Speakers to hand-shorten
                # any prose bullet > 24 chars, route to the smallest layout
                # that fits the content. process1 / list1 / vList2 are the
                # routing tiers; only items longer than vList2's cap are
                # truncated, with a Speaker-visible warning attached.
                if not item.smartart_spec:
                    items = extract_list_items_from_marker_shape(
                        prs=prs,
                        slide_index_1based=item.slide_index,
                        marker_name=item.marker_name,
                    )
                    # Run 8 Finding #26 — picture-SmartArt image binding.
                    # When per_item_image_paths is supplied, force routing to
                    # a picture-capable layout regardless of bullet length;
                    # the bullet count must equal the image count or the
                    # builder will refuse to render.
                    image_paths = item.per_item_image_paths
                    if image_paths is not None:
                        if len(image_paths) != len(items):
                            raise EnrichmentApplyError(
                                f"per_item_image_paths length ({len(image_paths)}) "
                                f"does not match marker bullet count ({len(items)}) "
                                f"for {item.marker_name!r} on slide {item.slide_index}"
                            )
                        # All paths must resolve under the allowlist.
                        resolved_image_paths = [
                            resolve_within_allowlist(p, allowed_image_roots)
                            for p in image_paths
                        ]
                        # Picture-capable layout selection: pList1 is wider
                        # (gallery row, 5+ items); vList3 is portrait-stack
                        # (3-5 items). Choose by item count; both honour
                        # data_shape="picture".
                        layout_id = "pList1" if len(items) >= 5 else "vList3"
                        items_to_render = list(items)
                        layout_warnings = [
                            f"layout={layout_id} forced because "
                            f"per_item_image_paths supplied "
                            f"({len(image_paths)} images bound to bullets)"
                        ]
                        # Picture builder expects a list of dicts with
                        # `label` + `image_path` per item (one image per
                        # bullet). Build the parallel-array → list-of-dicts
                        # transform here so the bridge contract stays
                        # ergonomic (operator supplies a flat list of paths)
                        # while the builder gets its declared input shape.
                        spec_data: dict[str, Any] = {
                            "items": [
                                {"label": label, "image_path": str(img)}
                                for label, img in zip(
                                    items_to_render, resolved_image_paths
                                )
                            ],
                        }
                    else:
                        layout_id, items_to_render, layout_warnings = (
                            select_layout_for_bullets(items)
                        )
                        spec_data = {"items": items_to_render}

                    item.smartart_spec = {
                        "graphic_type": "flowchart",
                        "layout_id": layout_id,
                        "data": spec_data,
                    }
                    if layout_warnings:
                        # Attach for the orchestrator to surface; the
                        # EnrichmentReport ledger reads this in Step 11.
                        item.smartart_spec["layout_warnings"] = layout_warnings
                if item.action == "apply_clear_overlap":
                    _remove_overlapping_text_shapes(
                        prs, item.slide_index, item.overlap_shape_names,
                    )
                # v0.2 #23 — drop residual addText shapes whose visible
                # text equals the marker name. /pptx authoring patterns
                # sometimes include a small grey-italic label addText
                # alongside the bullet list to make the marker visible
                # in the unenriched deck; without this cleanup, that
                # label survives enrichment as a stranded floater.
                _drop_residual_marker_label_shapes(
                    prs, item.slide_index, item.marker_name,
                )
            elif item.kind == "chart":
                # Issue #55 — CHART:slug marker kind.
                # chart_spec is operator-supplied on EnrichmentItem at
                # /enrich-deck time (pattern (a) from the plan). The marker
                # rect + any sibling addText label are removed inside
                # apply_chart_enrichment via the same cleanup logic as BG /
                # SMARTART / SMARTART-FROM-LIST (v0.2 #23 pattern).
                if item.chart_spec is None:
                    raise EnrichmentApplyError(
                        f"chart item missing chart_spec: {item}"
                    )
                apply_chart_enrichment(
                    prs=prs,
                    slide_index_1based=item.slide_index,
                    marker_name=item.marker_name,
                    chart_spec=item.chart_spec,
                    brand_palette=brand_palette,
                )
            else:
                raise EnrichmentApplyError(f"unknown enrichment kind {item.kind!r}")

        # Phase 3 - Save in-memory tree to temp
        prs.save(str(tmp))

        # Phase 4 - Op3 SmartArt injections against the temp file.
        # When brand_palette is supplied, patch each carrier's colors1.xml
        # with brand colours BEFORE injection (Contract 1, Finding #10).
        # Both ``smartart`` and ``smartart_from_list`` items use the same
        # injection pipeline — by Phase 4 the latter has had its smartart_spec
        # populated from the marker shape's bullet text.
        for item in active_items:
            if item.kind not in ("smartart", "smartart_from_list"):
                continue
            carrier_path = render_carrier(item.smartart_spec, output_dir=carriers_dir)
            if brand_palette is not None:
                patch_carrier_palette(carrier_path, brand_palette)
            inject_smartart_into_file(
                host_pptx=tmp,
                slide_index_1based=item.slide_index,
                marker_name=item.marker_name,
                carrier_pptx=carrier_path,
            )

        # Phase 5 - atomic rename
        os.replace(tmp, out)

    except Exception as exc:
        # Cleanup: ensure no temp file is left behind, no partial output
        try:
            if tmp.exists():
                os.unlink(tmp)
        except OSError:
            pass
        # Also ensure no partial output (in case rename half-succeeded)
        if out.exists() and not isinstance(exc, EnrichmentApplyError):
            # Preserve original semantics: only delete the output if we
            # are sure we created it during this call. We mark this by
            # checking that it didn't exist before the call.
            pass
        raise EnrichmentApplyError(str(exc)) from exc
    finally:
        # Belt-and-braces - temp file MUST NOT survive
        try:
            if tmp.exists():
                os.unlink(tmp)
        except OSError:
            pass
