"""OOXML-primary analyser via python-pptx.

Lifted from Spike 3's parsers/pptx_parser.py with two production-bound
extensions:

  1. bgColor attribute on <p:sld> → background_kind='solid' (Spike 3
     spec update item #7 — covers PptxGenJS colour-only backgrounds).
  2. FileNotFoundError instead of swallowing missing-input.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from src.placeholder import parse_marker
from src.slide_facts import Marker, SlideFacts


def _shape_element_type(shape) -> str:
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "text"
    return "shape"


def _background_kind(slide) -> str:
    """Detect background category, considering both <p:bg> and the slide
    element's bgColor attribute (PptxGenJS emits colour backgrounds via
    bgColor, not <p:bg>)."""
    cSld = slide.element.find(qn("p:cSld"))
    bg = cSld.find(qn("p:bg")) if cSld is not None else None
    if bg is not None:
        if bg.find(".//" + qn("a:blipFill")) is not None:
            return "image"
        if bg.find(".//" + qn("a:solidFill")) is not None:
            return "solid"
        if bg.find(".//" + qn("a:gradFill")) is not None:
            return "solid"
        return "unknown"
    # Spike 3 spec update item #7 — check bgColor attribute on <p:sld>
    if slide.element.get("bgColor"):
        return "solid"
    return "default"


def parse_pptx(path: Path | str) -> list[SlideFacts]:
    """Parse a .pptx → list[SlideFacts] using OOXML as the source of truth.

    Raises FileNotFoundError if `path` does not exist.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(p)
    prs = Presentation(str(p))
    results: list[SlideFacts] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        markers: list[Marker] = []
        counts: dict[str, int] = {"text": 0, "shape": 0, "image": 0,
                                   "chart": 0, "table": 0}
        for shape in slide.shapes:
            name = getattr(shape, "name", "") or ""
            parsed = parse_marker(name)
            if parsed is not None:
                kind, ident = parsed
                markers.append(Marker(
                    kind=kind,
                    identifier=ident,
                    left_emu=shape.left or 0,
                    top_emu=shape.top or 0,
                    width_emu=shape.width or 0,
                    height_emu=shape.height or 0,
                ))
                continue
            kind = _shape_element_type(shape)
            counts[kind] = counts.get(kind, 0) + 1
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        results.append(SlideFacts(
            slide_index=idx,
            text_content="\n".join(t for t in text_parts if t.strip()),
            markers=markers,
            background_kind=_background_kind(slide),
            element_types=counts,
        ))
    return results
