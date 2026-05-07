"""SMARTART marker-adjacent text overlap detection.

For every slide carrying a SMARTART:* marker, check whether any other
shape's bounding box on that slide intersects the marker's geometry.
Spec Section 3.1 — converts the unverifiable brief-side promise into
an analyser-side check.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from src.placeholder import parse_marker
from src.slide_facts import OverlapWarning


def _bbox(shape) -> tuple[int, int, int, int]:
    """Return (left, top, right, bottom) in EMU."""
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return left, top, left + width, top + height


def _intersects(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    al, at, ar, ab = a
    bl, bt, br, bb = b
    return not (ar <= bl or br <= al or ab <= bt or bb <= at)


def find_overlaps(pptx_path: Path | str) -> list[OverlapWarning]:
    """Return one OverlapWarning per slide whose SMARTART marker is overlapped
    by at least one other text-bearing shape."""
    prs = Presentation(str(pptx_path))
    warnings: list[OverlapWarning] = []
    for idx, slide in enumerate(prs.slides, start=1):
        smartart_markers: list[tuple[str, tuple[int, int, int, int]]] = []
        other_shapes: list[tuple[str, tuple[int, int, int, int]]] = []
        for shape in slide.shapes:
            name = getattr(shape, "name", "") or ""
            parsed = parse_marker(name)
            if parsed is not None and parsed[0] == "SMARTART":
                smartart_markers.append((f"{parsed[0]}:{parsed[1]}", _bbox(shape)))
                continue
            # Only consider shapes that bear text — empty rects don't bleed
            if shape.has_text_frame and shape.text_frame.text.strip():
                other_shapes.append((name, _bbox(shape)))
        for marker_id, marker_box in smartart_markers:
            overlap_names = sorted(
                name for name, box in other_shapes if _intersects(marker_box, box)
            )
            if overlap_names:
                warnings.append(OverlapWarning(
                    slide_index=idx,
                    marker_id=marker_id,
                    overlapping_shape_names=overlap_names,
                ))
    return warnings
