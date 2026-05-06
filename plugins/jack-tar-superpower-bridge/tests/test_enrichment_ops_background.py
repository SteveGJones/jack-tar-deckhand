import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from src.enrichment_ops.background import apply_background_in_memory
from src.security import AllowedPathError


def test_applies_blipfill_and_removes_marker(tmp_path, seed_variant_a, placeholder_png):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    img_dir = tmp_path / "generated"
    img_dir.mkdir()
    img = img_dir / "bg.png"
    shutil.copy(placeholder_png, img)

    prs = Presentation(str(work))
    apply_background_in_memory(
        prs=prs,
        slide_index_1based=1,
        image_path=img,
        marker_name="BG:title-hero-dark-grid",
        allowed_image_roots=[img_dir],
    )
    # Assert the in-memory tree carries a <p:bg> with a <a:blip>
    slide = prs.slides[0]
    cSld = slide.element.find(qn("p:cSld"))
    bg = cSld.find(qn("p:bg"))
    assert bg is not None
    blip = bg.find(".//" + qn("a:blip"))
    assert blip is not None
    # Marker shape removed
    names = [s.name for s in slide.shapes]
    assert "BG:title-hero-dark-grid" not in names


def test_image_path_outside_allowlist_raises(tmp_path, seed_variant_a, placeholder_png):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    rogue = tmp_path / "rogue.png"
    shutil.copy(placeholder_png, rogue)

    prs = Presentation(str(work))
    allowed = tmp_path / "generated"
    allowed.mkdir()
    with pytest.raises(AllowedPathError):
        apply_background_in_memory(
            prs=prs, slide_index_1based=1, image_path=rogue,
            marker_name=None, allowed_image_roots=[allowed],
        )


def test_idempotent_when_no_existing_bg(tmp_path, seed_variant_a, placeholder_png):
    """Calling apply_background twice on the same slide replaces, doesn't stack."""
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(placeholder_png, img)

    prs = Presentation(str(work))
    apply_background_in_memory(prs=prs, slide_index_1based=1, image_path=img,
                                marker_name=None, allowed_image_roots=[img_dir])
    apply_background_in_memory(prs=prs, slide_index_1based=1, image_path=img,
                                marker_name=None, allowed_image_roots=[img_dir])
    cSld = prs.slides[0].element.find(qn("p:cSld"))
    bgs = cSld.findall(qn("p:bg"))
    assert len(bgs) == 1


def test_removes_addtext_label_matching_marker_name(tmp_path, placeholder_png):
    """Run 5 Finding #17: /pptx authors a BG marker as a placeholder rect AND a
    separate ``addText`` shape rendering "BG:slug" as visible text inside the
    rect (so authors can see what the marker is). The bridge swaps the rect
    for the AI background, but the ``addText`` shape sits on top of the new
    image as residual cosmetic text. The bridge must auto-remove any text
    shape whose text content matches the marker name.

    The narrative-brief-architect agent now warns /pptx away from authoring
    the addText label at all (build.js authoring caveat); this test guards
    against decks authored before that warning was in effect.
    """
    from pptx import Presentation as NewPresentation
    from pptx.util import Inches

    work = tmp_path / "deck.pptx"
    prs = NewPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # The BG marker rect — same shape /pptx normally emits via addShape().
    marker_name = "BG:pivot-moment"
    rect = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    rect.name = marker_name

    # The cosmetic addText label sitting on top of the rect — exactly the
    # Run 5 anti-pattern. python-pptx does not let us tag a shape's NAME
    # as something different from the OOXML cNvPr; what /pptx authors is a
    # second shape whose visible TEXT is the marker string. The bridge
    # detects by text content.
    label_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.5))
    label_shape.text_frame.text = marker_name
    label_shape.name = "addText-residual"

    # An unrelated text shape that must NOT be removed
    keep_shape = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
    keep_shape.text_frame.text = "From opportunity. To decision."
    keep_shape.name = "pivot-headline"

    prs.save(str(work))

    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(placeholder_png, img)

    prs = Presentation(str(work))
    apply_background_in_memory(
        prs=prs, slide_index_1based=1, image_path=img,
        marker_name=marker_name, allowed_image_roots=[img_dir],
    )

    slide = prs.slides[0]
    shape_texts = [
        sh.text_frame.text for sh in slide.shapes
        if sh.has_text_frame
    ]
    shape_names = [sh.name for sh in slide.shapes]

    # Marker rect gone
    assert marker_name not in shape_names
    # addText residual matching the marker text gone
    assert marker_name not in shape_texts
    # Unrelated text preserved
    assert "From opportunity. To decision." in shape_texts
