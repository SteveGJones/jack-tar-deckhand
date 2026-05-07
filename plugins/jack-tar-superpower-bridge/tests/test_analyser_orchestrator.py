import pytest
from pathlib import Path

from src.analyser import analyse_pptx, AnalyserOptions
from src.security import PptxPreflightError
from src.slide_facts import AnalyserResult


def test_variant_a_uses_ooxml_only(seed_variant_a):
    result = analyse_pptx(seed_variant_a)
    assert isinstance(result, AnalyserResult)
    assert result.total_slides == 10
    assert result.total_markers == 8
    assert result.js_fallback_used is False


def test_variant_b_falls_back_to_js_when_buildjs_present(seed_variant_b):
    """Variant B's build.js sits alongside the .pptx; JS fallback should fire."""
    result = analyse_pptx(seed_variant_b)
    assert result.js_fallback_used is True
    assert result.total_markers >= 1


def test_no_buildjs_does_not_fall_back(seed_no_buildjs):
    result = analyse_pptx(seed_no_buildjs)
    # Control case has 1 named marker shape; OOXML finds it, no fallback needed
    assert result.js_fallback_used is False


def test_duplicate_marker_ids_surfaced(tmp_path):
    """Two slides both carrying IMAGE:foo must produce a duplicate warning."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        m = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(2))
        m.name = "IMAGE:foo"
    out = tmp_path / "dupes.pptx"
    prs.save(out)
    result = analyse_pptx(out)
    assert "IMAGE:foo" in result.duplicate_marker_ids


def test_overlap_warnings_collected_from_verifier(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    marker = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(5), Inches(3))
    marker.name = "SMARTART:foo"
    body = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(0.5))
    body.text_frame.text = "x"
    out = tmp_path / "overlap.pptx"
    prs.save(out)
    result = analyse_pptx(out)
    assert len(result.overlap_warnings) == 1


def test_preflight_failure_propagates(tmp_path):
    """A .pptx that fails preflight must raise PptxPreflightError, not be silently passed."""
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a zip")
    with pytest.raises(PptxPreflightError):
        analyse_pptx(bad)


def test_analyser_options_disable_js_fallback(seed_variant_b):
    result = analyse_pptx(seed_variant_b, options=AnalyserOptions(enable_js_fallback=False))
    assert result.js_fallback_used is False


def test_slide_count_mismatch_is_logged_in_notes(tmp_path):
    """If the JS fallback produces a different number of slides than OOXML,
    the analyser must log the discrepancy in notes (never silent)."""
    # Build a 2-slide deck with no markers in OOXML (triggers JS fallback).
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for _ in range(2):
        prs.slides.add_slide(prs.slide_layouts[6])
    out = tmp_path / "two-slide.pptx"
    prs.save(out)

    # Write a build.js that describes 3 slides (mismatch vs the .pptx's 2 slides).
    build_js = tmp_path / "build.js"
    build_js.write_text(
        "const pres = new pptxgen();\n"
        "const s1 = pres.addSlide();\n"
        "s1.addShape('rect', { objectName: 'IMAGE:a', x:1, y:1, w:1, h:1 });\n"
        "const s2 = pres.addSlide();\n"
        "s2.addShape('rect', { objectName: 'IMAGE:b', x:1, y:1, w:1, h:1 });\n"
        "const s3 = pres.addSlide();\n"
        "s3.addShape('rect', { objectName: 'IMAGE:c', x:1, y:1, w:1, h:1 });\n"
    )

    result = analyse_pptx(out)
    assert result.js_fallback_used is True
    # Notes must include the mismatch diagnostic
    assert any("produced 3 slides vs OOXML 2" in n for n in result.notes), (
        f"no slide-count mismatch note in {result.notes}"
    )
