"""Tests for colors_xml_builder — Contract 1 (Finding #10) brand-palette injection.

Caught during Run 3 dogfood (2026-04-26): SmartArt carriers ship with default
Microsoft palette regardless of brief. These tests verify:

- BrandPalette validates and normalises hex inputs.
- apply_brand_palette_to_colors_xml transforms schemeClr → srgbClr while
  preserving structure (tints, shades, alphas, all styleLbls).
- patch_carrier_palette round-trips a real carrier .pptx, atomically replacing
  ppt/diagrams/colors1.xml without disturbing other parts.
- derive_palette_from_brief_text picks up Section-B-style palette tables.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from lxml import etree

from src.colors_xml_builder import (
    BrandPalette,
    apply_brand_palette_to_colors_xml,
    derive_palette_from_brief_text,
    patch_carrier_palette,
)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


# --- BrandPalette validation -------------------------------------------------

def test_brand_palette_normalises_hex_input():
    p = BrandPalette(
        primary_fill_hex="#14213d",
        text_on_primary_hex="e8dec2",
        text_on_surface_hex=" #1B2333 ",
    )
    # Strips #, uppercases, trims whitespace
    assert p.primary_fill_hex == "14213D"
    assert p.text_on_primary_hex == "E8DEC2"
    assert p.text_on_surface_hex == "1B2333"


def test_brand_palette_rejects_invalid_hex():
    with pytest.raises(ValueError, match="primary_fill_hex"):
        BrandPalette(primary_fill_hex="zzzzzz",
                      text_on_primary_hex="000000",
                      text_on_surface_hex="ffffff")
    with pytest.raises(ValueError, match="text_on_primary_hex"):
        BrandPalette(primary_fill_hex="000000",
                      text_on_primary_hex="abc",       # too short
                      text_on_surface_hex="ffffff")
    with pytest.raises(ValueError, match="text_on_surface_hex"):
        BrandPalette(primary_fill_hex="000000",
                      text_on_primary_hex="ffffff",
                      text_on_surface_hex="1234567")   # too long


def test_brand_palette_is_frozen():
    p = BrandPalette("000000", "ffffff", "808080")
    with pytest.raises(Exception):  # FrozenInstanceError or AttributeError
        p.primary_fill_hex = "111111"  # type: ignore[misc]


# --- apply_brand_palette_to_colors_xml --------------------------------------

def _wrap_in_styledef(inner_xml: str) -> bytes:
    """Wrap inner XML in a minimal valid colorsDef root."""
    return (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<dgm:colorsDef xmlns:dgm="{DGM_NS}" xmlns:a="{A_NS}" uniqueId="urn:test">'
        f'<dgm:title val=""/><dgm:desc val=""/>'
        f'<dgm:catLst><dgm:cat type="accent1" pri="11200"/></dgm:catLst>'
        f'{inner_xml}'
        f'</dgm:colorsDef>'
    ).encode("utf-8")


def _palette_blueprint() -> BrandPalette:
    return BrandPalette(
        primary_fill_hex="14213D",   # navy
        text_on_primary_hex="E8DEC2", # warm light
        text_on_surface_hex="1B2333", # body text on cream
    )


def test_apply_replaces_accent1_with_srgb():
    inner = (
        '<dgm:styleLbl name="node0">'
        '<dgm:fillClrLst meth="repeat"><a:schemeClr val="accent1"/></dgm:fillClrLst>'
        '<dgm:linClrLst meth="repeat"><a:schemeClr val="lt1"/></dgm:linClrLst>'
        '<dgm:effectClrLst/><dgm:txLinClrLst/>'
        '<dgm:txFillClrLst><a:schemeClr val="dk1"/></dgm:txFillClrLst>'
        '<dgm:txEffectClrLst/>'
        '</dgm:styleLbl>'
    )
    result = apply_brand_palette_to_colors_xml(_wrap_in_styledef(inner), _palette_blueprint())
    root = etree.fromstring(result)

    # Three colour elements should now be srgbClr with the right hex
    srgbs = root.findall(f".//{{{A_NS}}}srgbClr")
    assert len(srgbs) == 3
    assert {srgb.get("val") for srgb in srgbs} == {"14213D", "E8DEC2", "1B2333"}

    # No schemeClr remains for accent1/lt1/dk1
    schemes = root.findall(f".//{{{A_NS}}}schemeClr")
    assert all(s.get("val") not in {"accent1", "lt1", "dk1"} for s in schemes)


def test_apply_preserves_tint_and_shade_modifiers():
    """Critical: tints/shades/alphas applied to a colour are children of the
    schemeClr element. After tag swap the modifiers must still apply."""
    inner = (
        '<dgm:styleLbl name="node0">'
        '<dgm:fillClrLst meth="repeat">'
        '<a:schemeClr val="accent1"><a:tint val="50000"/><a:alpha val="80000"/></a:schemeClr>'
        '</dgm:fillClrLst>'
        '<dgm:linClrLst meth="repeat">'
        '<a:schemeClr val="accent1"><a:shade val="60000"/></a:schemeClr>'
        '</dgm:linClrLst>'
        '<dgm:effectClrLst/><dgm:txLinClrLst/><dgm:txFillClrLst/><dgm:txEffectClrLst/>'
        '</dgm:styleLbl>'
    )
    result = apply_brand_palette_to_colors_xml(_wrap_in_styledef(inner), _palette_blueprint())
    root = etree.fromstring(result)

    # Both schemeClr elements became srgbClr
    srgbs = root.findall(f".//{{{A_NS}}}srgbClr")
    assert len(srgbs) == 2

    # Tint, shade, alpha children all preserved
    assert root.find(f".//{{{A_NS}}}srgbClr/{{{A_NS}}}tint") is not None
    assert root.find(f".//{{{A_NS}}}srgbClr/{{{A_NS}}}alpha") is not None
    assert root.find(f".//{{{A_NS}}}srgbClr/{{{A_NS}}}shade") is not None


def test_apply_leaves_unmapped_scheme_refs_alone():
    """phClr, hlink, folHlink are not in the slot map — preserve them."""
    inner = (
        '<dgm:styleLbl name="node0">'
        '<dgm:fillClrLst meth="repeat"><a:schemeClr val="phClr"/></dgm:fillClrLst>'
        '<dgm:linClrLst meth="repeat"><a:schemeClr val="hlink"/></dgm:linClrLst>'
        '<dgm:effectClrLst/><dgm:txLinClrLst/><dgm:txFillClrLst/><dgm:txEffectClrLst/>'
        '</dgm:styleLbl>'
    )
    result = apply_brand_palette_to_colors_xml(_wrap_in_styledef(inner), _palette_blueprint())
    root = etree.fromstring(result)

    schemes = root.findall(f".//{{{A_NS}}}schemeClr")
    scheme_vals = {s.get("val") for s in schemes}
    assert "phClr" in scheme_vals
    assert "hlink" in scheme_vals

    # No unintended srgbClr injection
    assert root.find(f".//{{{A_NS}}}srgbClr") is None


def test_apply_handles_all_accent_slots():
    """accent2-6 should also map to primary_fill (not just accent1)."""
    inner = "".join(
        f'<dgm:styleLbl name="node{i}">'
        f'<dgm:fillClrLst meth="repeat"><a:schemeClr val="accent{i}"/></dgm:fillClrLst>'
        '<dgm:linClrLst/><dgm:effectClrLst/><dgm:txLinClrLst/>'
        '<dgm:txFillClrLst/><dgm:txEffectClrLst/>'
        '</dgm:styleLbl>'
        for i in range(1, 7)
    )
    result = apply_brand_palette_to_colors_xml(_wrap_in_styledef(inner), _palette_blueprint())
    root = etree.fromstring(result)
    srgbs = root.findall(f".//{{{A_NS}}}srgbClr")
    assert len(srgbs) == 6
    assert all(s.get("val") == "14213D" for s in srgbs)


def test_apply_preserves_dgm_styleLbl_structure():
    """Don't accidentally reorder or strip styleLbl wrappers / catLst / title."""
    inner = (
        '<dgm:styleLbl name="node0">'
        '<dgm:fillClrLst meth="repeat"><a:schemeClr val="accent1"/></dgm:fillClrLst>'
        '<dgm:linClrLst/><dgm:effectClrLst/><dgm:txLinClrLst/>'
        '<dgm:txFillClrLst/><dgm:txEffectClrLst/>'
        '</dgm:styleLbl>'
    )
    result = apply_brand_palette_to_colors_xml(_wrap_in_styledef(inner), _palette_blueprint())
    root = etree.fromstring(result)
    # Title, desc, catLst still present
    assert root.find(f"{{{DGM_NS}}}title") is not None
    assert root.find(f"{{{DGM_NS}}}desc") is not None
    assert root.find(f"{{{DGM_NS}}}catLst") is not None
    # styleLbl preserved with name attr
    style = root.find(f"{{{DGM_NS}}}styleLbl")
    assert style is not None
    assert style.get("name") == "node0"


# --- patch_carrier_palette: carrier round-trip ------------------------------

def _make_carrier_with_smartart_via_engine(tmp_path: Path) -> Path:
    """Render a real SmartArt carrier via the bridge's smartart_bridge facade."""
    from src.smartart_bridge import render_carrier
    spec = {
        "graphic_type": "flowchart",
        "layout_id": "process1",
        "data": {"items": ["Alpha", "Beta", "Gamma"]},
    }
    return render_carrier(spec, output_dir=tmp_path)


def test_patch_carrier_palette_round_trip(tmp_path):
    carrier = _make_carrier_with_smartart_via_engine(tmp_path)
    palette = _palette_blueprint()

    # Original colors.xml uses schemeClr — confirm before patching
    with zipfile.ZipFile(carrier) as z:
        original_colors = z.read("ppt/diagrams/colors1.xml")
    assert b"schemeClr" in original_colors
    assert b"14213D" not in original_colors  # not yet brand-mapped

    # Patch in place
    patch_carrier_palette(carrier, palette)

    # After patching: brand colours present, default scheme refs replaced
    with zipfile.ZipFile(carrier) as z:
        patched_colors = z.read("ppt/diagrams/colors1.xml")
        # All other parts still present
        names = set(z.namelist())
        assert "ppt/diagrams/layout1.xml" in names
        assert "ppt/diagrams/quickStyle1.xml" in names
        assert "ppt/diagrams/data1.xml" in names

    assert b"14213D" in patched_colors    # primary fill
    assert b"E8DEC2" in patched_colors    # text on primary
    assert b"1B2333" in patched_colors    # text on surface
    # accent1/lt1/dk1 schemeClr should be gone (other unmapped scheme refs may remain)
    root = etree.fromstring(patched_colors)
    schemes = root.findall(f".//{{{A_NS}}}schemeClr")
    scheme_vals = {s.get("val") for s in schemes}
    assert "accent1" not in scheme_vals
    assert "lt1" not in scheme_vals
    assert "dk1" not in scheme_vals


def test_patch_carrier_palette_preserves_other_parts(tmp_path):
    """Carrier should be byte-identical except for ppt/diagrams/colors1.xml."""
    carrier = _make_carrier_with_smartart_via_engine(tmp_path)

    with zipfile.ZipFile(carrier) as z:
        before = {name: z.read(name) for name in z.namelist()}

    patch_carrier_palette(carrier, _palette_blueprint())

    with zipfile.ZipFile(carrier) as z:
        after = {name: z.read(name) for name in z.namelist()}

    assert set(before.keys()) == set(after.keys())
    for name in before:
        if name == "ppt/diagrams/colors1.xml":
            assert before[name] != after[name]  # the only file that changed
        else:
            assert before[name] == after[name], f"{name} should be unchanged"


def test_patch_carrier_palette_raises_for_missing_carrier(tmp_path):
    palette = _palette_blueprint()
    with pytest.raises(FileNotFoundError):
        patch_carrier_palette(tmp_path / "missing.pptx", palette)


def test_patch_carrier_palette_atomic_on_failure(tmp_path):
    """If patching fails partway, the original carrier must remain intact.

    Triggers a real failure path: corrupt the carrier's colors1.xml to
    invalid XML so lxml.etree.fromstring raises during the transform.
    """
    carrier = _make_carrier_with_smartart_via_engine(tmp_path)

    # Replace colors1.xml with malformed XML to force a parse error
    import shutil
    malformed = b"<?xml version=\"1.0\"?><not<valid<xml/>"
    corrupted = tmp_path / "corrupt.pptx"
    with zipfile.ZipFile(carrier) as src:
        with zipfile.ZipFile(corrupted, "w", zipfile.ZIP_DEFLATED) as dst:
            for info in src.infolist():
                data = src.read(info.filename)
                if info.filename == "ppt/diagrams/colors1.xml":
                    data = malformed
                dst.writestr(info, data)

    original_size = corrupted.stat().st_size
    with zipfile.ZipFile(corrupted) as z:
        original_colors = z.read("ppt/diagrams/colors1.xml")
    assert original_colors == malformed

    # Patching should fail because lxml can't parse malformed XML
    with pytest.raises(Exception):  # lxml raises XMLSyntaxError
        patch_carrier_palette(corrupted, _palette_blueprint())

    # Corrupted carrier must be untouched (atomicity property)
    assert corrupted.exists()
    assert corrupted.stat().st_size == original_size
    with zipfile.ZipFile(corrupted) as z:
        assert z.read("ppt/diagrams/colors1.xml") == original_colors


# --- derive_palette_from_brief_text ------------------------------------------

def test_derive_palette_from_blueprint_brief_text():
    """Run 3-style Blueprint Retrospective brief should yield correct palette."""
    text = """
    The deck adopts a Blueprint Retrospective register.
    Surface (cream/parchment) #F5EFE0 — primary slide background.
    Structural (deep navy) #14213D — section headers.
    Body text on cream #1B2333 — primary body copy.
    Accent (muted rust) #B7410E — single emphasis colour.
    """
    palette = derive_palette_from_brief_text(text)
    assert palette is not None
    assert palette.primary_fill_hex == "14213D"
    assert palette.text_on_primary_hex == "F5EFE0"
    assert palette.text_on_surface_hex == "1B2333"


def test_derive_palette_returns_none_when_brief_lacks_hex():
    palette = derive_palette_from_brief_text(
        "A clean monochrome layout with restrained typography."
    )
    assert palette is None


def test_derive_palette_returns_none_when_role_keywords_missing():
    """If brief has hex values but no recognised role keywords for all 3 slots,
    return None — caller must construct palette explicitly."""
    palette = derive_palette_from_brief_text(
        "Something orange #FF8800 and something blue #0088FF."
    )
    assert palette is None


def test_derive_palette_from_markdown_table_format():
    """Run 3 dogfood (Finding #11) caught: brief writers emit palette tables
    with backtick-wrapped hex (`#XXXXXX`) inside markdown table rows. The
    heuristic must handle this format — it's the canonical brief output shape.
    """
    text = """
    Visual personality: Blueprint Retrospective.

    | Role | Hex | Usage |
    |---|---|---|
    | Surface (cream/parchment) | `#F5EFE0` | Primary slide background |
    | Structural (deep navy) | `#14213D` | Section headers, dividers |
    | Body text on cream | `#1B2333` | Primary body copy |
    | Accent (muted rust) | `#B7410E` | Single emphasis colour |
    """
    palette = derive_palette_from_brief_text(text)
    assert palette is not None
    assert palette.primary_fill_hex == "14213D"      # structural beats accent (tiered)
    assert palette.text_on_primary_hex == "F5EFE0"   # surface
    assert palette.text_on_surface_hex == "1B2333"   # body text


def test_derive_palette_structural_beats_accent_for_primary_fill():
    """Tier rule: when both 'structural' and 'accent' are named, primary_fill
    should be the structural colour (the deck's dominant fill), not the
    accent (which is an emphasis-only colour). Critical on light-surface
    decks like Blueprint Retrospective where SmartArt fill must contrast."""
    text = """
    Accent (rust) #B7410E for emphasis.
    Structural (navy) #14213D for headers.
    Surface (cream) #F5EFE0 background.
    Body text #1B2333 on cream.
    """
    palette = derive_palette_from_brief_text(text)
    assert palette is not None
    # Even though accent appears first in the text, structural takes priority
    # because of the tiered keyword matching.
    assert palette.primary_fill_hex == "14213D"
    assert palette.text_on_primary_hex == "F5EFE0"
    assert palette.text_on_surface_hex == "1B2333"


def test_derive_palette_accent_only_brief_derives_palette(tmp_path):
    """Run 4 Finding #12 — accent-only briefs (a single accent colour named
    alongside surface and body text, no explicit "structural") should still
    derive a usable palette. The accent fills the primary_fill slot via the
    tier-1 keyword fallback."""
    text = """
    Visual personality: minimal. One accent colour for emphasis.

    | Role | Hex | Usage |
    |---|---|---|
    | Surface | `#FAFAFA` | Background |
    | Accent | `#0F4C75` | Single emphasis colour for headers and SmartArt |
    | Body text | `#1B1B1B` | Primary body copy |
    """
    palette = derive_palette_from_brief_text(text)
    assert palette is not None, (
        "accent-only briefs must still produce a palette via tier-1 fallback"
    )
    # accent fills primary_fill_hex (tier-1)
    assert palette.primary_fill_hex == "0F4C75"
    assert palette.text_on_primary_hex == "FAFAFA"   # surface
    assert palette.text_on_surface_hex == "1B1B1B"   # body text


def test_derive_palette_brand_only_brief_derives_palette():
    """Some briefs use 'brand' or 'hero' rather than 'structural' / 'accent'
    for the dominant fill colour. The heuristic should pick those up too."""
    text = """
    | Role | Hex | Usage |
    |---|---|---|
    | Surface (cream) | `#F5EFE0` | Background |
    | Brand colour (navy) | `#14213D` | Headers, SmartArt fills |
    | Body text on cream | `#1B2333` | Primary body copy |
    """
    palette = derive_palette_from_brief_text(text)
    assert palette is not None
    assert palette.primary_fill_hex == "14213D"   # brand → primary
    assert palette.text_on_primary_hex == "F5EFE0"
    assert palette.text_on_surface_hex == "1B2333"


def test_derive_palette_picks_first_recognised_role_when_ambiguous():
    """If only one role keyword is recognised, the heuristic should still
    return None rather than fabricate values — partial palettes are less
    useful than no palette (which prompts the Speaker to fix the brief)."""
    text = "Just an accent: #FF0000."
    palette = derive_palette_from_brief_text(text)
    assert palette is None  # no surface, no body text → cannot fill all 3 slots


# --- Integration: apply_enrichment + brand_palette --------------------------

def test_apply_enrichment_with_brand_palette_patches_smartart_colours(
    tmp_path, seed_variant_a,
):
    """End-to-end: a SMARTART item enriched with a brand_palette must produce
    an output .pptx whose injected SmartArt colors1.xml contains brand hex
    values, not the SDK fixture's default schemeClr references.
    """
    from pptx import Presentation
    from src.enrichment import EnrichmentItem, EnrichmentPlan, apply_enrichment

    # Find a real SMARTART marker on the seed deck
    prs = Presentation(str(seed_variant_a))
    target_marker = None
    target_slide_idx = None
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name and shape.name.startswith("SMARTART:"):
                target_marker = shape.name
                target_slide_idx = idx
                break
        if target_marker:
            break
    assert target_marker is not None, "seed fixture must contain a SMARTART marker"

    out = tmp_path / "branded.pptx"
    plan = EnrichmentPlan(
        source_pptx=seed_variant_a,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=target_slide_idx,
                kind="smartart",
                marker_name=target_marker,
                asset_path=None,
                action="apply",
                smartart_spec={
                    "graphic_type": "flowchart",
                    "layout_id": "process1",
                    "data": {"items": ["one", "two", "three"]},
                },
            ),
        ],
    )
    palette = BrandPalette(
        primary_fill_hex="14213D",   # navy
        text_on_primary_hex="E8DEC2",
        text_on_surface_hex="1B2333",
    )
    apply_enrichment(
        plan,
        allowed_image_roots=[tmp_path],
        brand_palette=palette,
    )
    assert out.exists()

    # Inspect the injected SmartArt's colors part — should contain brand hex,
    # not the default schemeClr accent1/lt1/dk1.
    with zipfile.ZipFile(out) as z:
        diagram_color_parts = [
            n for n in z.namelist()
            if n.startswith("ppt/diagrams/") and "colors" in n
        ]
        assert len(diagram_color_parts) >= 1
        # Inspect the most recent (last) colors part — that's the injected one
        injected_colors_xml = z.read(diagram_color_parts[-1])

    assert b"14213D" in injected_colors_xml, "primary fill (#14213D) should be injected"
    assert b"E8DEC2" in injected_colors_xml, "text on primary (#E8DEC2) should be injected"
    assert b"1B2333" in injected_colors_xml, "text on surface (#1B2333) should be injected"
    # All three default scheme refs should be replaced for the patched carrier
    root = etree.fromstring(injected_colors_xml)
    schemes = root.findall(f".//{{{A_NS}}}schemeClr")
    scheme_vals = {s.get("val") for s in schemes}
    assert "accent1" not in scheme_vals
    assert "lt1" not in scheme_vals
    assert "dk1" not in scheme_vals


def test_apply_enrichment_without_brand_palette_keeps_default_colours(
    tmp_path, seed_variant_a,
):
    """Backwards compat: when brand_palette is None (default), the SmartArt
    carrier's stock Microsoft palette must reach the output unchanged.
    Pre-Contract-1 callers continue to work."""
    from pptx import Presentation
    from src.enrichment import EnrichmentItem, EnrichmentPlan, apply_enrichment

    prs = Presentation(str(seed_variant_a))
    target_marker = None
    target_slide_idx = None
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name and shape.name.startswith("SMARTART:"):
                target_marker = shape.name
                target_slide_idx = idx
                break
        if target_marker:
            break
    assert target_marker is not None

    out = tmp_path / "default.pptx"
    plan = EnrichmentPlan(
        source_pptx=seed_variant_a,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=target_slide_idx,
                kind="smartart",
                marker_name=target_marker,
                asset_path=None,
                action="apply",
                smartart_spec={
                    "graphic_type": "flowchart",
                    "layout_id": "process1",
                    "data": {"items": ["one", "two", "three"]},
                },
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])  # NO brand_palette

    with zipfile.ZipFile(out) as z:
        diagram_color_parts = [
            n for n in z.namelist()
            if n.startswith("ppt/diagrams/") and "colors" in n
        ]
        injected_colors_xml = z.read(diagram_color_parts[-1])

    # Default schemeClr accent1 should be preserved (not replaced with srgbClr)
    assert b"schemeClr" in injected_colors_xml
    assert b'val="accent1"' in injected_colors_xml
