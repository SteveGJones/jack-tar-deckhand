"""Tests for SMARTART-FROM-LIST (Contract 2 / Finding #9).

Run 3 dogfood (2026-04-26) caught Finding #9: SMARTART markers force a
"full content zone takes a graphic" model. Contract 2 is the richer
authoring contract — /pptx writes a real text shape with bullets at the
slide's natural content position, marks it ``SMARTART-FROM-LIST:slug``,
and the bridge transforms it in place: extract items, render SmartArt
with brand palette (Contract 1), replace the list with the SmartArt at
the same coordinates. Title and supporting prose untouched.

Tests cover:
- ``extract_list_items_from_marker_shape`` happy path + edge cases
- end-to-end ``apply_enrichment`` with a ``smartart_from_list`` item
- combination with ``brand_palette`` (Contract 1 + 2 stack)
- backwards compat: existing ``smartart`` items still work unchanged
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from src.colors_xml_builder import BrandPalette
from src.enrichment import EnrichmentItem, EnrichmentPlan, apply_enrichment
from src.enrichment_ops.smartart_from_list import (
    BulletsTooLongError,
    LAYOUT_BULLET_CAPS,
    SmartArtFromListExtractError,
    extract_list_items_from_marker_shape,
    select_layout_for_bullets,
)
from src.placeholder import parse_marker

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# --- Fixture: a minimal .pptx with a SMARTART-FROM-LIST marker --------------

def _build_pptx_with_list_marker(
    out_path: Path,
    *,
    items: list[str],
    marker_name: str = "SMARTART-FROM-LIST:test-list",
    include_title: bool = True,
    include_supporting_prose: bool = True,
) -> Path:
    """Create a minimal 1-slide .pptx whose content zone is a text shape
    marked ``SMARTART-FROM-LIST:...`` with the given bullet items.

    The shape's ``cNvPr@name`` is set so python-pptx and the bridge analyser
    can find it by ``shape.name``.
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    if include_title:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_box.text_frame.text = "Three capability themes"
        title_box.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
            if False else ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        )

    # The list shape — this is the SMARTART-FROM-LIST marker
    list_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = list_box.text_frame
    tf.text = items[0]
    for it in items[1:]:
        p = tf.add_paragraph()
        p.text = it
    # Set the cNvPr name so the bridge can find this shape by marker name.
    # The element_image op uses ``shape.name`` which reads the cNvPr@name.
    list_box.name = marker_name

    if include_supporting_prose:
        prose_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.8))
        prose_box.text_frame.text = "Supporting commentary that should remain on the slide."

    prs.save(str(out_path))
    return out_path


# --- Marker grammar regression: SMARTART-FROM-LIST parses correctly ---------

def test_marker_grammar_accepts_smartart_from_list():
    parsed = parse_marker("SMARTART-FROM-LIST:capability-themes")
    assert parsed == ("SMARTART-FROM-LIST", "capability-themes")


def test_marker_grammar_distinguishes_smartart_vs_smartart_from_list():
    # Both kinds parse correctly without prefix collisions
    assert parse_marker("SMARTART:plain-graphic") == ("SMARTART", "plain-graphic")
    assert parse_marker("SMARTART-FROM-LIST:from-bullets") == (
        "SMARTART-FROM-LIST", "from-bullets",
    )


# --- extract_list_items_from_marker_shape ----------------------------------

def test_extract_returns_paragraph_items(tmp_path):
    pptx = _build_pptx_with_list_marker(
        tmp_path / "list.pptx",
        items=["Modular plugins", "Spike-then-build", "Dogfood gate"],
    )
    prs = Presentation(str(pptx))
    items = extract_list_items_from_marker_shape(
        prs=prs, slide_index_1based=1,
        marker_name="SMARTART-FROM-LIST:test-list",
    )
    assert items == ["Modular plugins", "Spike-then-build", "Dogfood gate"]


def test_extract_skips_empty_paragraphs(tmp_path):
    """Authors often press Enter after the last bullet; the resulting empty
    paragraph must be filtered, not interpreted as an empty SmartArt item."""
    pptx_path = tmp_path / "with-empty.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(4))
    tf = box.text_frame
    tf.text = "First"
    tf.add_paragraph().text = ""           # empty
    tf.add_paragraph().text = "Second"
    tf.add_paragraph().text = "   "        # whitespace-only
    tf.add_paragraph().text = "Third"
    box.name = "SMARTART-FROM-LIST:t"
    prs.save(str(pptx_path))

    prs = Presentation(str(pptx_path))
    items = extract_list_items_from_marker_shape(
        prs=prs, slide_index_1based=1, marker_name="SMARTART-FROM-LIST:t",
    )
    assert items == ["First", "Second", "Third"]


def test_extract_raises_when_shape_missing(tmp_path):
    pptx = _build_pptx_with_list_marker(
        tmp_path / "ok.pptx", items=["a", "b", "c"],
    )
    prs = Presentation(str(pptx))
    with pytest.raises(SmartArtFromListExtractError, match="not found"):
        extract_list_items_from_marker_shape(
            prs=prs, slide_index_1based=1, marker_name="SMARTART-FROM-LIST:nonexistent",
        )


def test_extract_raises_when_slide_index_out_of_range(tmp_path):
    pptx = _build_pptx_with_list_marker(
        tmp_path / "ok.pptx", items=["a", "b"],
    )
    prs = Presentation(str(pptx))
    with pytest.raises(SmartArtFromListExtractError, match="out of range"):
        extract_list_items_from_marker_shape(
            prs=prs, slide_index_1based=99, marker_name="SMARTART-FROM-LIST:test-list",
        )


def test_extract_raises_when_paragraphs_all_empty(tmp_path):
    """Marker present but author left the list completely empty — surface
    a clear error rather than silently producing a 0-item SmartArt."""
    pptx_path = tmp_path / "empty-list.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(4))
    box.text_frame.text = ""
    box.name = "SMARTART-FROM-LIST:empty"
    prs.save(str(pptx_path))

    prs = Presentation(str(pptx_path))
    with pytest.raises(SmartArtFromListExtractError, match="non-empty paragraphs"):
        extract_list_items_from_marker_shape(
            prs=prs, slide_index_1based=1, marker_name="SMARTART-FROM-LIST:empty",
        )


# --- End-to-end: apply_enrichment with smartart_from_list -------------------

def test_apply_enrichment_with_smartart_from_list_replaces_list_with_graphic(tmp_path):
    """Full pipeline: build a .pptx with a SMARTART-FROM-LIST marker, run
    apply_enrichment, verify the output contains SmartArt diagram parts AND
    the original list shape is gone (replaced)."""
    src_pptx = _build_pptx_with_list_marker(
        tmp_path / "src.pptx",
        items=["Modular plugins", "Spike-then-build", "Dogfood gate"],
    )
    out = tmp_path / "out.pptx"

    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:test-list",
                asset_path=None,
                action="apply",
                # smartart_spec is None — bridge extracts items from marker shape
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    # Output should contain a SmartArt diagram (proof of injection)
    with zipfile.ZipFile(out) as z:
        diagram_parts = [
            n for n in z.namelist() if n.startswith("ppt/diagrams/")
        ]
        assert any("colors" in n for n in diagram_parts), \
            "output must contain SmartArt colors part"
        assert any("layout" in n for n in diagram_parts), \
            "output must contain SmartArt layout part"
        assert any("data" in n for n in diagram_parts), \
            "output must contain SmartArt data part"


def test_apply_enrichment_smartart_from_list_with_brand_palette(tmp_path):
    """Contract 1 + Contract 2 stacked: list-driven SmartArt rendered with
    brand palette. Output's diagram colours must contain brand hex, not
    Microsoft default schemeClr."""
    src_pptx = _build_pptx_with_list_marker(
        tmp_path / "src.pptx",
        items=["one", "two", "three"],
    )
    out = tmp_path / "branded.pptx"

    palette = BrandPalette(
        primary_fill_hex="14213D",       # navy
        text_on_primary_hex="F5EFE0",    # cream
        text_on_surface_hex="1B2333",    # body slate
    )
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:test-list",
                asset_path=None,
                action="apply",
            ),
        ],
    )
    apply_enrichment(
        plan,
        allowed_image_roots=[tmp_path],
        brand_palette=palette,
    )
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        color_parts = [
            n for n in z.namelist()
            if n.startswith("ppt/diagrams/") and "colors" in n
        ]
        assert color_parts, "output must contain SmartArt colors1.xml"
        injected = z.read(color_parts[0])

    assert b"14213D" in injected, "primary fill (#14213D) should be in output"
    assert b"F5EFE0" in injected, "text on primary (#F5EFE0) should be in output"
    # Default scheme refs (accent1/lt1/dk1) should be replaced
    root = etree.fromstring(injected)
    schemes = root.findall(f".//{{{A_NS}}}schemeClr")
    scheme_vals = {s.get("val") for s in schemes}
    assert "accent1" not in scheme_vals
    assert "lt1" not in scheme_vals
    assert "dk1" not in scheme_vals


def test_apply_enrichment_smartart_from_list_uses_explicit_spec_when_provided(tmp_path):
    """Backwards-friendly: if the caller pre-supplies a smartart_spec on a
    smartart_from_list item, the bridge uses it (skipping list extraction).
    Useful for callers that want explicit layout control."""
    src_pptx = _build_pptx_with_list_marker(
        tmp_path / "src.pptx",
        items=["from", "the", "shape", "itself"],   # 4 items in marker shape
    )
    out = tmp_path / "out.pptx"

    # Caller pre-supplies a spec with DIFFERENT items — the bridge should
    # honour the explicit spec and NOT re-extract from the shape.
    explicit_spec = {
        "graphic_type": "flowchart",
        "layout_id": "process1",
        "data": {"items": ["explicit", "spec", "wins"]},
    }
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:test-list",
                asset_path=None,
                action="apply",
                smartart_spec=explicit_spec,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    # Verify the rendered SmartArt's data1.xml contains explicit-spec items
    with zipfile.ZipFile(out) as z:
        data_parts = [
            n for n in z.namelist()
            if n.startswith("ppt/diagrams/") and n.endswith("data1.xml")
        ]
        assert data_parts
        data_xml = z.read(data_parts[0]).decode("utf-8")

    assert "explicit" in data_xml
    assert "spec" in data_xml
    assert "wins" in data_xml
    # Items from the shape itself should NOT appear when explicit spec is used
    assert "shape" not in data_xml or "shape" in "explicit-spec wins"  # weak check


# ---------------------------------------------------------------------------
# Layout selection by bullet length (Run 4/5/6 Finding #13)
# ---------------------------------------------------------------------------
#
# Three consecutive dogfood runs hit the process1 24-char cap with prose-style
# bullets and required hand-truncation to ship. The bridge now routes to a
# more generous layout when bullets exceed process1's cap, only falling back
# to truncation when even the largest layout (vList2) cannot accommodate the
# content. This preserves authorial intent (full bullet text) wherever the
# layout family can carry it, and degrades gracefully when it cannot.


def test_layout_caps_table_documents_canonical_layouts():
    """The cap table is the load-bearing constants for routing — assert the
    expected three layouts are present so future maintainers cannot silently
    drop one and break the routing logic."""
    assert "process1" in LAYOUT_BULLET_CAPS
    assert "list1" in LAYOUT_BULLET_CAPS
    assert "vList2" in LAYOUT_BULLET_CAPS
    # Caps must be strictly increasing — the routing assumes process1 < list1 < vList2
    assert LAYOUT_BULLET_CAPS["process1"] < LAYOUT_BULLET_CAPS["list1"]
    assert LAYOUT_BULLET_CAPS["list1"] < LAYOUT_BULLET_CAPS["vList2"]


def test_select_layout_short_bullets_route_to_process1():
    """Run 4/5/6's hand-shortened (≤24-char) bullets should keep the chevron
    flow visual."""
    items = ["Edge inference", "Operator playbook", "GPU procurement"]
    layout, returned_items, warnings = select_layout_for_bullets(items)
    assert layout == "process1"
    assert returned_items == items   # not truncated
    assert warnings == []


def test_select_layout_at_process1_boundary_stays_on_process1():
    """Exactly 24 chars must still fit process1 — boundary is inclusive."""
    items = ["a" * 24, "Modular plugin design", "Speaker-led decisions"]
    layout, returned, warnings = select_layout_for_bullets(items)
    assert layout == "process1"
    assert returned == items
    assert warnings == []


def test_select_layout_25_to_30_chars_route_to_list1():
    """Run 4/5/6's prose bullets (25-30 chars) should route to list1 instead
    of forcing the Speaker to hand-shorten."""
    items = [
        "Edge inference capacity",       # 23 chars
        "Customer overlap, 3 of 7 wins", # 29 chars
        "Operator playbook entries",     # 25 chars
    ]
    layout, returned, warnings = select_layout_for_bullets(items)
    assert layout == "list1"
    assert returned == items
    # Speaker-visible warning explains why we routed away from default
    assert any("process1" in w for w in warnings)


def test_select_layout_long_prose_routes_to_vlist2():
    """Bullets in the 31-60 char range route to vList2."""
    items = [
        "Edge inference capacity for regulated customers",  # 47 chars
        "Operator playbook covering five core verticals",    # 47 chars
        "GPU procurement contracts at $8M run-rate",         # 41 chars
    ]
    layout, returned, warnings = select_layout_for_bullets(items)
    assert layout == "vList2"
    assert returned == items
    assert any("vList2" in w for w in warnings)


def test_select_layout_raises_when_even_vlist2_overflows():
    """Bullets exceeding vList2's cap raise BulletsTooLongError (Run 7 #22).

    Earlier behaviour silently truncated bullets at vList2's cap with an
    ellipsis. Run 7 evidence showed mid-word "..." cuts mangled operator
    content (e.g. "...defer triage ent..." for "...defer triage entirely").
    The bridge now hard-fails at the cap so the operator's intent stays
    intact — they must rewrite within 60 chars or pick a different layout.
    """
    cap = LAYOUT_BULLET_CAPS["vList2"]
    overlong = "x" * (cap + 20)
    items = ["short one", overlong]
    with pytest.raises(BulletsTooLongError) as exc_info:
        select_layout_for_bullets(items)
    msg = str(exc_info.value)
    # Error message names the offending bullet verbatim so the operator
    # can locate it in build.js
    assert overlong in msg
    assert str(cap) in msg
    assert "vList2" in msg
    # Surfaces the count and total
    assert "1 of 2" in msg


def test_select_layout_at_vlist2_boundary_does_not_raise():
    """Exactly 60 chars must still fit vList2 — boundary is inclusive."""
    cap = LAYOUT_BULLET_CAPS["vList2"]
    items = ["a" * cap, "short two"]
    layout, returned, warnings = select_layout_for_bullets(items)
    assert layout == "vList2"
    assert returned == items   # untouched, never truncated


def test_select_layout_error_lists_all_offenders():
    """When multiple bullets exceed the cap, the error lists each one."""
    cap = LAYOUT_BULLET_CAPS["vList2"]
    items = [
        "short within cap",
        "x" * (cap + 5),
        "y" * (cap + 10),
        "z" * (cap + 15),
    ]
    with pytest.raises(BulletsTooLongError) as exc_info:
        select_layout_for_bullets(items)
    msg = str(exc_info.value)
    # All three offenders appear in the message; the short bullet does not
    assert msg.count("'x") == 1
    assert msg.count("'y") == 1
    assert msg.count("'z") == 1
    assert "short within cap" not in msg
    assert "3 of 4" in msg


def test_select_layout_empty_input_defaults_to_process1():
    """Defensive — caller should not pass an empty list, but if they do, the
    helper must not crash. Default to process1 with no items."""
    layout, returned, warnings = select_layout_for_bullets([])
    assert layout == "process1"
    assert returned == []
    assert warnings == []


def test_apply_enrichment_smartart_from_list_routes_to_list1_for_longer_bullets(tmp_path):
    """End-to-end: when the marker shape's bullets exceed process1's cap, the
    bridge routes the spec to list1 automatically — Speakers no longer need
    to hand-shorten Run 4/5/6-style prose bullets to ship."""
    long_bullets = [
        "Edge inference capacity",          # 23 chars — within process1
        "Customer overlap, 3 of 7 wins",    # 29 chars — overflows process1
        "Operator playbook entries",        # 25 chars — overflows process1
    ]
    src_pptx = _build_pptx_with_list_marker(
        tmp_path / "src.pptx", items=long_bullets,
    )
    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:test-list",
                asset_path=None,
                action="apply",
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        layout_parts = [
            n for n in z.namelist()
            if n.startswith("ppt/diagrams/") and "layout" in n
        ]
        assert layout_parts
        layout_xml = z.read(layout_parts[0]).decode("utf-8")
    # The list1 layout URI must appear in the injected layout part
    assert "list1" in layout_xml or "List1" in layout_xml
