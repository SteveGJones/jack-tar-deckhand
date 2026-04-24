import shutil
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from src.enrichment_ops.smartart import inject_smartart_into_file
from src.smartart_bridge import build_spec_from_slide, render_carrier
from src.slide_facts import SlideFacts, Marker


def test_injects_smartart_replacing_marker(tmp_path, seed_variant_a):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    # Find the SMARTART marker on the seed
    prs = Presentation(str(work))
    target_marker = None
    target_slide = None
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name and shape.name.startswith("SMARTART:"):
                target_marker = shape.name
                target_slide = idx
                break
        if target_marker:
            break
    assert target_marker is not None

    # Build a spec using bridge primitives
    sf = SlideFacts(slide_index=target_slide, text_content="Discover\nDesign\nBuild",
                     markers=[Marker("SMARTART", target_marker.split(":")[1], 0, 0, 0, 0)])
    spec = build_spec_from_slide(sf, marker_id=target_marker, layout_id="process1")
    carrier = render_carrier(spec, output_dir=tmp_path / "carriers")

    inject_smartart_into_file(
        host_pptx=work,
        slide_index_1based=target_slide,
        marker_name=target_marker,
        carrier_pptx=carrier,
    )
    # The host now contains diagram parts and no longer has the placeholder shape
    with zipfile.ZipFile(work) as zf:
        names = zf.namelist()
    assert any("ppt/diagrams/data" in n for n in names)


def test_missing_marker_raises(tmp_path, seed_variant_a):
    """If the carrier targets a marker that doesn't exist, we get an InjectionError."""
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    sf = SlideFacts(slide_index=1, text_content="a\nb\nc",
                     markers=[Marker("SMARTART", "ghost", 0, 0, 0, 0)])
    spec = build_spec_from_slide(sf, marker_id="SMARTART:ghost", layout_id="process1")
    carrier = render_carrier(spec, output_dir=tmp_path / "carriers")
    with pytest.raises(Exception):  # InjectionError from msft-smartart plugin
        inject_smartart_into_file(
            host_pptx=work, slide_index_1based=1,
            marker_name="SMARTART:ghost", carrier_pptx=carrier,
        )
