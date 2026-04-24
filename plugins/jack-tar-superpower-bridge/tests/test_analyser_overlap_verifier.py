import pytest
from pptx import Presentation
from pptx.util import Inches

from src.analyser.overlap_verifier import find_overlaps
from src.slide_facts import SlideFacts, Marker, OverlapWarning


def _build_deck_with_smartart_marker_and_text(tmp_path, *, overlap: bool):
    from pptx.util import Emu
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Marker rectangle at (1in, 1in, 5in, 3in)
    marker = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(1), Inches(1), Inches(5), Inches(3),
    )
    marker.name = "SMARTART:three-pillars"
    # Body text — overlap if requested
    if overlap:
        body = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(0.5))
    else:
        body = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(2), Inches(0.5))
    body.name = "Body 1"
    body.text_frame.text = "Hello"
    out = tmp_path / "deck.pptx"
    prs.save(out)
    return out


def test_overlap_detected_when_text_intersects_marker(tmp_path):
    deck = _build_deck_with_smartart_marker_and_text(tmp_path, overlap=True)
    warnings = find_overlaps(deck)
    assert len(warnings) == 1
    w = warnings[0]
    assert w.marker_id == "SMARTART:three-pillars"
    assert "Body 1" in w.overlapping_shape_names


def test_no_overlap_returns_empty_list(tmp_path):
    deck = _build_deck_with_smartart_marker_and_text(tmp_path, overlap=False)
    warnings = find_overlaps(deck)
    assert warnings == []


def test_only_smartart_markers_trigger_check(tmp_path):
    """An IMAGE marker overlapping body text is NOT a warning — only SMARTART
    markers cause cosmetic bleed because they get replaced with a graphic."""
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    marker = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(5), Inches(3))
    marker.name = "IMAGE:foo"
    body = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(0.5))
    body.name = "Body 1"
    body.text_frame.text = "x"
    out = tmp_path / "image.pptx"
    prs.save(out)
    assert find_overlaps(out) == []


def test_multiple_overlapping_shapes_collected(tmp_path):
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    marker = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(5), Inches(3))
    marker.name = "SMARTART:foo"
    for i in range(3):
        b = slide.shapes.add_textbox(Inches(2), Inches(1.5 + i*0.6), Inches(2), Inches(0.5))
        b.name = f"Body {i}"
        b.text_frame.text = "x"
    out = tmp_path / "many.pptx"
    prs.save(out)
    warnings = find_overlaps(out)
    assert len(warnings) == 1
    assert len(warnings[0].overlapping_shape_names) == 3
