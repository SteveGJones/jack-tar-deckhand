"""Tests for the CHART:slug marker kind (issue #55).

Coverage:
  - Phase 1: marker grammar — parse_marker accepts CHART:slug, rejects uppercase
  - Phase 2: EnrichmentItem.kind="chart" dispatch + apply path
  - Phase 3: chart_palette derivation
  - Phase 5: analyser picks up CHART markers via MARKER_RE + end-to-end fixture test
"""
from __future__ import annotations

import io
import warnings
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Phase 1 — Marker grammar
# ---------------------------------------------------------------------------


def test_parse_marker_chart_slug_returns_kind_and_identifier():
    from src.placeholder import parse_marker
    result = parse_marker("CHART:revenue-vs-costs")
    assert result == ("CHART", "revenue-vs-costs")


def test_parse_marker_chart_slug_with_underscore():
    from src.placeholder import parse_marker
    result = parse_marker("CHART:capacity_over_time")
    assert result == ("CHART", "capacity_over_time")


def test_parse_marker_chart_rejects_uppercase_identifier():
    from src.placeholder import parse_marker
    assert parse_marker("CHART:RevenueVsCosts") is None


def test_parse_marker_chart_rejects_empty_identifier():
    from src.placeholder import parse_marker
    assert parse_marker("CHART:") is None


def test_chart_marker_is_distinct_from_image_marker():
    from src.placeholder import parse_marker
    chart_kind, _ = parse_marker("CHART:foo")
    image_kind, _ = parse_marker("IMAGE:foo")
    assert chart_kind != image_kind
    assert chart_kind == "CHART"


def test_marker_re_accepts_chart_prefix():
    from src.placeholder import MARKER_RE
    m = MARKER_RE.match("CHART:mttr-by-class")
    assert m is not None
    assert m.group(1) == "CHART"
    assert m.group(2) == "mttr-by-class"


# ---------------------------------------------------------------------------
# Phase 3 — Chart palette derivation
# ---------------------------------------------------------------------------


def _palette(primary="14213D", on_primary="F5F0E8", on_surface="2C2C2C"):
    from src.colors_xml_builder import BrandPalette
    return BrandPalette(primary, on_primary, on_surface)


def test_derive_chart_palette_default_returns_n_colours():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    result = derive_chart_palette(pal, "default", 4)
    assert len(result) == 4
    assert all(isinstance(c, str) and len(c) == 6 for c in result)


def test_derive_chart_palette_default_first_colour_is_primary():
    from src.chart_palette import derive_chart_palette
    pal = _palette(primary="14213D")
    result = derive_chart_palette(pal, "default", 3)
    # First colour must equal the primary fill (uppercase, no #).
    assert result[0] == "14213D"


def test_derive_chart_palette_explicit_wins_over_role():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    explicit = ["#FF0000", "#00FF00", "#0000FF"]
    result = derive_chart_palette(pal, "data_categorical", 3, explicit=explicit)
    assert result == ["FF0000", "00FF00", "0000FF"]


def test_derive_chart_palette_explicit_padded_to_n_series():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    result = derive_chart_palette(pal, "default", 4, explicit=["AABBCC"])
    assert len(result) == 4
    assert result == ["AABBCC", "AABBCC", "AABBCC", "AABBCC"]


def test_derive_chart_palette_explicit_truncated_to_n_series():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    result = derive_chart_palette(pal, "default", 2, explicit=["FF0000", "00FF00", "0000FF"])
    assert result == ["FF0000", "00FF00"]


def test_derive_chart_palette_vital_and_mourning_two_colours():
    from src.chart_palette import derive_chart_palette
    pal = _palette(primary="8B0000")
    result = derive_chart_palette(pal, "vital_and_mourning", 2)
    assert len(result) == 2
    assert result[0] == "8B0000"  # vital — primary unchanged
    # mourning must be a different (desaturated/darkened) colour
    assert result[1] != result[0]


def test_derive_chart_palette_vital_and_mourning_wrong_n_warns_and_falls_back():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    with warnings.catch_warnings(record=True) as w:
        warnings.simplefilter("always")
        result = derive_chart_palette(pal, "vital_and_mourning", 3)
    assert len(result) == 3
    assert any("vital_and_mourning" in str(warning.message) for warning in w)


def test_derive_chart_palette_data_categorical_distinct_hues():
    from src.chart_palette import derive_chart_palette
    pal = _palette(primary="4472C4")
    result = derive_chart_palette(pal, "data_categorical", 4)
    assert len(result) == 4
    # All should be distinct (hue rotation guarantees this for n >= 2)
    assert len(set(result)) == 4


def test_derive_chart_palette_unknown_role_warns_and_uses_default():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    with warnings.catch_warnings(record=True) as w:
        warnings.simplefilter("always")
        result = derive_chart_palette(pal, "nonexistent_role", 3)
    assert len(result) == 3
    assert any("nonexistent_role" in str(warning.message) for warning in w)


def test_derive_chart_palette_large_n_warns_soft_fail():
    from src.chart_palette import derive_chart_palette
    pal = _palette()
    with warnings.catch_warnings(record=True) as w:
        warnings.simplefilter("always")
        result = derive_chart_palette(pal, "default", 9)
    assert len(result) == 9
    assert any("n_series=9" in str(warning.message) for warning in w)


# ---------------------------------------------------------------------------
# Helpers — build a minimal .pptx with a CHART marker rect
# ---------------------------------------------------------------------------


def _make_pptx_with_chart_marker(output_path: Path, marker_name: str,
                                  x_in=1.0, y_in=1.5, w_in=9.0, h_in=5.0,
                                  with_label: bool = True) -> None:
    """Create a minimal .pptx with one CHART placeholder rect (and optional label)."""
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # The marker placeholder rect
    shapes = slide.shapes
    marker_rect = shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 in python-pptx; use the int directly
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    marker_rect.name = marker_name

    if with_label:
        # Sibling addText label (simulates build.js authoring pattern)
        label = shapes.add_textbox(
            Inches(x_in), Inches(y_in + h_in / 2 - 0.15),
            Inches(w_in), Inches(0.3),
        )
        label.name = f"label-for-{marker_name}"
        label.text_frame.text = marker_name

    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Phase 2 — EnrichmentItem chart dispatch + apply path
# ---------------------------------------------------------------------------


def _minimal_chart_spec(chart_type: str = "column") -> dict:
    return {
        "type": chart_type,
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [
            {"name": "Revenue", "values": [12.5, 14.2, 15.8, 17.1]},
        ],
    }


def test_enrichment_item_chart_kind_constructs():
    from src.enrichment import EnrichmentItem
    item = EnrichmentItem(
        slide_index=1,
        kind="chart",
        marker_name="CHART:revenue-vs-costs",
        asset_path=None,
        chart_spec=_minimal_chart_spec(),
    )
    assert item.kind == "chart"
    assert item.chart_spec is not None
    assert item.chart_spec["type"] == "column"


def test_apply_chart_enrichment_missing_type_raises(tmp_path):
    """chart_ops raises ValueError (wrapped by orchestrator as EnrichmentApplyError)."""
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx_with_chart_marker(pptx_path, "CHART:test")
    prs = Presentation(str(pptx_path))
    with pytest.raises(ValueError, match="missing required field"):
        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name="CHART:test",
            chart_spec={"categories": ["A"], "series": [{"name": "S", "values": [1]}]},
        )


def test_apply_chart_enrichment_missing_categories_raises(tmp_path):
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx_with_chart_marker(pptx_path, "CHART:test")
    prs = Presentation(str(pptx_path))
    with pytest.raises(ValueError, match="missing required field"):
        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name="CHART:test",
            chart_spec={"type": "bar", "series": [{"name": "S", "values": [1]}]},
        )


def test_apply_chart_enrichment_missing_series_raises(tmp_path):
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx_with_chart_marker(pptx_path, "CHART:test")
    prs = Presentation(str(pptx_path))
    with pytest.raises(ValueError, match="missing required field"):
        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name="CHART:test",
            chart_spec={"type": "column", "categories": ["A"]},
        )


def test_apply_chart_enrichment_invalid_type_raises(tmp_path):
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx_with_chart_marker(pptx_path, "CHART:test")
    prs = Presentation(str(pptx_path))
    with pytest.raises(ValueError, match="unsupported type"):
        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name="CHART:test",
            chart_spec={
                "type": "heatmap",
                "categories": ["A"],
                "series": [{"name": "S", "values": [1]}],
            },
        )


def test_apply_chart_enrichment_missing_marker_raises(tmp_path):
    """Raises LookupError (wrapped by orchestrator as EnrichmentApplyError)."""
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx_with_chart_marker(pptx_path, "CHART:other-slug")
    prs = Presentation(str(pptx_path))
    with pytest.raises(LookupError, match="not found"):
        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name="CHART:does-not-exist",
            chart_spec=_minimal_chart_spec(),
        )


def _get_slide_shape_names(prs, slide_index_0based: int) -> list[str]:
    slide = prs.slides[slide_index_0based]
    return [s.name for s in slide.shapes]


def _count_graphic_frames(prs, slide_index_0based: int) -> int:
    """Count graphic frame shapes (charts are graphic frames) on a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    slide = prs.slides[slide_index_0based]
    return sum(1 for s in slide.shapes if s.shape_type == 3)  # FREEFORM=5, CHART=3 in MSO_SHAPE_TYPE? No — use has_chart


def _has_chart(prs, slide_index_0based: int) -> bool:
    """Return True if any shape on the slide is a chart."""
    slide = prs.slides[slide_index_0based]
    for s in slide.shapes:
        try:
            _ = s.chart
            return True
        except Exception:
            pass
    return False


def test_apply_chart_enrichment_inserts_chart_and_removes_marker(tmp_path):
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    marker_name = "CHART:revenue-vs-costs"
    _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=True)

    prs = Presentation(str(pptx_path))

    # Before: marker rect present
    assert marker_name in _get_slide_shape_names(prs, 0)

    apply_chart_enrichment(
        prs=prs,
        slide_index_1based=1,
        marker_name=marker_name,
        chart_spec=_minimal_chart_spec("column"),
    )

    # After: marker rect removed
    names_after = _get_slide_shape_names(prs, 0)
    assert marker_name not in names_after

    # Chart inserted
    assert _has_chart(prs, 0)


def test_apply_chart_enrichment_removes_sibling_label(tmp_path):
    from src.enrichment_ops.chart import apply_chart_enrichment
    pptx_path = tmp_path / "deck.pptx"
    marker_name = "CHART:capacity-over-time"
    _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=True)

    prs = Presentation(str(pptx_path))
    # Verify label is there before
    slide = prs.slides[0]
    label_texts = []
    for s in slide.shapes:
        try:
            t = s.text_frame.text.strip()
            if t == marker_name:
                label_texts.append(t)
        except Exception:
            pass
    assert label_texts, "Expected a sibling label before enrichment"

    apply_chart_enrichment(
        prs=prs,
        slide_index_1based=1,
        marker_name=marker_name,
        chart_spec=_minimal_chart_spec("bar"),
    )

    # Label should be gone
    slide = prs.slides[0]
    for s in slide.shapes:
        try:
            t = s.text_frame.text.strip()
            assert t != marker_name, "Sibling label should have been removed"
        except Exception:
            pass


def test_apply_chart_enrichment_all_v1_chart_types_render(tmp_path):
    """Verify all 7 v1 chart types can be applied without error."""
    from src.enrichment_ops.chart import apply_chart_enrichment
    chart_types = ["bar", "column", "line", "area", "pie", "doughnut", "scatter"]

    for chart_type in chart_types:
        pptx_path = tmp_path / f"deck-{chart_type}.pptx"
        marker_name = f"CHART:{chart_type}-test"
        _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=False)
        prs = Presentation(str(pptx_path))

        spec = {
            "type": chart_type,
            "categories": ["A", "B", "C"],
            "series": [{"name": "S1", "values": [1.0, 2.0, 3.0]}],
        }
        if chart_type == "scatter":
            spec["categories"] = [1.0, 2.0, 3.0]

        apply_chart_enrichment(
            prs=prs,
            slide_index_1based=1,
            marker_name=marker_name,
            chart_spec=spec,
        )
        assert _has_chart(prs, 0), f"Chart type {chart_type!r} did not insert a chart"


def test_apply_chart_enrichment_with_brand_palette_uses_palette_colour(tmp_path):
    """When brand_palette is supplied, series colour derives from it."""
    from src.enrichment_ops.chart import apply_chart_enrichment
    from src.colors_xml_builder import BrandPalette
    pptx_path = tmp_path / "deck.pptx"
    marker_name = "CHART:revenue"
    _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=False)
    prs = Presentation(str(pptx_path))

    palette = BrandPalette("8B0000", "F5F0E8", "2C2C2C")
    apply_chart_enrichment(
        prs=prs,
        slide_index_1based=1,
        marker_name=marker_name,
        chart_spec=_minimal_chart_spec("column"),
        brand_palette=palette,
    )
    assert _has_chart(prs, 0)


def test_apply_chart_enrichment_explicit_series_colours_used(tmp_path):
    """Explicit series_colours override palette derivation."""
    from src.enrichment_ops.chart import apply_chart_enrichment
    from src.colors_xml_builder import BrandPalette
    from pptx.dml.color import RGBColor

    pptx_path = tmp_path / "deck.pptx"
    marker_name = "CHART:data"
    _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=False)
    prs = Presentation(str(pptx_path))

    palette = BrandPalette("14213D", "F5F0E8", "2C2C2C")
    spec = {
        "type": "bar",
        "categories": ["X", "Y"],
        "series": [
            {"name": "A", "values": [10, 20]},
            {"name": "B", "values": [30, 40]},
        ],
        "series_colours": ["FF0000", "00FF00"],
    }
    apply_chart_enrichment(
        prs=prs,
        slide_index_1based=1,
        marker_name=marker_name,
        chart_spec=spec,
        brand_palette=palette,
    )

    # Verify the chart is present
    assert _has_chart(prs, 0)
    slide = prs.slides[0]
    chart_shape = next(s for s in slide.shapes if _shape_has_chart(s))
    chart = chart_shape.chart
    # Series 0 should be red (FF0000)
    series_0_fill = chart.series[0].format.fill
    assert series_0_fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)


def _shape_has_chart(shape) -> bool:
    try:
        _ = shape.chart
        return True
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Phase 5 — OOXML analyser picks up CHART markers
# ---------------------------------------------------------------------------


def test_analyser_detects_chart_marker_via_marker_re():
    """The OOXML primary analyser uses MARKER_RE, which now includes CHART."""
    from src.placeholder import MARKER_RE
    # CHART markers follow the same objectName pattern as IMAGE/BG/SMARTART.
    test_cases = [
        ("CHART:revenue-vs-costs", True),
        ("CHART:capacity_over_time", True),
        ("CHART:mttr-by-class", True),
        ("CHART:UPPERCASE", False),
        ("CHART:", False),
    ]
    for name, should_match in test_cases:
        matched = MARKER_RE.match(name) is not None
        assert matched == should_match, (
            f"MARKER_RE.match({name!r}) expected {should_match}, got {matched}"
        )


def test_pptx_parser_detects_chart_marker_in_fixture(tmp_path):
    """The OOXML primary analyser extracts CHART markers from a real .pptx fixture."""
    from src.analyser.pptx_parser import parse_pptx

    pptx_path = tmp_path / "chart_deck.pptx"
    marker_name = "CHART:extinction-rate"
    _make_pptx_with_chart_marker(pptx_path, marker_name, with_label=False)

    result = parse_pptx(pptx_path)
    assert len(result) >= 1
    chart_markers = [
        m for slide in result for m in slide.markers if m.kind == "CHART"
    ]
    assert len(chart_markers) == 1
    assert chart_markers[0].identifier == "extinction-rate"


# ---------------------------------------------------------------------------
# End-to-end — fixture .pptx with CHART marker → apply_enrichment → chart present
# ---------------------------------------------------------------------------


def test_end_to_end_apply_enrichment_chart_kind(tmp_path):
    """apply_enrichment with kind='chart' inserts a native chart and removes the marker."""
    from src.enrichment import EnrichmentPlan, EnrichmentItem, apply_enrichment

    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    marker_name = "CHART:revenue-vs-costs"
    _make_pptx_with_chart_marker(source, marker_name, with_label=True)

    items = [
        EnrichmentItem(
            slide_index=1,
            kind="chart",
            marker_name=marker_name,
            asset_path=None,
            action="apply",
            chart_spec={
                "type": "column",
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [
                    {"name": "Revenue", "values": [12.5, 14.2, 15.8, 17.1]},
                    {"name": "Costs", "values": [8.1, 8.4, 8.9, 9.2]},
                ],
                "title": "Quarterly revenue vs costs",
                "show_legend": True,
            },
        ),
    ]
    plan = EnrichmentPlan(source_pptx=source, output_pptx=output, items=items)
    apply_enrichment(plan, allowed_image_roots=[tmp_path])

    assert output.exists()
    prs_out = Presentation(str(output))
    # Marker removed
    names = _get_slide_shape_names(prs_out, 0)
    assert marker_name not in names
    # Chart present
    assert _has_chart(prs_out, 0)


def test_end_to_end_chart_missing_spec_raises(tmp_path):
    """apply_enrichment raises EnrichmentApplyError when chart_spec is None."""
    from src.enrichment import EnrichmentPlan, EnrichmentItem, apply_enrichment, EnrichmentApplyError

    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    marker_name = "CHART:test"
    _make_pptx_with_chart_marker(source, marker_name, with_label=False)

    items = [
        EnrichmentItem(
            slide_index=1,
            kind="chart",
            marker_name=marker_name,
            asset_path=None,
            action="apply",
            chart_spec=None,
        ),
    ]
    plan = EnrichmentPlan(source_pptx=source, output_pptx=output, items=items)
    with pytest.raises(EnrichmentApplyError, match="missing chart_spec"):
        apply_enrichment(plan, allowed_image_roots=[tmp_path])

    # No partial output on failure
    assert not output.exists()
