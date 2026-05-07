import pytest
from pptx import Presentation
from pptx.util import Inches

from src.analyser.pptx_parser import parse_pptx
from src.slide_facts import SlideFacts


def test_parses_variant_a_extracts_eight_markers(seed_variant_a):
    facts = parse_pptx(seed_variant_a)
    assert len(facts) == 10
    total_markers = sum(len(s.markers) for s in facts)
    assert total_markers == 8


def test_parses_variant_a_marker_kinds(seed_variant_a):
    facts = parse_pptx(seed_variant_a)
    kinds = {m.kind for s in facts for m in s.markers}
    assert kinds <= {"IMAGE", "SMARTART", "BG"}
    assert "IMAGE" in kinds
    assert "BG" in kinds


def test_marker_geometry_is_emu(seed_variant_a):
    facts = parse_pptx(seed_variant_a)
    for slide in facts:
        for marker in slide.markers:
            assert isinstance(marker.left_emu, int)
            assert isinstance(marker.top_emu, int)
            assert marker.width_emu > 0
            assert marker.height_emu > 0


def test_variant_b_returns_zero_markers(seed_variant_b):
    """Variant B uses `name:` which PptxGenJS drops — OOXML sees no marker shape names."""
    facts = parse_pptx(seed_variant_b)
    total = sum(len(s.markers) for s in facts)
    assert total == 0


def test_text_content_extracted_from_text_frames(seed_variant_a):
    facts = parse_pptx(seed_variant_a)
    assert any(s.text_content.strip() for s in facts)


def test_element_types_counts_charts_images(seed_variant_a):
    facts = parse_pptx(seed_variant_a)
    # Every slide must have all 5 keys present (uniform shape contract; downstream
    # consumers index by key without .get() defaults).
    for slide in facts:
        for key in ("text", "shape", "image", "chart", "table"):
            assert key in slide.element_types, (
                f"slide {slide.slide_index} missing element_types[{key!r}]"
            )
    # Variant A is text+shape heavy with no charts/tables/images outside markers
    assert all(s.element_types["chart"] == 0 for s in facts)
    assert all(s.element_types["table"] == 0 for s in facts)
    # At least one slide has body text shapes
    assert sum(s.element_types["text"] for s in facts) > 0


def test_bgcolor_attribute_recognised_as_solid_background(tmp_path):
    """If a slide carries a bgColor attribute (PptxGenJS colour backgrounds), the
    parser must report background_kind='solid' instead of 'default'."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Inject the PptxGenJS-style bgColor attribute on <p:sld>
    slide.element.set("bgColor", "0E1513")
    out = tmp_path / "bgcolor.pptx"
    prs.save(out)

    facts = parse_pptx(out)
    assert facts[0].background_kind == "solid"


def test_image_background_recognised_via_blipfill(tmp_path):
    """If a slide has <p:bg><a:blipFill>, background_kind must be 'image'."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    slide = prs.slides[0]
    cSld = slide.element.find(qn("p:cSld"))
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:bgPr><a:blipFill><a:blip r:embed="rId99"/><a:stretch><a:fillRect/></a:stretch>'
        '</a:blipFill></p:bgPr></p:bg>'
    )
    cSld.insert(0, etree.fromstring(bg_xml))
    out = tmp_path / "blipfill.pptx"
    prs.save(out)

    facts = parse_pptx(out)
    assert facts[0].background_kind == "image"


def test_nonexistent_file_raises_filenotfound(tmp_path):
    with pytest.raises(FileNotFoundError):
        parse_pptx(tmp_path / "nope.pptx")
