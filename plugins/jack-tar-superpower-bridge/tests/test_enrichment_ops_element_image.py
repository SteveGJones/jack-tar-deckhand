import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from src.enrichment_ops.element_image import replace_image_marker_in_memory
from src.security import AllowedPathError


def test_replaces_image_marker_with_picture_at_same_geometry(
    tmp_path, seed_variant_a, placeholder_png
):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "elem.png"; shutil.copy(placeholder_png, img)

    prs = Presentation(str(work))
    # Find the IMAGE marker on Variant A
    target_marker: str | None = None
    target_geom = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name and shape.name.startswith("IMAGE:"):
                target_marker = shape.name
                target_geom = (shape.left, shape.top, shape.width, shape.height)
                break
        if target_marker:
            break
    assert target_marker is not None, "Variant A must have at least one IMAGE marker"

    replace_image_marker_in_memory(
        prs=prs,
        marker_name=target_marker,
        image_path=img,
        allowed_image_roots=[img_dir],
    )

    # The marker shape is gone; a PICTURE shape now occupies the same geometry
    found_marker = False
    found_picture_at_geom = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == target_marker:
                found_marker = True
            if (shape.left, shape.top, shape.width, shape.height) == target_geom:
                if shape.shape_type is not None and shape.shape_type == 13:  # PICTURE
                    found_picture_at_geom = True
    assert not found_marker
    assert found_picture_at_geom


def test_unknown_marker_raises_lookuperror(tmp_path, seed_variant_a, placeholder_png):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "x.png"; shutil.copy(placeholder_png, img)

    prs = Presentation(str(work))
    with pytest.raises(LookupError, match="not found"):
        replace_image_marker_in_memory(
            prs=prs, marker_name="IMAGE:nonexistent",
            image_path=img, allowed_image_roots=[img_dir],
        )


def test_image_outside_allowlist_raises(tmp_path, seed_variant_a, placeholder_png):
    work = tmp_path / "deck.pptx"
    shutil.copy(seed_variant_a, work)
    rogue = tmp_path / "rogue.png"; shutil.copy(placeholder_png, rogue)

    prs = Presentation(str(work))
    allowed = tmp_path / "generated"; allowed.mkdir()
    with pytest.raises(AllowedPathError):
        replace_image_marker_in_memory(
            prs=prs, marker_name="IMAGE:agent-architecture",
            image_path=rogue, allowed_image_roots=[allowed],
        )
