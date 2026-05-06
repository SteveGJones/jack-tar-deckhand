"""Regression test for Run 9 Finding #29 — picture-SmartArt media collision.

Run 9 surfaced: when a host .pptx has both slide-level IMAGE/BG markers
(media added by python-pptx during apply_enrichment Phase 2) AND a
picture-SmartArt-from-list (carrier media grafted in Phase 4), the
carrier graft overwrote ``ppt/media/image1.png`` etc. Slide-level rels
still pointed at ``image1.png`` but the file at that path was now a
species portrait, not the silent-forest watercolour the operator had
authored.

v0.1.4 fix: ``msft-smartart/assembler_patch.inject()`` now scans the
host's existing media for ``max(imageN.ext)`` and renumbers carrier
media to start at ``max + 1``. Carrier rels Targets are rewritten to
match. This test reproduces the Run 9 deck shape (combined IMAGE +
picture-SmartArt + BG markers on the same .pptx) and asserts:

1. The slide-level IMAGE/BG markers' rels still resolve to the original
   image bytes (not overwritten).
2. The picture-SmartArt carrier media files are present in the host
   under renumbered filenames.
3. The carrier's data rels point at the renumbered targets.
4. Both image sets coexist; no filename collisions.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from src.enrichment import EnrichmentItem, EnrichmentPlan, apply_enrichment


def _make_image(path: Path, colour: str, size: int = 64) -> Path:
    img = Image.new("RGB", (size, size), color=colour)
    img.save(str(path), "PNG")
    return path


def _build_combined_marker_pptx(out_path: Path) -> Path:
    """Build a minimal 2-slide deck with BOTH a SmartArt-from-list AND a
    slide-level IMAGE marker. This is the Run 9 reproduction recipe."""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Slide 1: SmartArt-from-list marker (will get picture binding)
    s1 = prs.slides.add_slide(blank)
    list_box = s1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5))
    tf = list_box.text_frame
    tf.text = "Alpha"
    for label in ["Beta", "Gamma"]:
        p = tf.add_paragraph()
        p.text = label
    list_box.name = "SMARTART-FROM-LIST:rollcall"

    # Slide 2: slide-level IMAGE marker (will get a slide-level image)
    s2 = prs.slides.add_slide(blank)
    img_rect = s2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(6), Inches(4),
    )
    img_rect.name = "IMAGE:atmospheric"

    prs.save(str(out_path))
    return out_path


def test_combined_picture_smartart_and_slide_image_marker_no_collision(
    tmp_path: Path,
) -> None:
    """The Run 9 combined-marker reproduction: assert no media collision.

    This is the load-bearing v0.1.4 verification. Without the fix, slide
    2's IMAGE marker would resolve to a species portrait (the carrier's
    image1.png overwriting python-pptx's image1.png). With the fix, the
    carrier's media is renumbered so both sets of images coexist.
    """
    src_pptx = _build_combined_marker_pptx(tmp_path / "src.pptx")

    # Three SmartArt portraits (one per bullet) — distinct colours so
    # bytes are distinct and we can detect overwrites.
    portraits_dir = tmp_path / "portraits"
    portraits_dir.mkdir()
    portrait_paths = [
        _make_image(portraits_dir / "portrait-1.png", "navy"),
        _make_image(portraits_dir / "portrait-2.png", "maroon"),
        _make_image(portraits_dir / "portrait-3.png", "purple"),
    ]

    # One slide-level atmospheric image — distinct colour
    atmo_dir = tmp_path / "atmo"
    atmo_dir.mkdir()
    atmo_image = _make_image(atmo_dir / "forest.png", "olive", size=128)

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            # Picture-SmartArt with 3 species portraits — populates ppt/media/
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:rollcall",
                asset_path=None,
                per_item_image_paths=portrait_paths,
            ),
            # Slide-level IMAGE marker — separately populates ppt/media/
            EnrichmentItem(
                slide_index=2,
                kind="image",
                marker_name="IMAGE:atmospheric",
                asset_path=atmo_image,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        names = set(z.namelist())
        media_files = {
            n: z.read(n) for n in names if n.startswith("ppt/media/")
        }

        # Assertion 1: AT LEAST 4 media files present (1 atmospheric + 3 portraits)
        assert len(media_files) >= 4, (
            f"expected >= 4 media files (1 atmo + 3 portraits), "
            f"got {len(media_files)}: {sorted(media_files)}"
        )

        # Assertion 2: every distinct image must coexist (no overwrites)
        # — count distinct file sizes
        size_set = {len(b) for b in media_files.values()}
        assert len(size_set) >= 2, (
            f"all media files have same size — likely overwrites. "
            f"sizes: {[len(b) for b in media_files.values()]}"
        )

        # Assertion 3: slide 2 rels point at a media file whose bytes match
        # the atmospheric image, NOT a portrait.
        slide2_rels_xml = z.read("ppt/slides/_rels/slide2.xml.rels").decode("utf-8")
        # Find the image Target in slide 2 rels
        import re
        image_targets = re.findall(
            r'Type="[^"]*relationships/image"\s+Target="([^"]+)"',
            slide2_rels_xml,
        )
        assert image_targets, f"slide 2 has no image rel: {slide2_rels_xml}"
        target = image_targets[0]
        # Resolve the relative target: ../media/imageN.png from
        # ppt/slides/_rels/ → ppt/media/imageN.png
        resolved = "ppt/media/" + target.rsplit("/", 1)[-1]
        assert resolved in media_files, (
            f"slide 2 rel resolves to {resolved}, but file not in package: "
            f"{sorted(media_files)}"
        )
        atmospheric_bytes = atmo_image.read_bytes()
        assert media_files[resolved] == atmospheric_bytes, (
            f"slide 2 IMAGE marker resolves to {resolved} but its bytes "
            f"differ from the atmospheric image. Run 9 Finding #29: carrier "
            f"graft overwrote the slide-level marker's media. "
            f"Expected {len(atmospheric_bytes)} bytes, "
            f"got {len(media_files[resolved])} bytes."
        )

        # Assertion 4: carrier's data1.xml.rels points at media that exists
        # AND has portrait bytes (not atmospheric bytes).
        data_rels_name = next(
            (n for n in names if "data1.xml.rels" in n), None
        )
        assert data_rels_name is not None, (
            "carrier's data1.xml.rels missing from output"
        )
        data_rels_xml = z.read(data_rels_name).decode("utf-8")
        carrier_targets = re.findall(
            r'Type="[^"]*relationships/image"\s+Target="([^"]+)"',
            data_rels_xml,
        )
        assert len(carrier_targets) == 3, (
            f"expected 3 image rels in carrier (one per portrait), "
            f"got {len(carrier_targets)}: {carrier_targets}"
        )
        # All carrier targets must resolve to media files in the host
        # AND those media files must NOT be the atmospheric image.
        for ct in carrier_targets:
            resolved_ct = "ppt/media/" + ct.rsplit("/", 1)[-1]
            assert resolved_ct in media_files, (
                f"carrier rel target {ct} resolves to {resolved_ct} "
                f"but file not in package"
            )
            assert media_files[resolved_ct] != atmospheric_bytes, (
                f"carrier rel {ct} resolves to atmospheric image bytes — "
                f"the rels were not renumbered correctly during injection"
            )

        # Assertion 5: carrier targets and slide-level target are DIFFERENT
        # filenames (no overlap).
        slide2_target_leaf = target.rsplit("/", 1)[-1]
        carrier_target_leaves = {ct.rsplit("/", 1)[-1] for ct in carrier_targets}
        assert slide2_target_leaf not in carrier_target_leaves, (
            f"slide 2 IMAGE rel and carrier rel both point at "
            f"{slide2_target_leaf} — collision not avoided."
        )


def test_picture_smartart_alone_still_works_after_fix(tmp_path: Path) -> None:
    """The fix must not regress the picture-SmartArt-only case (Run 8 Finding #26).

    A deck with ONLY a picture-SmartArt (no slide-level IMAGE markers)
    should still graft the carrier's media correctly. After the fix, the
    media renumbers from image1.png (since the host has no images yet).
    """
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    list_box = s1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5))
    tf = list_box.text_frame
    tf.text = "Alpha"
    for label in ["Beta", "Gamma"]:
        p = tf.add_paragraph()
        p.text = label
    list_box.name = "SMARTART-FROM-LIST:alone"
    src_pptx = tmp_path / "src.pptx"
    prs.save(str(src_pptx))

    portraits_dir = tmp_path / "portraits"
    portraits_dir.mkdir()
    portrait_paths = [
        _make_image(portraits_dir / f"p{i}.png", "olive")
        for i in range(3)
    ]

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:alone",
                asset_path=None,
                per_item_image_paths=portrait_paths,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        media_files = [n for n in z.namelist() if n.startswith("ppt/media/")]
        # Three portraits in carrier + at least no overwrites
        assert len(media_files) >= 3, (
            f"expected at least 3 portrait media files, got {media_files}"
        )
