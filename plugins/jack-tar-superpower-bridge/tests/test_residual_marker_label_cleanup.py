"""Tests for v0.2 #23 — extend the v0.1.x #17 BG-cleanup pattern to all marker kinds.

Run 8 + Run 9 dogfood revealed: when /pptx authors a small ``addText``
label alongside a SMARTART-FROM-LIST or SMARTART marker (e.g. a
grey-italic ``"SMARTART-FROM-LIST:tooling-gaps"`` label floating beside
the bullet list, so the marker is visible in the unenriched deck), the
label survives enrichment as a stranded cosmetic floater on top of the
rendered SmartArt.

The v0.1.x #17 fix solved this for ``BG:`` markers via
``apply_background_in_memory``. v0.2 generalises the cleanup to all
marker kinds via ``_drop_residual_marker_label_shapes`` in
``enrichment.py``.

Strict-equality test: only shapes whose visible text exactly equals the
marker name are removed. Shapes that merely CONTAIN the marker string
as a substring stay (so legitimate body prose mentioning the marker
name doesn't get silently deleted).

Marker shape itself is NOT removed by this helper (the inject path
handles its replacement); only auxiliary label addText shapes go.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.enrichment import EnrichmentItem, EnrichmentPlan, apply_enrichment


def _build_pptx_with_marker_and_label(
    out_path: Path,
    *,
    marker_name: str,
    items: list[str],
    add_residual_label: bool = True,
) -> Path:
    """Create a 1-slide deck with a SMARTART-FROM-LIST marker + (optionally)
    a separate addText shape labelling it (the v0.2 #23 reproduction)."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Marker shape: bullet list named with the marker name
    list_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(4))
    tf = list_box.text_frame
    tf.text = items[0]
    for it in items[1:]:
        p = tf.add_paragraph()
        p.text = it
    list_box.name = marker_name

    if add_residual_label:
        # Residual marker-name label — the cosmetic floater the cleanup
        # should remove. Different shape, different cNvPr name, but the
        # visible text equals the marker name.
        label = slide.shapes.add_textbox(Inches(6.5), Inches(1.0), Inches(3), Inches(0.4))
        label.text_frame.text = marker_name
        # Default name is auto-generated like "TextBox 5" — important
        # that it's NOT the marker name, otherwise the inject path would
        # replace it.

    prs.save(str(out_path))
    return out_path


def _shape_with_text_count(pptx_path: Path, slide_index_1based: int, target: str) -> int:
    """Return the count of <a:t> runs whose stripped text equals ``target``
    on the named slide. Used to assert the residual label was removed."""
    a_t_qn = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    with zipfile.ZipFile(pptx_path) as z:
        slide_xml = z.read(
            f"ppt/slides/slide{slide_index_1based}.xml"
        ).decode("utf-8")
    from lxml import etree
    root = etree.fromstring(slide_xml.encode("utf-8"))
    count = 0
    for t in root.iter(a_t_qn):
        if t.text and t.text.strip() == target:
            count += 1
    return count


def test_residual_smartart_from_list_label_removed(tmp_path: Path) -> None:
    """Run 8/9 reproduction: SMARTART-FROM-LIST marker + residual label
    addText. After enrichment, only the marker (now SmartArt) should
    remain; the label addText should be gone."""
    src = _build_pptx_with_marker_and_label(
        tmp_path / "src.pptx",
        marker_name="SMARTART-FROM-LIST:tooling-gaps",
        items=["one", "two", "three"],
        add_residual_label=True,
    )

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:tooling-gaps",
                asset_path=None,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    # The residual label's text "SMARTART-FROM-LIST:tooling-gaps" must
    # NOT appear anywhere on slide 1's <a:t> runs after enrichment.
    count = _shape_with_text_count(
        out, slide_index_1based=1, target="SMARTART-FROM-LIST:tooling-gaps",
    )
    assert count == 0, (
        f"residual marker label not removed — found {count} shapes "
        f"with text 'SMARTART-FROM-LIST:tooling-gaps' on slide 1. "
        f"v0.2 #23 cleanup should have dropped them."
    )


def test_marker_without_residual_label_unchanged(tmp_path: Path) -> None:
    """Negative regression — when there's NO residual label, the
    cleanup pass shouldn't break anything."""
    src = _build_pptx_with_marker_and_label(
        tmp_path / "src.pptx",
        marker_name="SMARTART-FROM-LIST:cleanmarker",
        items=["a", "b"],
        add_residual_label=False,
    )

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:cleanmarker",
                asset_path=None,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()


def test_legitimate_text_mentioning_marker_substring_kept(tmp_path: Path) -> None:
    """Strict equality only — body prose that merely CONTAINS the marker
    name as a substring stays. Only shapes whose entire stripped text
    equals the marker name are removed."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Marker
    marker_name = "SMARTART-FROM-LIST:my-list"
    list_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(4))
    list_box.text_frame.text = "one"
    list_box.text_frame.add_paragraph().text = "two"
    list_box.name = marker_name

    # Body prose that MENTIONS the marker name as a substring of a
    # longer sentence — must be preserved.
    body = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    body.text_frame.text = (
        f"Note: this slide uses {marker_name} for the bullet rendering."
    )

    src = tmp_path / "src.pptx"
    prs.save(str(src))

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name=marker_name,
                asset_path=None,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    # The body prose must still be present
    with zipfile.ZipFile(out) as z:
        slide_xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "this slide uses" in slide_xml, (
        "legitimate body prose mentioning the marker name as a substring "
        "was incorrectly removed — the cleanup should be strict-equality only"
    )


def test_residual_smartart_label_removed(tmp_path: Path) -> None:
    """Same #23 cleanup applies to plain SMARTART markers (full content zone),
    not just SMARTART-FROM-LIST. Verifies the cleanup is in both branches."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    marker_name = "SMARTART:flowchart-stages"
    # SMARTART placeholder (full content zone) — bridge expects a shape
    # with this name. Use add_shape (rectangle) like a placeholder.
    from pptx.enum.shapes import MSO_SHAPE
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(8), Inches(5),
    )
    placeholder.name = marker_name

    # Residual label
    label = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.4))
    label.text_frame.text = marker_name

    src = tmp_path / "src.pptx"
    prs.save(str(src))

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart",
                marker_name=marker_name,
                asset_path=None,
                smartart_spec={
                    "graphic_type": "flowchart",
                    "layout_id": "process1",
                    "data": {"items": ["alpha", "beta", "gamma"]},
                },
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    count = _shape_with_text_count(
        out, slide_index_1based=1, target=marker_name,
    )
    assert count == 0, (
        f"residual SMARTART marker label not removed — found {count}."
    )
