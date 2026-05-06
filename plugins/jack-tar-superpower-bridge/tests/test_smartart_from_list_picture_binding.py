"""Tests for Run 8 Finding #26 — picture-SmartArt image-binding.

When ``EnrichmentItem.kind == "smartart_from_list"`` is supplied alongside
``per_item_image_paths``, the bridge must:

1. Force the routed layout to a picture-capable target (vList3 / pList1).
2. Validate that the image-paths length matches the bullet-item count.
3. Resolve every image path against the security allowlist.
4. Construct the picture-builder's expected ``items: [{label, image_path}]``
   shape so the engine binds each child node to its image via blipFill.
5. Produce a carrier .pptx whose ``ppt/diagrams/data1.xml`` contains
   blipFill elements AND whose ``ppt/diagrams/_rels/data1.xml.rels``
   contains image relationships AND whose ``ppt/media/`` carries the
   embedded image files.

This is the load-bearing v0.1.3 verification for Finding #26.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from src.enrichment import (
    EnrichmentApplyError,
    EnrichmentItem,
    EnrichmentPlan,
    apply_enrichment,
)


def _build_pptx_with_picture_list_marker(
    out_path: Path,
    *,
    items: list[str],
    marker_name: str = "SMARTART-FROM-LIST:species-rollcall",
) -> Path:
    """Create a minimal 1-slide .pptx with a SMARTART-FROM-LIST marker."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    list_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5))
    tf = list_box.text_frame
    tf.text = items[0]
    for it in items[1:]:
        p = tf.add_paragraph()
        p.text = it
    list_box.name = marker_name
    prs.save(str(out_path))
    return out_path


def _create_sample_image(path: Path, colour: str = "navy") -> Path:
    """Create a tiny solid-colour PNG so blipFill embeds have real bytes."""
    img = Image.new("RGB", (64, 64), color=colour)
    img.save(str(path), "PNG")
    return path


def test_picture_binding_routes_to_pList1_for_six_items(tmp_path: Path) -> None:
    """Six species + six images → pList1 (gallery row, max_nodes 8).

    The Run 8 motivating case: a roll-call of 5-8 lost species, each
    with a portrait. With 6 items, the bridge picks pList1 (5+ items
    threshold).
    """
    species = [
        "Passenger Pigeon (1914)",
        "Great Auk (1844)",
        "Thylacine (1936)",
        "Baiji River Dolphin (2006)",
        "Golden Toad (1989)",
        "Spix's Macaw (2000)",
    ]
    src_pptx = _build_pptx_with_picture_list_marker(
        tmp_path / "src.pptx",
        items=species,
    )

    images_dir = tmp_path / "portraits"
    images_dir.mkdir()
    image_paths = [
        _create_sample_image(images_dir / f"species-{i}.png", "navy")
        for i in range(len(species))
    ]

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:species-rollcall",
                asset_path=None,
                per_item_image_paths=image_paths,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])

    assert out.exists()

    # The carrier .pptx grafted into the output should have blipFill
    # elements in data1.xml and image-rels in data1.xml.rels.
    with zipfile.ZipFile(out) as z:
        diagram_parts = [n for n in z.namelist() if n.startswith("ppt/diagrams/")]
        media_parts = [n for n in z.namelist() if n.startswith("ppt/media/")]

        # Diagram parts present
        assert any("data1.xml" in n for n in diagram_parts), \
            f"data1.xml missing from {diagram_parts}"
        assert any("layout1.xml" in n for n in diagram_parts), \
            f"layout1.xml missing from {diagram_parts}"

        # data1.xml contains blipFill elements (proof of image binding)
        data_part_name = next(
            n for n in diagram_parts if n.endswith("data1.xml")
        )
        data_xml = z.read(data_part_name).decode("utf-8")
        assert "blipFill" in data_xml, (
            f"data1.xml missing <a:blipFill> elements — picture binding "
            f"failed; got XML head:\n{data_xml[:500]}"
        )

        # rels file binds image relationships
        rels_name = next(
            (n for n in diagram_parts if "data1.xml.rels" in n), None
        )
        assert rels_name is not None, \
            f"missing data1.xml.rels — image relationships not wired"
        rels_xml = z.read(rels_name).decode("utf-8")
        # Each image gets one relationship; 6 images expected
        assert rels_xml.count('Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"') == len(species), (
            f"expected {len(species)} image relationships, "
            f"got: {rels_xml}"
        )

        # Media files embedded
        # The injection pipeline may rename media on graft; assert we
        # have at least the right count of media files.
        assert len(media_parts) >= len(species), (
            f"expected at least {len(species)} embedded images, "
            f"got {len(media_parts)}: {media_parts}"
        )


def test_picture_binding_routes_to_vList3_for_three_items(tmp_path: Path) -> None:
    """Three items + three images → vList3 (portrait stack, fewer items)."""
    items = ["Alpha", "Beta", "Gamma"]
    src_pptx = _build_pptx_with_picture_list_marker(
        tmp_path / "src.pptx",
        items=items,
        marker_name="SMARTART-FROM-LIST:test-rollcall",
    )

    images_dir = tmp_path / "imgs"
    images_dir.mkdir()
    image_paths = [
        _create_sample_image(images_dir / f"img-{i}.png", "olive")
        for i in range(len(items))
    ]

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:test-rollcall",
                asset_path=None,
                per_item_image_paths=image_paths,
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        # vList3 is a picture-data-shape layout; assert the picture-binding
        # plumbing fired.
        diagram_parts = [n for n in z.namelist() if n.startswith("ppt/diagrams/")]
        data_part_name = next(
            n for n in diagram_parts if n.endswith("data1.xml")
        )
        data_xml = z.read(data_part_name).decode("utf-8")
        assert "blipFill" in data_xml, "data1.xml missing blipFill"


def test_image_count_mismatch_raises(tmp_path: Path) -> None:
    """If images count doesn't match bullet count, apply_enrichment fails fast."""
    items = ["one", "two", "three"]
    src_pptx = _build_pptx_with_picture_list_marker(
        tmp_path / "src.pptx",
        items=items,
        marker_name="SMARTART-FROM-LIST:mismatch",
    )

    images_dir = tmp_path / "imgs"
    images_dir.mkdir()
    image_paths = [
        _create_sample_image(images_dir / "a.png"),
        _create_sample_image(images_dir / "b.png"),
        # Only two images for three bullets — mismatch
    ]

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:mismatch",
                asset_path=None,
                per_item_image_paths=image_paths,
            ),
        ],
    )
    with pytest.raises(EnrichmentApplyError) as exc_info:
        apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert "per_item_image_paths length" in str(exc_info.value)
    assert "2" in str(exc_info.value)  # image count
    assert "3" in str(exc_info.value)  # bullet count
    assert not out.exists()


def test_image_paths_must_resolve_within_allowlist(tmp_path: Path) -> None:
    """Image paths outside the allowlist raise — security gate is honoured."""
    items = ["a", "b"]
    src_pptx = _build_pptx_with_picture_list_marker(
        tmp_path / "src.pptx",
        items=items,
        marker_name="SMARTART-FROM-LIST:security-gate",
    )

    # Image OUTSIDE tmp_path — security gate should reject
    rogue_dir = tmp_path.parent / "outside-allowlist"
    rogue_dir.mkdir(exist_ok=True)
    rogue_image = _create_sample_image(rogue_dir / "rogue.png")

    safe_dir = tmp_path / "imgs"
    safe_dir.mkdir()
    safe_image = _create_sample_image(safe_dir / "safe.png")

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:security-gate",
                asset_path=None,
                per_item_image_paths=[safe_image, rogue_image],
            ),
        ],
    )
    with pytest.raises(EnrichmentApplyError):
        # AllowedPathError wraps as EnrichmentApplyError via the
        # transactional wrapper.
        apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert not out.exists()


def test_no_image_paths_falls_through_to_text_routing(tmp_path: Path) -> None:
    """When per_item_image_paths is None (default), routing is text-only.

    This is the v0.1.2 behaviour — picture-binding is opt-in. Confirms
    the v0.1.3 extension doesn't change the no-images path.
    """
    items = ["short one", "short two", "short three"]
    src_pptx = _build_pptx_with_picture_list_marker(
        tmp_path / "src.pptx",
        items=items,
        marker_name="SMARTART-FROM-LIST:text-only",
    )

    out = tmp_path / "out.pptx"
    plan = EnrichmentPlan(
        source_pptx=src_pptx,
        output_pptx=out,
        items=[
            EnrichmentItem(
                slide_index=1,
                kind="smartart_from_list",
                marker_name="SMARTART-FROM-LIST:text-only",
                asset_path=None,
                # per_item_image_paths defaults to None
            ),
        ],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()

    with zipfile.ZipFile(out) as z:
        diagram_parts = [n for n in z.namelist() if n.startswith("ppt/diagrams/")]
        data_part_name = next(
            n for n in diagram_parts if n.endswith("data1.xml")
        )
        data_xml = z.read(data_part_name).decode("utf-8")
        # Text-only path: no blipFill expected
        assert "blipFill" not in data_xml, (
            "text-only routing should not emit blipFill; per_item_image_paths "
            f"should be opt-in. Got: {data_xml[:200]}"
        )
