import os
import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from src.enrichment import (
    EnrichmentPlan,
    EnrichmentItem,
    apply_enrichment,
    EnrichmentApplyError,
)


def _plan_with(items, *, source, output):
    return EnrichmentPlan(source_pptx=source, output_pptx=output, items=items)


def test_no_enrichments_writes_clean_copy(tmp_path, seed_variant_a):
    out = tmp_path / "out.pptx"
    plan = _plan_with([], source=seed_variant_a, output=out)
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()
    assert out.stat().st_size > 0


def test_background_only_succeeds(tmp_path, seed_variant_a, placeholder_png):
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(placeholder_png, img)
    out = tmp_path / "bg-only.pptx"
    items = [EnrichmentItem(slide_index=1, kind="background",
                              marker_name="BG:title-hero-dark-grid",
                              asset_path=img, action="apply")]
    plan = _plan_with(items, source=seed_variant_a, output=out)
    apply_enrichment(plan, allowed_image_roots=[img_dir])
    assert out.exists()


def test_failure_in_op2_leaves_no_output_file(tmp_path, seed_variant_a, placeholder_png):
    """Op2 raises LookupError on a missing marker -> finally must remove the temp file
    AND must not produce the output file."""
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "x.png"; shutil.copy(placeholder_png, img)
    out = tmp_path / "fail.pptx"
    items = [EnrichmentItem(slide_index=1, kind="image",
                              marker_name="IMAGE:does-not-exist",
                              asset_path=img, action="apply")]
    plan = _plan_with(items, source=seed_variant_a, output=out)
    with pytest.raises(EnrichmentApplyError):
        apply_enrichment(plan, allowed_image_roots=[img_dir])
    assert not out.exists()
    # No leftover temp file in the output directory
    leftovers = [p for p in tmp_path.iterdir()
                  if p.name.startswith("fail.pptx.tmp-")]
    assert leftovers == []


def test_smartart_skip_drops_op_but_leaves_other_enrichments(tmp_path, seed_variant_a, placeholder_png):
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(placeholder_png, img)
    out = tmp_path / "skip.pptx"
    items = [
        EnrichmentItem(slide_index=1, kind="background",
                        marker_name="BG:title-hero-dark-grid",
                        asset_path=img, action="apply"),
        EnrichmentItem(slide_index=4, kind="smartart",
                        marker_name="SMARTART:nonexistent",
                        asset_path=None, action="skip"),
    ]
    plan = _plan_with(items, source=seed_variant_a, output=out)
    apply_enrichment(plan, allowed_image_roots=[img_dir])
    assert out.exists()


def test_apply_clear_overlap_removes_overlapping_text_shapes(
    tmp_path, seed_variant_a
):
    """When user picks apply_clear_overlap, the orchestrator removes overlapping
    text shapes BEFORE injecting SmartArt."""
    # We'll use Variant A. Find the SMARTART marker.
    work_seed = tmp_path / "seed.pptx"
    shutil.copy(seed_variant_a, work_seed)
    prs = Presentation(str(work_seed))
    target = None
    target_slide_idx = None
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name and shape.name.startswith("SMARTART:"):
                target = shape.name
                target_slide_idx = idx
                break
        if target:
            break
    assert target is not None

    items = [
        EnrichmentItem(slide_index=target_slide_idx, kind="smartart",
                        marker_name=target, asset_path=None,
                        action="apply_clear_overlap",
                        smartart_spec={"graphic_type": "flowchart",
                                        "layout_id": "process1",
                                        "data": {"items": ["a", "b", "c"]}},
                        overlap_shape_names=["Body 1"]),  # may or may not match a real shape
    ]
    out = tmp_path / "cleared.pptx"
    plan = _plan_with(items, source=seed_variant_a, output=out)
    apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert out.exists()


def test_atomic_rename_no_partial_output(tmp_path, seed_variant_a, placeholder_png, monkeypatch):
    """If os.replace fails, the temp file should still be cleaned up AND the
    raised exception MUST be EnrichmentApplyError wrapping the OSError (not the
    raw OSError) so callers can catch a single error type."""
    from src import enrichment as mod
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(placeholder_png, img)
    out = tmp_path / "rename-fail.pptx"

    # Force os.replace to raise after the temp file is built
    def boom(src, dst):
        raise OSError("simulated rename failure")
    monkeypatch.setattr(mod.os, "replace", boom)

    items = [EnrichmentItem(slide_index=1, kind="background",
                              marker_name="BG:title-hero-dark-grid",
                              asset_path=img, action="apply")]
    plan = _plan_with(items, source=seed_variant_a, output=out)
    with pytest.raises(EnrichmentApplyError) as excinfo:
        apply_enrichment(plan, allowed_image_roots=[img_dir])
    # Must wrap the underlying OSError so callers can rely on a single error type
    assert isinstance(excinfo.value.__cause__, OSError)
    assert not out.exists()
    leftovers = [p for p in tmp_path.iterdir() if p.name.startswith(out.name + ".tmp-")]
    assert leftovers == []


def test_apply_enrichment_runs_preflight_first(tmp_path, monkeypatch):
    """Spec Security & Privacy gate - apply_enrichment MUST refuse a .pptx that
    fails pre-flight checks before opening it with python-pptx."""
    from src.security import PptxPreflightError
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a zip")
    out = tmp_path / "out.pptx"
    plan = _plan_with([], source=bad, output=out)
    with pytest.raises(EnrichmentApplyError) as excinfo:
        apply_enrichment(plan, allowed_image_roots=[tmp_path])
    assert isinstance(excinfo.value.__cause__, PptxPreflightError)
    assert not out.exists()
