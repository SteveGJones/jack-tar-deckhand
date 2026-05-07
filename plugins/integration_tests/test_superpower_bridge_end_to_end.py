"""Integration test: end-to-end happy path on Spike 1 Variant A.

Walks: analyse → build a tiny enrichment plan with one BG and one SmartArt
→ apply transactionally → verify outputs exist + open cleanly via python-pptx.
"""
import shutil
import sys
from pathlib import Path

WORKTREE = Path(__file__).resolve().parents[2]
BRIDGE_ROOT = WORKTREE / "plugins" / "jack-tar-superpower-bridge"
MSFT_ROOT = WORKTREE / "plugins" / "jack-tar-msft-smartart"
SPIKE1 = WORKTREE / "docs" / "spikes" / "2026-04-23-pptx-marker-adherence" / "outputs" / "variant-a"
SPIKE2 = WORKTREE / "docs" / "spikes" / "2026-04-23-python-pptx-enrichment"


def test_end_to_end_happy_path(tmp_path, monkeypatch):
    sys.path.insert(0, str(BRIDGE_ROOT))
    for k in list(sys.modules):
        if k == "src" or k.startswith("src."):
            del sys.modules[k]
    monkeypatch.setenv("JACK_TAR_SUPERPOWER_BRIDGE_FORCE_MSFT_ROOT", str(MSFT_ROOT))

    from src.analyser import analyse_pptx
    from src.enrichment import EnrichmentPlan, EnrichmentItem, apply_enrichment
    from src.smartart_bridge import build_spec_from_slide

    # Stage the seed
    work_seed = tmp_path / "src.pptx"
    shutil.copy(SPIKE1 / "presentation.pptx", work_seed)
    shutil.copy(SPIKE1 / "build.js", tmp_path / "build.js")
    img_dir = tmp_path / "generated"; img_dir.mkdir()
    img = img_dir / "bg.png"; shutil.copy(SPIKE2 / "seed" / "placeholder.png", img)

    result = analyse_pptx(work_seed)
    assert result.total_markers >= 1

    # Find a BG marker
    bg_marker = None
    bg_slide_idx = None
    for s in result.slides:
        for m in s.markers:
            if m.kind == "BG":
                bg_marker = f"BG:{m.identifier}"
                bg_slide_idx = s.slide_index
                break
        if bg_marker:
            break
    assert bg_marker is not None

    plan = EnrichmentPlan(
        source_pptx=work_seed,
        output_pptx=tmp_path / "enriched.pptx",
        items=[EnrichmentItem(slide_index=bg_slide_idx, kind="background",
                                marker_name=bg_marker, asset_path=img,
                                action="apply")],
    )
    apply_enrichment(plan, allowed_image_roots=[img_dir])
    assert (tmp_path / "enriched.pptx").exists()

    from pptx import Presentation
    prs = Presentation(str(tmp_path / "enriched.pptx"))
    assert len(prs.slides) == 10


def test_end_to_end_smartart_path(tmp_path, monkeypatch):
    """Caveat-fix #13 — exercise the cross-plugin SmartArt path end-to-end.

    The BG-only end-to-end test above does not exercise loader → carrier →
    inject. That chain is the single most likely place for sys.modules
    contamination to surface; this test forces it.
    """
    sys.path.insert(0, str(BRIDGE_ROOT))
    for k in list(sys.modules):
        if k == "src" or k.startswith("src."):
            del sys.modules[k]
    monkeypatch.setenv("JACK_TAR_SUPERPOWER_BRIDGE_FORCE_MSFT_ROOT", str(MSFT_ROOT))

    from src.analyser import analyse_pptx
    from src.enrichment import EnrichmentPlan, EnrichmentItem, apply_enrichment
    from src.smartart_bridge import build_spec_from_slide, select_layout_for_slide

    work_seed = tmp_path / "src.pptx"
    shutil.copy(SPIKE1 / "presentation.pptx", work_seed)

    result = analyse_pptx(work_seed)
    smartart_marker_id = None
    smartart_slide_idx = None
    smartart_slide = None
    for s in result.slides:
        for m in s.markers:
            if m.kind == "SMARTART":
                smartart_marker_id = f"SMARTART:{m.identifier}"
                smartart_slide_idx = s.slide_index
                smartart_slide = s
                break
        if smartart_marker_id:
            break
    assert smartart_marker_id is not None, \
        "Variant A must include at least one SMARTART marker for this test"

    layout_id = select_layout_for_slide(smartart_slide, marker_id=smartart_marker_id)
    spec = build_spec_from_slide(smartart_slide, marker_id=smartart_marker_id,
                                   layout_id=layout_id)

    plan = EnrichmentPlan(
        source_pptx=work_seed,
        output_pptx=tmp_path / "enriched.pptx",
        items=[EnrichmentItem(
            slide_index=smartart_slide_idx, kind="smartart",
            marker_name=smartart_marker_id, asset_path=None,
            action="apply", smartart_spec=spec,
        )],
    )
    apply_enrichment(plan, allowed_image_roots=[tmp_path])

    out = tmp_path / "enriched.pptx"
    assert out.exists()

    # The injected SmartArt's diagram parts must be present
    import zipfile
    with zipfile.ZipFile(out) as zf:
        names = zf.namelist()
    assert any("ppt/diagrams/data" in n for n in names), \
        "no diagram parts found in enriched deck — SmartArt injection failed silently"

    # Bridge's src.placeholder must STILL be importable as the bridge's module
    # after the cross-plugin call chain ran (sys.modules contamination guard).
    from src.placeholder import MARKER_RE
    assert MARKER_RE.match("IMAGE:foo") is not None
