"""Template-driven deck assembly via python-pptx.

Opens a corporate .pptx template, strips example slides, adds new slides
from the outline using the template's mapped layouts, and populates
placeholders with content.
"""

import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.annotation_payload import estimate_label_box, segment_box_entry
from src.process_image import compute_content_hash, get_dimensions

# Tag under <p:spTree> in the PresentationML namespace. Picture shapes must
# sit AFTER <p:nvGrpSpPr> and <p:grpSpPr> for PowerPoint to render the slide.
_PML_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_SPTREE_HEAD_TAGS = (f"{_PML_NS}nvGrpSpPr", f"{_PML_NS}grpSpPr")


def _resolve_layout(prs, profile, slide_type):
    """Find the python-pptx SlideLayout object for a given slide type."""
    mapping = profile.get('layout_mapping', {})
    entry = mapping.get(slide_type) or profile.get('unmapped_fallback')
    if not entry:
        return prs.slide_layouts[0]

    master_index = profile.get('master_index', 0)
    master = prs.slide_masters[master_index]
    layout_index = entry['layout_index']
    return master.slide_layouts[layout_index]


def _strip_existing_slides(prs):
    """Remove all existing slides from the presentation."""
    slide_list = prs.slides._sldIdLst
    for sldId in list(slide_list):
        slide_list.remove(sldId)


def _find_placeholder_by_type(slide, ph_type, profile_layout):
    """Find a slide placeholder matching the given type from profile layout metadata."""
    if not profile_layout:
        return None

    target_indices = set()
    for ph_info in profile_layout.get('placeholders', []):
        if ph_info['type'] == ph_type:
            target_indices.add(ph_info['idx'])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx in target_indices:
            return ph
    return None


def _populate_text(placeholder, text):
    """Set text on a placeholder, preserving template formatting."""
    if placeholder is None or not text:
        return
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text


def _populate_body_points(placeholder, body_points):
    """Set bulleted body points on a content or body placeholder."""
    if placeholder is None or not body_points:
        return
    tf = placeholder.text_frame
    tf.clear()
    for i, point in enumerate(body_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point


def _get_profile_layout(profile, layout_name):
    """Find a layout entry in the profile by name."""
    for layout in profile.get('layouts', []):
        if layout['name'] == layout_name:
            return layout
    return None


def _get_mapped_layout_name(profile, slide_type):
    """Get the layout name for a slide type from the mapping."""
    mapping = profile.get('layout_mapping', {})
    entry = mapping.get(slide_type) or profile.get('unmapped_fallback')
    if entry:
        return entry['layout_name']
    return None


def _apply_full_bleed(slide, image_path, slide_w, slide_h):
    """Render a full_bleed slide: strip all shapes, add picture at canvas size.

    Implements issue #88. Promotes the fullbleed_deck.py post-processor to
    first-class assembler behaviour. Image goes edge-to-edge with zero
    chrome — no title, no body, no footer placeholders.

    Args:
        slide: python-pptx Slide object (already added via prs.slides.add_slide).
        image_path: Absolute path to the picture file, or None.
        slide_w: Slide width in EMU (from prs.slide_width).
        slide_h: Slide height in EMU (from prs.slide_height).

    Edge case: when ``image_path`` is None or the file is missing, the slide
    is left empty (every shape stripped). Speakers spot the gap during
    review rather than seeing a misleading title placeholder.
    """
    for shape in list(slide.shapes):
        el = shape._element
        el.getparent().remove(el)

    if not image_path or not os.path.isfile(image_path):
        return

    pic = slide.shapes.add_picture(image_path, 0, 0, width=slide_w, height=slide_h)
    _hoist_picture_to_sptree_head(pic)


def _hoist_picture_to_sptree_head(pic):
    """Move a just-added picture shape to right after nvGrpSpPr/grpSpPr.

    PowerPoint requires picture shapes to sit after the two head tags
    (``_SPTREE_HEAD_TAGS``) for the slide to render. Shared by
    ``_apply_full_bleed`` and the annotate-figure v2 (#142 v2) contain-fit
    placement helper below, both of which strip all existing shapes and
    add a single picture that must land first in z-order.
    """
    sp_tree = pic._element.getparent()
    sp_tree.remove(pic._element)
    children = list(sp_tree)
    insert_at = 0
    for i, ch in enumerate(children):
        if ch.tag in _SPTREE_HEAD_TAGS:
            insert_at = i + 1
    sp_tree.insert(insert_at, pic._element)


# ---------------------------------------------------------------------------
# annotate-figure v2 — deck-native annotations (issue #142 v2, T5)
#
# Design: docs/superpowers/plans/2026-07-17-annotate-figure-v2.md §5.
# ---------------------------------------------------------------------------


def _contain_fit_rect(img_w, img_h, zone_rect):
    """Aspect-preserving contain-fit rect of an image within a zone (§3.2).

    Shared by both assembler paths conceptually (this is the python-pptx
    side; the JS assembler has its own equivalent in build_deck.js's
    existing content/diagram fit helpers). All units are whatever the
    caller passes for ``zone_rect`` (EMU here) — the maths is unit-agnostic.

    Args:
        img_w, img_h: native image dimensions (any consistent unit, e.g.
            pixels — only the RATIO is used).
        zone_rect: (x, y, w, h) of the placement zone, same unit as the
            returned rect.

    Returns:
        (fx, fy, fw, fh): the fitted rect, aspect-matched to img_w/img_h,
        centred within zone_rect on whichever axis has slack.
    """
    zx, zy, zw, zh = zone_rect
    img_ratio = img_w / img_h
    zone_ratio = zw / zh
    if img_ratio > zone_ratio:
        # Image wider than the zone -> fit to width, letterbox top/bottom.
        fw = zw
        fh = zw / img_ratio
        fx = zx
        fy = zy + (zh - fh) / 2
    else:
        # Fit to height, letterbox left/right.
        fh = zh
        fw = zh * img_ratio
        fx = zx + (zw - fw) / 2
        fy = zy
    return fx, fy, fw, fh


def _map_norm_point(norm_xy, fit_rect):
    """Map an image-normalized [0-1] point into a fitted rect's coordinate
    space (§3.2's ``X = fx + nx*fw``, ``Y = fy + ny*fh``)."""
    fx, fy, fw, fh = fit_rect
    nx, ny = norm_xy
    return (fx + nx * fw, fy + ny * fh)


def _place_contain_fit_picture(slide, image_path, zone_rect):
    """Strip all shapes and place ``image_path`` contain-fit inside
    ``zone_rect`` (EMU), hoisted to the front of z-order.

    Shared by the native-annotation overlay path and the raster / no-payload
    fallback paths (F11): every annotated image — whether it will carry a
    native overlay or already has labels baked into its pixels by v1's
    ``annotate()`` — is placed with aspect preserved, never cover-cropped
    (JS builders' ``cover``) or stretched (this module's un-annotated
    ``_apply_full_bleed``). Mirrors ``_apply_full_bleed``'s strip-then-place
    shape, but sizes+positions the picture to the aspect-preserving fit rect
    instead of the full zone.

    Args:
        slide: python-pptx Slide (chrome not yet stripped).
        image_path: absolute path to the image file, or None/missing.
        zone_rect: (x, y, w, h) EMU rect of the placement zone.

    Returns:
        The picture shape, or None if image_path is missing/absent (mirrors
        ``_apply_full_bleed``'s empty-slide edge case).
    """
    for shape in list(slide.shapes):
        el = shape._element
        el.getparent().remove(el)

    if not image_path or not os.path.isfile(image_path):
        return None

    img_w, img_h = get_dimensions(image_path)
    fx, fy, fw, fh = _contain_fit_rect(img_w, img_h, zone_rect)

    pic = slide.shapes.add_picture(
        image_path,
        int(round(fx)), int(round(fy)),
        width=int(round(fw)), height=int(round(fh)),
    )
    _hoist_picture_to_sptree_head(pic)
    return pic


def _add_annotation_connector(slide, p0, p1, color, width_emu, name):
    """Add a straight connector from p0 to p1 (EMU points) — python-pptx's
    ``add_connector`` takes explicit begin/end points, no flip arithmetic
    needed (§5.1, a genuine asymmetry with the JS path's bbox+flipV leader).
    """
    x0, y0 = p0
    x1, y1 = p1
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        int(round(x0)), int(round(y0)), int(round(x1)), int(round(y1)),
    )
    conn.line.color.rgb = color
    conn.line.width = int(round(width_emu))
    conn.name = name
    return conn


def _add_annotation_dot(slide, centre_pt, radius_emu, color, name):
    """Add a filled, borderless oval of the given radius centred on
    centre_pt (EMU) — used for both terminus dots and their casing rings
    (§5.2)."""
    cx, cy = centre_pt
    left = int(round(cx - radius_emu))
    top = int(round(cy - radius_emu))
    size = int(round(radius_emu * 2))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    oval.line.fill.background()
    oval.name = name
    return oval


def _add_annotation_label_box(slide, box_rect, text, fill_color, border_color,
                              border_width_emu, text_color, font_face,
                              font_size_pt, name):
    """Add a filled, bordered textbox centred in box_rect (EMU
    left/top/right/bottom) with the label text centred inside (§5.2)."""
    left, top, right, bottom = box_rect
    w = right - left
    h = bottom - top
    tb = slide.shapes.add_textbox(
        int(round(left)), int(round(top)), int(round(w)), int(round(h)))
    tb.fill.solid()
    tb.fill.fore_color.rgb = fill_color
    tb.line.color.rgb = border_color
    tb.line.width = int(round(border_width_emu))

    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.name = font_face
    run.font.color.rgb = text_color

    tb.name = name
    return tb


def _apply_native_annotation(slide, base_image_path, payload, slide_w, slide_h,
                             zone_rect_emu):
    """Render a native-annotated slide (§5): pure figure (F2) — a contain-fit
    base image with an editable overlay of leader connectors, casing
    underlays, terminus dots, and label textboxes, all resolved from a
    SlideAnnotations payload (annotations.schema.json) that already carries
    RESOLVED image-normalized coordinates. This function draws; it never
    re-runs placement.

    Modelled on ``_apply_full_bleed`` — same strip-all-shapes, single
    picture, no title/body/footer chrome — but the picture is contain-fit
    (§3.1/§3.2) rather than stretched, and connectors/ovals/textboxes are
    drawn on top afterward so they land after the picture in spTree
    (= on top in z-order), matching the pass order in §4.4:
    casing (leader + dot ring) -> dark cores (leader + dot) -> label boxes.

    Hash gate (F4, step 0): the picture is placed contain-fit regardless,
    then ``base_image_path``'s current content hash is compared against
    ``payload['base_image_hash']``. On mismatch the function returns
    immediately after placing the bare image — warn-refuse: a stale overlay
    on a changed image is worse than none. AN-01 (deck-qa) is the QA-side
    tripwire for this case.

    Args:
        slide: python-pptx Slide (already added via prs.slides.add_slide;
            chrome not yet stripped).
        base_image_path: absolute path to the CURRENT on-disk base image
            (already resolved relative to deck_dir by the caller).
        payload: the loaded SlideAnnotations dict.
        slide_w, slide_h: slide dimensions in EMU. Not used directly (the
            picture is sized from ``zone_rect_emu``) but kept for signature
            symmetry with ``_apply_full_bleed`` and to make full-slide
            callers' intent explicit at the call site.
        zone_rect_emu: (x, y, w, h) EMU rect of the placement zone the base
            image occupies. Always ``(0, 0, slide_w, slide_h)`` for every
            strategy v2 wires (§3.3 — full-slide strategies only; the
            composed/``annotated_image_zone`` zone is schema-allowed but
            deferred).

    Returns:
        The picture shape, or None if base_image_path is missing/absent
        (mirrors ``_apply_full_bleed``'s empty-slide edge case).
    """
    pic = _place_contain_fit_picture(slide, base_image_path, zone_rect_emu)
    if pic is None:
        return None

    current_hash = compute_content_hash(base_image_path)
    expected_hash = payload.get('base_image_hash')
    if current_hash != expected_hash:
        print(
            f"WARNING: annotate-figure v2 payload hash mismatch for slide "
            f"{payload.get('slide_number')} — expected {expected_hash}, "
            f"found {current_hash}. Refusing overlay; base image placed "
            f"without annotations.",
            file=sys.stderr,
        )
        return pic

    fit_rect = (pic.left, pic.top, pic.width, pic.height)
    style = payload['style']
    slide_number = payload['slide_number']

    leader_color = RGBColor.from_string(style['leader_color'])
    casing_color = RGBColor.from_string(style['casing_color'])
    box_fill = RGBColor.from_string(style['box_fill'])
    box_border = RGBColor.from_string(style['box_border'])
    text_color = RGBColor.from_string(style['text_color'])

    leader_width_emu = Pt(style['leader_width_pt'])
    casing_width_emu = Pt(style['casing_width_pt'])
    dot_radius_emu = Pt(style['dot_radius_pt'])
    box_border_width_emu = Pt(style['box_border_width_pt'])
    font_size_pt = style['font_size_pt']
    font_face = style['font_face']

    # Precompute every label's geometry up front (anchor point, the box
    # rect sized by the shared estimator, and the leader's termination
    # point on its own box edge nearest the anchor — §4.6) before drawing
    # anything, mirroring annotate_figure.annotate()'s geometry pass.
    geometry = []
    for i, label in enumerate(payload['labels']):
        anchor_pt = _map_norm_point(label['anchor'], fit_rect)
        label_centre_pt = _map_norm_point(label['label_pos'], fit_rect)

        box_w_in, box_h_in = estimate_label_box(label['text'], font_size_pt)
        box_w_emu, box_h_emu = Inches(box_w_in), Inches(box_h_in)
        box_left = label_centre_pt[0] - box_w_emu / 2
        box_top = label_centre_pt[1] - box_h_emu / 2
        box_rect = (box_left, box_top, box_left + box_w_emu, box_top + box_h_emu)

        leader_end = segment_box_entry(anchor_pt, label_centre_pt, box_rect)

        geometry.append({
            'index': i,
            'text': label['text'],
            'anchor_pt': anchor_pt,
            'box_rect': box_rect,
            'leader_end': leader_end,
        })

    draw_casing = style['casing_width_pt'] > 0
    # No schema field pins the casing ring's radius margin explicitly (v1's
    # PIL path uses a fixed-pixel DOT_CASING_EXTRA constant with no vector
    # equivalent in annotations.schema.json's style block) — the ring is
    # sized so its edge extends casing_width_pt/2 beyond the dot's own
    # edge, i.e. the same half-width the casing line extends beyond the
    # leader core, keeping the two casing elements visually consistent.
    ring_radius_emu = dot_radius_emu + casing_width_emu / 2

    # Pass 1a: ALL casing underlays (leader casing lines + dot casing
    # rings) first, so one leader's casing never cuts another's dark core.
    if draw_casing:
        for g in geometry:
            _add_annotation_connector(
                slide, g['anchor_pt'], g['leader_end'], casing_color,
                casing_width_emu, f"annotation_casing_{slide_number}_{g['index']}")
            _add_annotation_dot(
                slide, g['anchor_pt'], ring_radius_emu, casing_color,
                f"annotation_dotring_{slide_number}_{g['index']}")

    # Pass 1b: ALL dark leader cores + terminus dots, over the casing layer
    # but before every label box drawn in pass 2.
    for g in geometry:
        _add_annotation_connector(
            slide, g['anchor_pt'], g['leader_end'], leader_color,
            leader_width_emu, f"annotation_leader_{slide_number}_{g['index']}")
        _add_annotation_dot(
            slide, g['anchor_pt'], dot_radius_emu, leader_color,
            f"annotation_dot_{slide_number}_{g['index']}")

    # Pass 2: ALL label boxes + text, painted over the leader layer.
    for g in geometry:
        _add_annotation_label_box(
            slide, g['box_rect'], g['text'], box_fill, box_border,
            box_border_width_emu, text_color, font_face, font_size_pt,
            f"annotation_label_{slide_number}_{g['index']}")

    return pic


def _emit_smartart_placeholder(slide, slide_number, profile_layout):
    """Add a named rectangle placeholder for pptx_native SmartArt injection.

    Uses the content placeholder bounds from the profile layout.
    The existing assembler_patch.py finds this by name and replaces it.
    """
    content_info = None
    if profile_layout:
        for ph_info in profile_layout.get('placeholders', []):
            if ph_info['type'] == 'content':
                content_info = ph_info
                break

    if not content_info:
        x, y, w, h = Inches(0.6), Inches(2.3), Inches(12.13), Inches(4.57)
    else:
        x = Inches(content_info['x'])
        y = Inches(content_info['y'])
        w = Inches(content_info['w'])
        h = Inches(content_info['h'])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.name = f'pptx_native_placeholder_{slide_number}'
    shape.fill.background()
    shape.line.fill.background()


def build_deck(deck_dir, template_path, template_profile):
    """Assemble a .pptx deck using a template's slide layouts.

    Args:
        deck_dir: Path to the DeckContext directory.
        template_path: Path to the source .pptx template.
        template_profile: dict conforming to TemplateProfile schema.

    Returns:
        str: Path to the output .pptx file.
    """
    with open(os.path.join(deck_dir, 'outline.json')) as f:
        outline = json.load(f)

    with open(os.path.join(deck_dir, 'image-manifest.json')) as f:
        image_manifest = json.load(f)

    notes_path = os.path.join(deck_dir, 'speaker-notes.json')
    speaker_notes = {}
    if os.path.isfile(notes_path):
        with open(notes_path) as f:
            notes_data = json.load(f)
        speaker_notes = {n['slide_number']: n['text'] for n in notes_data.get('notes', [])}

    image_lookup = {}
    annotations_path_lookup = {}
    for img in image_manifest.get('images', []):
        if img.get('status') in ('generated', 'accepted', 'accepted_with_issues'):
            sn = img['slide_number']
            if sn not in image_lookup:
                image_lookup[sn] = img['file_path']
            if sn not in annotations_path_lookup and img.get('annotations_path'):
                annotations_path_lookup[sn] = img['annotations_path']

    # Build a set of slide numbers that need pptx_native SmartArt placeholders
    pptx_native_slides = set()
    sa_manifest_path = os.path.join(deck_dir, 'smartart-manifest.json')
    if os.path.isfile(sa_manifest_path):
        with open(sa_manifest_path) as f:
            sa_manifest = json.load(f)
        for graphic in sa_manifest.get('graphics', []):
            if graphic.get('engine_used') == 'pptx_native':
                pptx_native_slides.add(graphic['slide_number'])

    # Build a set of slide numbers marked full_bleed (issue #88), and a map
    # of slide numbers with an annotate-figure v2 annotation_mode (#142 v2,
    # T5). speaker_override wins when present, mirroring the JS assembler.
    # annotation_mode is orthogonal to speaker_override/strategy (§6.1) — it
    # is read as-is, not overridden.
    full_bleed_slides = set()
    annotation_slides = {}  # slide_number -> 'native' | 'raster'
    strategy_map_path = os.path.join(deck_dir, 'strategy-map.json')
    if os.path.isfile(strategy_map_path):
        with open(strategy_map_path) as f:
            strategy_map = json.load(f)
        for entry in strategy_map.get('slides', []):
            effective = entry.get('speaker_override') or entry.get('strategy')
            if effective == 'full_bleed':
                full_bleed_slides.add(entry['slide_number'])
            annotation_mode = entry.get('annotation_mode')
            if annotation_mode in ('native', 'raster'):
                annotation_slides[entry['slide_number']] = annotation_mode

    prs = Presentation(template_path)
    _strip_existing_slides(prs)
    slide_w, slide_h = prs.slide_width, prs.slide_height

    for slide_data in outline.get('slides', []):
        slide_number = slide_data['slide_number']
        slide_type = slide_data.get('slide_type', 'content')
        headline = slide_data.get('headline', '')
        body_points = slide_data.get('body_points', [])

        layout = _resolve_layout(prs, template_profile, slide_type)
        slide = prs.slides.add_slide(layout)

        # annotate-figure v2 (#142 v2) short-circuits all standard
        # population, same as full_bleed: pure figure (F2), no headline,
        # no body, no footer. Checked BEFORE full_bleed so a slide whose
        # base strategy is full_bleed but which also carries
        # annotation_mode gets the dedicated contain-fit + overlay builder
        # instead of the plain stretch-to-canvas path (§3.1).
        if slide_number in annotation_slides:
            mode = annotation_slides[slide_number]
            abs_image_path = None
            if slide_number in image_lookup:
                raw = image_lookup[slide_number]
                abs_image_path = raw if os.path.isabs(raw) else os.path.join(deck_dir, raw)

            if mode == 'native':
                payload = None
                rel_ann_path = annotations_path_lookup.get(slide_number)
                if rel_ann_path:
                    abs_ann_path = (rel_ann_path if os.path.isabs(rel_ann_path)
                                    else os.path.join(deck_dir, rel_ann_path))
                    if os.path.isfile(abs_ann_path):
                        with open(abs_ann_path) as f:
                            payload = json.load(f)
                if payload is None:
                    # Payload absent (F3c/F5) — place contain-fit without an
                    # overlay and warn; AN-01 raises the error at QA time.
                    print(
                        f"WARNING: slide {slide_number} has "
                        f"annotation_mode=native but no annotations payload "
                        f"was found — placing base image without overlay.",
                        file=sys.stderr,
                    )
                    _place_contain_fit_picture(slide, abs_image_path, (0, 0, slide_w, slide_h))
                else:
                    _apply_native_annotation(
                        slide, abs_image_path, payload, slide_w, slide_h,
                        (0, 0, slide_w, slide_h))
            else:
                # 'raster' — labels already baked into pixels (F11);
                # contain-fit placement only, no overlay renderer.
                _place_contain_fit_picture(slide, abs_image_path, (0, 0, slide_w, slide_h))

            if slide_number in speaker_notes:
                slide.notes_slide.notes_text_frame.text = speaker_notes[slide_number]
            continue

        # full_bleed short-circuits all standard population. Picture only,
        # zero chrome. Issue #88.
        if slide_number in full_bleed_slides:
            abs_image_path = None
            if slide_number in image_lookup:
                raw = image_lookup[slide_number]
                abs_image_path = raw if os.path.isabs(raw) else os.path.join(deck_dir, raw)
            _apply_full_bleed(slide, abs_image_path, slide_w, slide_h)
            if slide_number in speaker_notes:
                slide.notes_slide.notes_text_frame.text = speaker_notes[slide_number]
            continue

        layout_name = _get_mapped_layout_name(template_profile, slide_type)
        profile_layout = _get_profile_layout(template_profile, layout_name)

        # Populate title
        title_ph = _find_placeholder_by_type(slide, 'title', profile_layout)
        _populate_text(title_ph, headline)

        # Populate body/content
        content_ph = _find_placeholder_by_type(slide, 'content', profile_layout)
        body_ph = _find_placeholder_by_type(slide, 'body', profile_layout)
        if content_ph and body_points:
            _populate_body_points(content_ph, body_points)
        elif body_ph and body_points:
            _populate_body_points(body_ph, body_points)

        # SmartArt placeholder for pptx_native injection
        if slide_number in pptx_native_slides:
            _emit_smartart_placeholder(slide, slide_number, profile_layout)

        # Populate picture placeholder
        if slide_number in image_lookup:
            pic_ph = _find_placeholder_by_type(slide, 'picture', profile_layout)
            if pic_ph:
                image_path = image_lookup[slide_number]
                abs_image_path = os.path.join(deck_dir, image_path) if not os.path.isabs(image_path) else image_path
                if os.path.isfile(abs_image_path):
                    pic_ph.insert_picture(abs_image_path)

        # Speaker notes
        if slide_number in speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes[slide_number]

    output_dir = os.path.join(deck_dir, 'output')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'presentation.pptx')
    prs.save(output_path)

    return output_path
