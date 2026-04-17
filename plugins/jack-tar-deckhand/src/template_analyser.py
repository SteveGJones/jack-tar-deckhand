"""Template analyser — extract layouts and placeholders from a .pptx template.

Reads a corporate .pptx template via python-pptx and produces a TemplateProfile
JSON structure for template-driven assembly.
"""

import hashlib
import json
import os
import re

from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400


def _emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    return round(emu / EMU_PER_INCH, 2)


# Map OOXML placeholder type strings to our vocabulary
_TYPE_MAP = {
    'TITLE': 'title',
    'CENTER_TITLE': 'title',
    'SUBTITLE': 'subtitle',
    'BODY': 'body',
    'OBJECT': 'content',
    'PICTURE': 'picture',
    'CHART': 'chart',
    'TABLE': 'table',
    'DATE': 'date',
    'FOOTER': 'footer',
    'SLIDE_NUMBER': 'slide_number',
}

# Name patterns that override type-based classification
_NAME_OVERRIDES = [
    (re.compile(r'chapter', re.IGNORECASE), 'chapter_box'),
    (re.compile(r'logo', re.IGNORECASE), 'other'),
]


def classify_placeholder(ph_type_str, ph_name):
    """Classify a placeholder into our type vocabulary.

    Args:
        ph_type_str: String like 'TITLE (1)' or 'BODY (2)' from python-pptx.
        ph_name: The shape name string from the template.

    Returns:
        One of: title, subtitle, body, content, picture, chart, table,
                footer, slide_number, date, chapter_box, other.
    """
    # Check name-based overrides first
    for pattern, override_type in _NAME_OVERRIDES:
        if pattern.search(ph_name):
            return override_type

    # Extract the type keyword before the parenthesized number
    type_key = ph_type_str.split('(')[0].strip().split()[-1] if ph_type_str else ''
    return _TYPE_MAP.get(type_key, 'other')


# Convention table: (regex pattern on layout name, slide_type)
_CONVENTION_TABLE = [
    (re.compile(r'^Title Slide', re.IGNORECASE), 'title'),
    (re.compile(r'^Title Only', re.IGNORECASE), 'title'),
    (re.compile(r'^Divider|^Section', re.IGNORECASE), 'section_divider'),
    (re.compile(r'^Content.*Photo|^Content.*Image', re.IGNORECASE), 'content_with_image'),
    (re.compile(r'^Comparison', re.IGNORECASE), 'comparison'),
    (re.compile(r'^Conclusion|^End', re.IGNORECASE), 'closing'),
    (re.compile(r'^Agenda', re.IGNORECASE), 'agenda'),
    (re.compile(r'^Text', re.IGNORECASE), 'quote'),
    (re.compile(r'^Blank$', re.IGNORECASE), 'blank'),
    # Content N patterns — must come after Content*Photo
    (re.compile(r'^Content\s*1\b', re.IGNORECASE), 'content'),
    (re.compile(r'^Content\s*[2-8]\b', re.IGNORECASE), 'two_column'),
]


def _has_picture_placeholder(layout):
    return any(p['type'] == 'picture' for p in layout.get('placeholders', []))


def _largest_content_area(layout):
    """Return the area of the largest single content placeholder in the layout."""
    max_area = 0
    for p in layout.get('placeholders', []):
        if p['type'] == 'content':
            area = p['w'] * p['h']
            if area > max_area:
                max_area = area
    return max_area


def auto_map_layouts(layouts):
    """Auto-map extracted layouts to slide types by naming convention.

    Args:
        layouts: List of layout dicts from extract_layouts().

    Returns:
        Tuple of (mapping_dict, fallback_dict).
        mapping_dict: {slide_type: {'layout_name': str, 'layout_index': int}}
        fallback_dict: {'layout_name': str, 'layout_index': int} or None.
    """
    if not layouts:
        return {}, None

    # Collect candidates: {slide_type: [layout, ...]}
    candidates = {}
    for layout in layouts:
        name = layout['name']
        for pattern, slide_type in _CONVENTION_TABLE:
            if pattern.search(name):
                candidates.setdefault(slide_type, []).append(layout)
                break
        else:
            # Check if it's a Content layout with a PICTURE placeholder
            if re.search(r'^Content', name, re.IGNORECASE) and _has_picture_placeholder(layout):
                candidates.setdefault('content_with_image', []).append(layout)

    # For each slide type, pick the simplest layout (fewest placeholders)
    mapping = {}
    for slide_type, options in candidates.items():
        options.sort(key=lambda l: l['placeholder_count'])
        winner = options[0]
        mapping[slide_type] = {
            'layout_name': winner['name'],
            'layout_index': winner['index'],
        }

    # Fallback: layout with the largest single content placeholder
    best_fallback = None
    best_area = 0
    for layout in layouts:
        area = _largest_content_area(layout)
        if area > best_area:
            best_area = area
            best_fallback = layout

    fallback = None
    if best_fallback:
        fallback = {
            'layout_name': best_fallback['name'],
            'layout_index': best_fallback['index'],
        }

    return mapping, fallback


def _compute_file_hash(file_path):
    """Compute SHA-256 hash of a file."""
    sha256 = hashlib.sha256()
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            sha256.update(chunk)
    return sha256.hexdigest()


def analyse_template(template_path, master_index=0):
    """Analyse a .pptx template and produce a TemplateProfile.

    Args:
        template_path: Path to the .pptx template file.
        master_index: Which slide master to use (default 0).

    Returns:
        dict conforming to the TemplateProfile schema.
    """
    prs = Presentation(template_path)

    layouts = extract_layouts(template_path, master_index=master_index)
    mapping, fallback = auto_map_layouts(layouts)

    profile = {
        'template_path': template_path,
        'template_hash': _compute_file_hash(template_path),
        'slide_width_inches': _emu_to_inches(prs.slide_width),
        'slide_height_inches': _emu_to_inches(prs.slide_height),
        'master_index': master_index,
        'layouts': layouts,
        'layout_mapping': mapping,
        'unmapped_fallback': fallback,
        'constrains_strategies': True,
        'speaker_approved': False,
    }

    return profile


def extract_layouts(template_path, master_index=0):
    """Extract all slide layouts from a template .pptx file.

    Args:
        template_path: Path to the .pptx template file.
        master_index: Which slide master to use (default 0).

    Returns:
        List of layout dicts with name, index, placeholders, and decorative shape count.
    """
    prs = Presentation(template_path)
    master = prs.slide_masters[master_index]
    results = []

    for i, layout in enumerate(master.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type_str = str(ph.placeholder_format.type) if ph.placeholder_format.type is not None else ''
            placeholders.append({
                'idx': ph.placeholder_format.idx,
                'type': classify_placeholder(ph_type_str, ph.name),
                'name': ph.name,
                'x': _emu_to_inches(ph.left),
                'y': _emu_to_inches(ph.top),
                'w': _emu_to_inches(ph.width),
                'h': _emu_to_inches(ph.height),
            })

        non_placeholder_shapes = [s for s in layout.shapes if not s.is_placeholder]

        results.append({
            'name': layout.name,
            'index': i,
            'placeholder_count': len(placeholders),
            'placeholders': placeholders,
            'decorative_shape_count': len(non_placeholder_shapes),
        })

    return results
