"""Tests for the annotate-figure v2 native-annotation builder in the
python-pptx template assembler path (issue #142 v2, T5).

Design: docs/superpowers/plans/2026-07-17-annotate-figure-v2.md §5, §8.3.

Covers:
- ``_apply_native_annotation``: chrome-stripping (pure figure, F2), exact
  shape counts per name prefix (F10), verbatim label text (AN-02's whole
  point), contain-fit anchor mapping (§3.2) including the off-aspect
  letterbox case, the F4 hash-mismatch warn-refuse, and spTree z-order
  (overlay shapes land after the picture).
- ``build_deck`` wiring: native end-to-end (with the template fixture),
  raster-mode contain-fit routing (F11), and backward compatibility when
  no ``annotation_mode`` is present (``full_bleed`` still stretches, as
  pinned by the sibling ``test_full_bleed_scale.py``).

Do not ``Read`` PNG files directly — every image fixture here is a
synthetic solid-colour PIL image built on disk in ``tmp_path``; only OOXML
geometry and text are asserted, never pixels.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
PROJECT_ROOT = PLUGIN_ROOT.parent.parent
sys.path.insert(0, str(PLUGIN_ROOT))

from src.annotation_payload import (  # noqa: E402
    build_annotation_payload,
    write_annotation_payload,
)
from src.assembler.build_deck_template import (  # noqa: E402
    _apply_native_annotation,
    _place_contain_fit_picture,
    _resolve_annotation_zone_rect,
    build_deck,
)

TEMPLATE_FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "templates" / "metamirror-template.pptx"
HERO_IMAGE_FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "minimal_deck" / "images" / "slide-01-hero.png"

# 16:9 EMU zone, matching the fixture template's canvas.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FULL_SLIDE_ZONE = (0, 0, SLIDE_W, SLIDE_H)

EMU_TOLERANCE = 5000  # ~0.005in — comfortably above int-rounding noise


# --- helpers -----------------------------------------------------------------


def _make_image(path, size, color=(40, 40, 40)):
    Image.new("RGB", size, color).save(path)
    return str(path)


def _blank_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return prs, slide


def _shapes_named(slide, prefix):
    return [s for s in slide.shapes if s.name.startswith(prefix)]


def _count_named(slide, prefix):
    return len(_shapes_named(slide, prefix))


# --- _apply_native_annotation: chrome + pure figure (F2) --------------------


def test_apply_native_annotation_strips_chrome_then_adds_picture(tmp_path):
    prs, slide = _blank_slide()
    chrome = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    chrome.text_frame.text = "Pre-existing chrome must not survive"
    assert len(slide.shapes) == 1

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=1, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors={"Rudder": [0.5, 0.5]},
    )

    pic = _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    assert pic is not None
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1

    # No pre-existing chrome, and no headline textbox was introduced — every
    # remaining text_frame-bearing shape is an annotation label.
    text_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert all(s.name.startswith("annotation_label_") for s in text_shapes)
    assert all("Pre-existing chrome" not in (s.text_frame.text or "") for s in text_shapes)


def test_apply_native_annotation_refuses_overlay_when_image_missing(tmp_path):
    prs, slide = _blank_slide()
    payload = {
        "slide_number": 1, "source": "generated",
        "base_image_path": str(tmp_path / "nope.png"),
        "base_image_hash": "0" * 64,
        "image_dimensions": {"width": 100, "height": 100},
        "placement_zone": "annotated_full_slide", "fit": "contain",
        "labels": [{"text": "X", "anchor": [0.5, 0.5], "label_pos": [0.5, 0.1]}],
        "style": {
            "leader_width_pt": 1.5, "casing_width_pt": 3.5, "casing_color": "FFFFFF",
            "leader_color": "141414", "dot_radius_pt": 3.0, "box_fill": "FFFFFF",
            "box_border": "141414", "box_border_width_pt": 1.0, "text_color": "141414",
            "font_face": "Calibri", "font_size_pt": 18,
        },
    }
    pic = _apply_native_annotation(
        slide, str(tmp_path / "nope.png"), payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)
    assert pic is None
    assert len(slide.shapes) == 0


# --- F10: exact shape counts per name prefix --------------------------------


def test_apply_native_annotation_emits_exact_shape_counts(tmp_path):
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))
    anchors = {"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15], "Keel": [0.5, 0.95]}
    payload = build_annotation_payload(
        slide_number=2, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors=anchors,
    )
    assert payload["style"]["casing_width_pt"] > 0  # sanity: casing is on by default

    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    n = len(anchors)
    assert _count_named(slide, "annotation_label_") == n
    assert _count_named(slide, "annotation_leader_") == n
    assert _count_named(slide, "annotation_dot_") == n
    assert _count_named(slide, "annotation_casing_") == n
    assert _count_named(slide, "annotation_dotring_") == n

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    # Total shape count is exactly picture + 5 overlay shapes per label.
    assert len(slide.shapes) == 1 + 5 * n


def test_apply_native_annotation_no_casing_when_casing_width_zero(tmp_path):
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=3, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors={"Rudder": [0.8, 0.7]},
        style_overrides={"casing_width_pt": 0},
    )
    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    assert _count_named(slide, "annotation_casing_") == 0
    assert _count_named(slide, "annotation_dotring_") == 0
    assert _count_named(slide, "annotation_leader_") == 1
    assert _count_named(slide, "annotation_dot_") == 1
    assert _count_named(slide, "annotation_label_") == 1


# --- AN-02's whole point: verbatim label text -------------------------------


def test_apply_native_annotation_label_text_is_verbatim(tmp_path):
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))
    anchors = {
        "Bow thruster (fwd)": [0.15, 0.5],
        "Mizzen — aft mast": [0.85, 0.2],
    }
    payload = build_annotation_payload(
        slide_number=4, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors=anchors,
    )
    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    label_texts = sorted(
        s.text_frame.text for s in _shapes_named(slide, "annotation_label_")
    )
    assert label_texts == sorted(anchors)


# --- §3.2 contain-fit coordinate mapping ------------------------------------


def test_apply_native_annotation_maps_anchor_into_fitted_rect(tmp_path):
    """A centre anchor always maps to the centre of the placement zone,
    regardless of any letterbox — the fit rect is always zone-centred on
    whichever axis has slack, so [0.5, 0.5] -> zone centre exactly."""
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))  # 16:9, ~matches zone
    payload = build_annotation_payload(
        slide_number=5, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors={"Centre": [0.5, 0.5]},
    )
    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    dots = _shapes_named(slide, "annotation_dot_")
    assert len(dots) == 1
    dot = dots[0]
    dot_cx = dot.left + dot.width / 2
    dot_cy = dot.top + dot.height / 2

    assert dot_cx == pytest.approx(SLIDE_W / 2, abs=EMU_TOLERANCE)
    assert dot_cy == pytest.approx(SLIDE_H / 2, abs=EMU_TOLERANCE)


def test_apply_native_annotation_letterboxes_off_aspect_image(tmp_path):
    """A square (1:1) image inside a 16:9 zone fits to height, letterboxing
    left/right — anchor [0, 0] must land on the fitted rect's left EDGE,
    not the slide's true origin."""
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (900, 900))  # 1:1
    payload = build_annotation_payload(
        slide_number=6, source="generated", base_image_path=base,
        image_dimensions={"width": 900, "height": 900},
        placement_zone="annotated_full_slide", anchors={"Corner": [0.0, 0.0]},
    )
    pic = _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    # Fitted to height: picture height == zone height, width narrower than
    # the zone (letterboxed left/right).
    assert pic.height == pytest.approx(SLIDE_H, abs=EMU_TOLERANCE)
    assert pic.width < SLIDE_W - EMU_TOLERANCE
    assert pic.left > EMU_TOLERANCE  # NOT flush with the slide's true x=0 origin

    dots = _shapes_named(slide, "annotation_dot_")
    dot = dots[0]
    dot_cx = dot.left + dot.width / 2
    dot_cy = dot.top + dot.height / 2

    # [0,0] maps to the fitted rect's top-left corner, i.e. (pic.left, pic.top)
    # — which is NOT the slide origin because of the letterbox.
    assert dot_cx == pytest.approx(pic.left, abs=EMU_TOLERANCE)
    assert dot_cy == pytest.approx(pic.top, abs=EMU_TOLERANCE)
    assert dot_cx > EMU_TOLERANCE


# --- F4: hash-mismatch warn-refuse ------------------------------------------


def test_apply_native_annotation_refuses_overlay_on_hash_mismatch(tmp_path, capsys):
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=7, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors={"Rudder": [0.8, 0.7]},
    )

    # Mutate the base image AFTER the payload was built — the anchors are
    # now stale relative to the on-disk file's content hash.
    _make_image(tmp_path / "base.png", (1920, 1080), color=(210, 15, 15))

    pic = _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    assert pic is not None
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    # Zero annotation_* shapes of any kind — warn-refuse, not partial draw.
    annotation_shapes = [s for s in slide.shapes if s.name.startswith("annotation_")]
    assert annotation_shapes == []
    assert len(slide.shapes) == 1  # picture only

    err = capsys.readouterr().err
    assert "hash mismatch" in err.lower()


# --- spTree z-order: overlay shapes land after the picture ------------------


def test_apply_native_annotation_z_order_labels_after_picture(tmp_path):
    prs, slide = _blank_slide()
    base = _make_image(tmp_path / "base.png", (1920, 1080))
    anchors = {"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15]}
    payload = build_annotation_payload(
        slide_number=8, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide", anchors=anchors,
    )
    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, FULL_SLIDE_ZONE)

    sp_tree = slide.shapes._spTree
    children = list(sp_tree)

    def _index_of(shape):
        return children.index(shape._element)

    pic_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pic_shapes) == 1
    pic_index = _index_of(pic_shapes[0])

    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_",
                   "annotation_casing_", "annotation_dotring_"):
        for shape in _shapes_named(slide, prefix):
            assert _index_of(shape) > pic_index


# --- build_deck end-to-end wiring -------------------------------------------


def _stub_template_profile():
    """Minimal TemplateProfile pointing at the fixture template (mirrors
    test_full_bleed_scale.py's helper of the same name)."""
    return {
        "master_index": 0,
        "layout_mapping": {
            "content": {"layout_index": 1, "layout_name": "Title and Content"},
            "title": {"layout_index": 0, "layout_name": "Title Slide"},
        },
        "unmapped_fallback": {"layout_index": 1, "layout_name": "Title and Content"},
        "layouts": [
            {
                "name": "Title and Content",
                "placeholders": [
                    {"idx": 0, "type": "title", "name": "Title", "x": 0.5, "y": 0.3, "w": 12.0, "h": 1.0},
                    {"idx": 1, "type": "content", "name": "Body", "x": 0.5, "y": 1.5, "w": 12.0, "h": 5.5},
                ],
            },
            {
                "name": "Title Slide",
                "placeholders": [
                    {"idx": 0, "type": "title", "name": "Title", "x": 0.5, "y": 3.0, "w": 12.0, "h": 1.5},
                ],
            },
        ],
    }


def _write_outline(deck_dir, second_slide_headline="This headline must not appear"):
    outline = {
        "narrative_arc": "infographic-narrative",
        "estimated_duration_minutes": 6,
        "total_slides": 2,
        "slides": [
            {
                "slide_number": 1, "slide_type": "content",
                "headline": "Composed slide retains chrome",
                "body_points": ["alpha", "beta"], "visual_type": "none",
                "layout_template": "content",
            },
            {
                "slide_number": 2, "slide_type": "content",
                "headline": second_slide_headline,
                "body_points": ["this body must not appear"],
                "visual_type": "hero_image", "layout_template": "content",
            },
        ],
    }
    with open(deck_dir / "outline.json", "w") as f:
        json.dump(outline, f)


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_native_annotation_end_to_end(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_outline(deck_dir)

    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [
            {
                "image_id": "slide-02-hero", "slide_number": 2,
                "file_path": str(HERO_IMAGE_FIXTURE),
                "placement_zone": "annotated_full_slide",
                "annotations_path": "annotations/slide-02-annotations.json",
                "dimensions": {"width": 1920, "height": 1080},
                "source_prompt": "test", "model_used": "test", "alt_text": "hero",
                "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0,
            },
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    payload = build_annotation_payload(
        slide_number=2, source="generated", base_image_path=str(HERO_IMAGE_FIXTURE),
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide",
        anchors={"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15]},
    )
    write_annotation_payload(str(deck_dir), 2, payload)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "keep chrome",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "annotated figure",
             "render_funnel": ["ollama", "cloud_low", "cloud_full"],
             "speaker_override": None, "annotation_mode": "native",
             "annotation": {"labels": [
                 {"text": "Rudder", "target": "the rudder"},
                 {"text": "Mainsail", "target": "the mainsail"},
             ]}},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    assert len(prs.slides) == 2

    composed_slide = prs.slides[0]
    has_title_text = any(
        getattr(s, "has_text_frame", False)
        and "Composed slide retains chrome" in (s.text_frame.text or "")
        for s in composed_slide.shapes
    )
    assert has_title_text, "sibling composed slide should keep its title text"

    native_slide = prs.slides[1]
    pictures = [s for s in native_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1

    label_texts = sorted(
        s.text_frame.text for s in native_slide.shapes if s.name.startswith("annotation_label_")
    )
    assert label_texts == ["Mainsail", "Rudder"]

    # Pure figure (F2): the outline's headline/body strings for slide 2 are
    # absent from the assembled slide.
    all_text = " ".join(
        s.text_frame.text for s in native_slide.shapes
        if getattr(s, "has_text_frame", False) and s.text_frame.text
    )
    assert "This headline must not appear" not in all_text
    assert "this body must not appear" not in all_text


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_raster_annotation_contains_not_stretches(tmp_path):
    """F11: a raster-mode annotated slide's picture is contain-fit (aspect
    preserved), not stretched to canvas like a plain full_bleed slide."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_outline(deck_dir, second_slide_headline="raster slide headline")

    images_dir = deck_dir / "images"
    images_dir.mkdir()
    _make_image(images_dir / "slide-02-square.png", (900, 900))

    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [
            {
                "image_id": "slide-02-hero", "slide_number": 2,
                "file_path": "images/slide-02-square.png",
                "placement_zone": "annotated_full_slide",
                "dimensions": {"width": 900, "height": 900},
                "source_prompt": "test", "model_used": "test", "alt_text": "square",
                "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0,
            },
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "keep chrome",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "annotated figure",
             "render_funnel": ["ollama"], "speaker_override": None,
             "annotation_mode": "raster",
             "annotation": {"labels": [{"text": "Rudder", "target": "the rudder"}]}},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    raster_slide = prs.slides[1]

    pictures = [s for s in raster_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]

    # Contain-fit of a 1:1 image in the 16:9 canvas: height matches slide
    # height, width is narrower — NOT stretched to slide_w x slide_h like
    # plain full_bleed.
    assert pic.height == pytest.approx(prs.slide_height, abs=EMU_TOLERANCE)
    assert pic.width < prs.slide_width - EMU_TOLERANCE
    # No overlay shapes — raster mode has no payload, labels are baked into
    # the pixels already.
    assert not any(s.name.startswith("annotation_") for s in raster_slide.shapes)


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_without_annotation_unchanged(tmp_path):
    """Backward-compat (F2/F11 non-regression): a full_bleed slide with NO
    annotation_mode key still stretches to canvas exactly as before."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_outline(deck_dir, second_slide_headline="plain full_bleed headline")

    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [
            {
                "image_id": "slide-02-hero", "slide_number": 2,
                "file_path": str(HERO_IMAGE_FIXTURE),
                "placement_zone": "background",
                "dimensions": {"width": 1920, "height": 1080},
                "source_prompt": "test", "model_used": "test", "alt_text": "hero",
                "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0,
            },
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "keep chrome",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "infographic register",
             "render_funnel": ["ollama", "cloud_low", "cloud_full"], "speaker_override": None},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    full_bleed_slide = prs.slides[1]

    pictures = [s for s in full_bleed_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    non_pictures = [s for s in full_bleed_slide.shapes if s.shape_type != MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert len(non_pictures) == 0
    # Stretched (not contain-fit): exact canvas size, no letterbox.
    assert pictures[0].left == 0 and pictures[0].top == 0
    assert pictures[0].width == prs.slide_width
    assert pictures[0].height == prs.slide_height


# =============================================================================
# v2.1 Feature A — composed strategy wiring (issue #142 v2.1, T4)
#
# Design: docs/superpowers/plans/2026-07-23-annotate-figure-v2.1.md §2.3, §5.2.
#
# These operate directly on _resolve_annotation_zone_rect / _apply_native_
# annotation / _place_contain_fit_picture against python-pptx's bundled
# default template (no external fixture needed — its "Picture with Caption"
# layout at index 8 has a real PICTURE placeholder with the same idx layout
# as the project's own template fixture; "Title and Content" at index 1 has
# none, exercising the F-06 fallback).
# =============================================================================


def _slide_with_layout(layout_index):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    return prs, slide


def _picture_profile_layout():
    """Profile-layout metadata for python-pptx's bundled 'Picture with
    Caption' layout (index 8): title idx0, picture idx1, body idx2 — the
    F-06 happy path (step 1: a real picture placeholder is found)."""
    return {
        "name": "Picture with Caption",
        "placeholders": [
            {"idx": 0, "type": "title", "name": "Title", "x": 1.96, "y": 5.25, "w": 6.0, "h": 0.62},
            {"idx": 1, "type": "picture", "name": "Picture", "x": 1.96, "y": 0.67, "w": 6.0, "h": 4.5},
            {"idx": 2, "type": "body", "name": "Text", "x": 1.96, "y": 5.87, "w": 6.0, "h": 0.88},
        ],
    }


def _content_only_profile_layout():
    """Profile-layout metadata mirroring _stub_template_profile()'s 'Title
    and Content' layout (index 1): title + content, NO picture type
    declared -- exercises the F-06 fallback chain's step 2."""
    return {
        "name": "Title and Content",
        "placeholders": [
            {"idx": 0, "type": "title", "name": "Title", "x": 0.5, "y": 0.3, "w": 12.0, "h": 1.0},
            {"idx": 1, "type": "content", "name": "Body", "x": 0.5, "y": 1.5, "w": 12.0, "h": 5.5},
        ],
    }


def test_composed_native_retains_headline_and_body(tmp_path):
    prs, slide = _slide_with_layout(8)
    profile_layout = _picture_profile_layout()
    slide.placeholders[0].text_frame.text = "Composed headline"
    slide.placeholders[2].text_frame.text = "body point one"

    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 10)

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=10, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors={"Rudder": [0.5, 0.5]},
    )
    _apply_native_annotation(slide, base, payload, SLIDE_W, SLIDE_H, zone_rect, strip_chrome=False)

    assert "Composed headline" in slide.placeholders[0].text_frame.text
    assert "body point one" in slide.placeholders[2].text_frame.text
    assert _count_named(slide, "annotation_label_") == 1
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_composed_native_maps_anchor_into_image_zone(tmp_path):
    """Anchor [0.5, 0.5] on a 16:9 image maps to the CENTRE of the
    picture-placeholder rect (EMU tolerance), NOT the slide centre."""
    prs, slide = _slide_with_layout(8)
    profile_layout = _picture_profile_layout()

    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 11)
    zx, zy, zw, zh = zone_rect
    zone_cx, zone_cy = zx + zw / 2, zy + zh / 2
    # Sanity: the picture-zone centre is NOT the slide centre (the fixture
    # layout's picture placeholder is off-centre / smaller than the canvas).
    assert abs(zone_cx - prs.slide_width / 2) > EMU_TOLERANCE or \
        abs(zone_cy - prs.slide_height / 2) > EMU_TOLERANCE

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=11, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors={"Centre": [0.5, 0.5]},
    )
    _apply_native_annotation(slide, base, payload, prs.slide_width, prs.slide_height,
                             zone_rect, strip_chrome=False)

    dots = _shapes_named(slide, "annotation_dot_")
    assert len(dots) == 1
    dot = dots[0]
    dot_cx = dot.left + dot.width / 2
    dot_cy = dot.top + dot.height / 2
    assert dot_cx == pytest.approx(zone_cx, abs=EMU_TOLERANCE)
    assert dot_cy == pytest.approx(zone_cy, abs=EMU_TOLERANCE)


def test_composed_native_removes_empty_picture_placeholder(tmp_path):
    prs, slide = _slide_with_layout(8)
    profile_layout = _picture_profile_layout()
    assert any(ph.placeholder_format.idx == 1 for ph in slide.placeholders)

    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 12)

    # The picture placeholder (idx 1) is gone -- no unfilled placeholder ships.
    assert not any(ph.placeholder_format.idx == 1 for ph in slide.placeholders)

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=12, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors={"Rudder": [0.5, 0.5]},
    )
    _apply_native_annotation(slide, base, payload, prs.slide_width, prs.slide_height,
                             zone_rect, strip_chrome=False)

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_composed_raster_places_contain_fit_in_zone_no_overlay(tmp_path):
    prs, slide = _slide_with_layout(8)
    profile_layout = _picture_profile_layout()
    slide.placeholders[0].text_frame.text = "Composed headline"

    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 13)
    base = _make_image(tmp_path / "square.png", (900, 900))
    _place_contain_fit_picture(slide, base, zone_rect, strip=False)

    assert "Composed headline" in slide.placeholders[0].text_frame.text
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert not any(s.name.startswith("annotation_") for s in slide.shapes)


def test_composed_native_exact_shape_counts(tmp_path):
    prs, slide = _slide_with_layout(8)
    profile_layout = _picture_profile_layout()
    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 14)

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    anchors = {"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15], "Keel": [0.5, 0.95]}
    payload = build_annotation_payload(
        slide_number=14, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors=anchors,
    )
    _apply_native_annotation(slide, base, payload, prs.slide_width, prs.slide_height,
                             zone_rect, strip_chrome=False)

    n = len(anchors)
    assert _count_named(slide, "annotation_label_") == n
    assert _count_named(slide, "annotation_leader_") == n
    assert _count_named(slide, "annotation_dot_") == n
    assert _count_named(slide, "annotation_casing_") == n
    assert _count_named(slide, "annotation_dotring_") == n


def test_composed_native_no_picture_placeholder_falls_back(tmp_path, capsys):
    """F-06: a layout with no picture placeholder falls back to the content
    placeholder rect, with a warning, and the overlay still draws."""
    prs, slide = _slide_with_layout(1)  # "Title and Content" -- no picture ph
    profile_layout = _content_only_profile_layout()

    zone_rect = _resolve_annotation_zone_rect(slide, profile_layout, 15)
    assert zone_rect == (Inches(0.5), Inches(1.5), Inches(12.0), Inches(5.5))

    err = capsys.readouterr().err
    assert "slide 15" in err and "content placeholder" in err.lower()

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    payload = build_annotation_payload(
        slide_number=15, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors={"Rudder": [0.5, 0.5]},
    )
    _apply_native_annotation(slide, base, payload, prs.slide_width, prs.slide_height,
                             zone_rect, strip_chrome=False)

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert _count_named(slide, "annotation_label_") == 1


# --- F-03: composed annotated slide with a non-content slide_type ----------


def _template_path_from_default(tmp_path):
    """A synthetic .pptx template built from python-pptx's own bundled
    default presentation -- avoids depending on the project's template
    fixture for tests that only need 'Picture with Caption' (index 8)."""
    path = tmp_path / "default-template.pptx"
    Presentation().save(str(path))
    return str(path)


def test_composed_native_diagram_slide_type_uses_content_chrome(tmp_path):
    """F-03 (option b): a composed native slide with slide_type='diagram'
    still renders through the standard content-chrome population (title +
    body populated, overlay drawn) -- python's build_deck has no separate
    diagram-specific builder to accidentally bypass, so this pins that a
    'diagram' slide_type composed-annotated slide is never silently
    short-circuited or left unannotated."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    template_path = _template_path_from_default(tmp_path)

    outline = {
        "narrative_arc": "test", "estimated_duration_minutes": 4, "total_slides": 1,
        "slides": [
            {"slide_number": 1, "slide_type": "diagram", "headline": "Diagram headline",
             "body_points": ["diagram body point"], "visual_type": "hero_image",
             "layout_template": "diagram"},
        ],
    }
    with open(deck_dir / "outline.json", "w") as f:
        json.dump(outline, f)

    base = _make_image(tmp_path / "base.png", (1920, 1080))
    image_manifest = {
        "generated_at": "2026-07-23T00:00:00Z", "image_backend": "ollama",
        "images": [
            {"image_id": "slide-01-figure", "slide_number": 1, "file_path": base,
             "placement_zone": "annotated_image_zone",
             "annotations_path": "annotations/slide-01-annotations.json",
             "dimensions": {"width": 1920, "height": 1080},
             "source_prompt": "test", "model_used": "test", "alt_text": "figure",
             "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0},
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    payload = build_annotation_payload(
        slide_number=1, source="generated", base_image_path=base,
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone", anchors={"Rudder": [0.5, 0.5]},
    )
    write_annotation_payload(str(deck_dir), 1, payload)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "diagram figure",
             "render_funnel": ["ollama"], "speaker_override": None,
             "annotation_mode": "native",
             "annotation": {"labels": [{"text": "Rudder", "target": "the rudder"}]}},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    profile = {
        "master_index": 0,
        "layout_mapping": {
            "diagram": {"layout_index": 8, "layout_name": "Picture with Caption"},
        },
        "unmapped_fallback": {"layout_index": 8, "layout_name": "Picture with Caption"},
        "layouts": [_picture_profile_layout()],
    }

    output_path = build_deck(str(deck_dir), template_path, profile)
    prs = Presentation(output_path)
    slide = prs.slides[0]

    has_headline = any(
        getattr(s, "has_text_frame", False) and "Diagram headline" in (s.text_frame.text or "")
        for s in slide.shapes
    )
    has_body = any(
        getattr(s, "has_text_frame", False) and "diagram body point" in (s.text_frame.text or "")
        for s in slide.shapes
    )
    assert has_headline, "composed native diagram slide must keep its headline"
    assert has_body, "composed native diagram slide must keep its body_points"
    assert _count_named(slide, "annotation_label_") == 1


# --- E2E (skipif no template fixture) ---------------------------------------


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_composed_native_end_to_end(tmp_path):
    """Full build_deck run: a composed native slide keeps headline + body +
    figure + overlay; a sibling PLAIN composed slide keeps its own chrome
    unaffected (backward-compat)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()

    outline = {
        "narrative_arc": "test", "estimated_duration_minutes": 6, "total_slides": 2,
        "slides": [
            {"slide_number": 1, "slide_type": "content", "headline": "Plain composed sibling",
             "body_points": ["alpha", "beta"], "visual_type": "none", "layout_template": "content"},
            {"slide_number": 2, "slide_type": "content", "headline": "Composed annotated headline",
             "body_points": ["composed annotated body point"], "visual_type": "hero_image",
             "layout_template": "content"},
        ],
    }
    with open(deck_dir / "outline.json", "w") as f:
        json.dump(outline, f)

    image_manifest = {
        "generated_at": "2026-07-23T00:00:00Z", "image_backend": "ollama",
        "images": [
            {"image_id": "slide-02-hero", "slide_number": 2,
             "file_path": str(HERO_IMAGE_FIXTURE),
             "placement_zone": "annotated_image_zone",
             "annotations_path": "annotations/slide-02-annotations.json",
             "dimensions": {"width": 1920, "height": 1080},
             "source_prompt": "test", "model_used": "test", "alt_text": "hero",
             "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0},
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    payload = build_annotation_payload(
        slide_number=2, source="generated", base_image_path=str(HERO_IMAGE_FIXTURE),
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_image_zone",
        anchors={"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15]},
    )
    write_annotation_payload(str(deck_dir), 2, payload)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "plain",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "composed", "rationale": "composed annotated figure",
             "render_funnel": ["ollama", "cloud_low", "cloud_full"],
             "speaker_override": None, "annotation_mode": "native",
             "annotation": {"labels": [
                 {"text": "Rudder", "target": "the rudder"},
                 {"text": "Mainsail", "target": "the mainsail"},
             ]}},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    assert len(prs.slides) == 2

    plain_slide = prs.slides[0]
    has_plain_title = any(
        getattr(s, "has_text_frame", False)
        and "Plain composed sibling" in (s.text_frame.text or "")
        for s in plain_slide.shapes
    )
    assert has_plain_title, "sibling plain composed slide must keep its own chrome"

    composed_slide = prs.slides[1]
    has_headline = any(
        getattr(s, "has_text_frame", False)
        and "Composed annotated headline" in (s.text_frame.text or "")
        for s in composed_slide.shapes
    )
    has_body = any(
        getattr(s, "has_text_frame", False)
        and "composed annotated body point" in (s.text_frame.text or "")
        for s in composed_slide.shapes
    )
    assert has_headline, "composed native slide must retain its headline"
    assert has_body, "composed native slide must retain its body_points"

    pictures = [s for s in composed_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1

    label_texts = sorted(
        s.text_frame.text for s in composed_slide.shapes if s.name.startswith("annotation_label_")
    )
    assert label_texts == ["Mainsail", "Rudder"]


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_full_slide_native_default_unchanged(tmp_path):
    """Regression: a non-headline, non-composed native slide is byte-parity
    with v2 -- this simply re-runs the v2 end-to-end test to pin that the
    v2.1 composed/headline branching didn't perturb the default path."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_outline(deck_dir)

    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [
            {
                "image_id": "slide-02-hero", "slide_number": 2,
                "file_path": str(HERO_IMAGE_FIXTURE),
                "placement_zone": "annotated_full_slide",
                "annotations_path": "annotations/slide-02-annotations.json",
                "dimensions": {"width": 1920, "height": 1080},
                "source_prompt": "test", "model_used": "test", "alt_text": "hero",
                "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0,
            },
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    payload = build_annotation_payload(
        slide_number=2, source="generated", base_image_path=str(HERO_IMAGE_FIXTURE),
        image_dimensions={"width": 1920, "height": 1080},
        placement_zone="annotated_full_slide",
        anchors={"Rudder": [0.85, 0.7], "Mainsail": [0.4, 0.15]},
    )
    write_annotation_payload(str(deck_dir), 2, payload)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "keep chrome",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "annotated figure",
             "render_funnel": ["ollama", "cloud_low", "cloud_full"],
             "speaker_override": None, "annotation_mode": "native",
             "annotation": {"labels": [
                 {"text": "Rudder", "target": "the rudder"},
                 {"text": "Mainsail", "target": "the mainsail"},
             ]}},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    native_slide = prs.slides[1]

    pictures = [s for s in native_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    non_annotation_non_picture = [
        s for s in native_slide.shapes
        if s.shape_type != MSO_SHAPE_TYPE.PICTURE and not s.name.startswith("annotation_")
    ]
    assert len(pictures) == 1
    assert non_annotation_non_picture == []  # pure figure -- no headline band, no chrome
    assert _shapes_named(native_slide, "native_headline_") == []
