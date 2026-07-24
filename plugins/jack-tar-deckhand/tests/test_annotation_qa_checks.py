"""Tests for the annotate-figure v2 deck-qa checks (issue #142 v2, T7).

Design: docs/superpowers/plans/2026-07-17-annotate-figure-v2.md §7, §8.4.

Covers:
- AN-01 (check_annotation_contract): absent payload, base_image_hash
  mismatch, exact shape-count enforcement (wrong/extra counts, exact pass).
- AN-02 (check_label_text_verbatim): character-exact label text.
- AN-03 (check_labels_within_bounds): label shapes fully on-slide.
- run_qa integration: native-only routing (raster/none slides emit no AN
  findings), the annotation_* structural exemption (AP-02 skips label
  shapes but still flags a non-annotation textbox), and the strategy-map
  loader refactor that retains full per-slide entries (F3b).

All pptx fixtures are built programmatically via python-pptx — no image
file is ever Read into context (discipline hook, CLAUDE.md).
"""
from __future__ import annotations

import hashlib
import json
import sys
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu, Inches, Pt

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PLUGIN_ROOT))

from src.qa.checks.annotation_checks import (  # noqa: E402
    check_annotation_contract,
    check_label_text_verbatim,
    check_labels_within_bounds,
)
from src.qa.run_qa import _load_slide_strategy_entries, run_qa  # noqa: E402


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- fixture helpers ---------------------------------------------------


def _make_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return prs, slide


def _add_label(slide, name, text, left=Inches(1), top=Inches(1),
               width=Inches(1.5), height=Inches(0.4), font_pt=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.name = name
    tb.text_frame.text = text
    if font_pt is not None:
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(font_pt)
    return tb


def _add_connector(slide, name):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(2))
    conn.name = name
    return conn


def _add_oval(slide, name):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(0.1), Inches(0.1))
    shp.name = name
    return shp


def _valid_style(casing_width_pt=3.5):
    return {
        "leader_width_pt": 1.5,
        "casing_width_pt": casing_width_pt,
        "casing_color": "FFFFFF",
        "leader_color": "141414",
        "dot_radius_pt": 3.0,
        "box_fill": "FFFFFF",
        "box_border": "141414",
        "box_border_width_pt": 1.0,
        "text_color": "141414",
        "font_face": "Calibri",
        "font_size_pt": 18,
    }


def _make_payload(labels, base_image_hash="e" * 64, casing_width_pt=3.5):
    return {
        "slide_number": 1,
        "source": "generated",
        "base_image_path": "images/slide-01-base.png",
        "base_image_hash": base_image_hash,
        "image_dimensions": {"width": 1920, "height": 1080},
        "placement_zone": "annotated_full_slide",
        "fit": "contain",
        "labels": [
            {"text": text, "anchor": [0.5, 0.5], "label_pos": [0.6, 0.6]}
            for text in labels
        ],
        "style": _valid_style(casing_width_pt=casing_width_pt),
    }


def _build_full_shape_set(slide, n, casing=True):
    """Add exactly n of each annotation shape kind (label/leader/dot[/casing/dotring])."""
    for i in range(n):
        _add_label(slide, f"annotation_label_1_{i}", f"Label{i}")
        _add_connector(slide, f"annotation_leader_1_{i}")
        _add_oval(slide, f"annotation_dot_1_{i}")
        if casing:
            _add_connector(slide, f"annotation_casing_1_{i}")
            _add_oval(slide, f"annotation_dotring_1_{i}")


# --- AN-01: check_annotation_contract -----------------------------------


def test_an01_errors_when_payload_absent_for_native_slide():
    _, slide = _make_presentation()
    findings = check_annotation_contract(slide, 1, None)
    assert len(findings) == 1
    assert findings[0]["severity"] == "error"
    assert findings[0]["category"] == "annotation"
    assert "absent" in findings[0]["description"].lower()


def test_an01_errors_on_base_image_hash_mismatch(tmp_path):
    _, slide = _make_presentation()
    img_path = tmp_path / "base.png"
    Image.new("RGB", (10, 10), color=(1, 2, 3)).save(img_path)
    actual_hash = hashlib.sha256(img_path.read_bytes()).hexdigest()
    assert actual_hash != "e" * 64  # sanity: our fixture hash is deliberately wrong

    payload = _make_payload(["Rudder"], base_image_hash="e" * 64)
    image_entry = {"file_path": str(img_path)}

    findings = check_annotation_contract(slide, 1, payload, image_entry=image_entry)
    assert len(findings) == 1
    assert findings[0]["severity"] == "error"
    assert findings[0]["category"] == "annotation"
    assert "hash mismatch" in findings[0]["description"].lower()


def test_an01_flags_wrong_leader_count():
    _, slide = _make_presentation()
    payload = _make_payload(["Rudder", "Mainsail"], casing_width_pt=0)  # no casing required
    # 2 labels/dots but only 1 leader.
    _add_label(slide, "annotation_label_1_0", "Rudder")
    _add_label(slide, "annotation_label_1_1", "Mainsail")
    _add_connector(slide, "annotation_leader_1_0")
    _add_oval(slide, "annotation_dot_1_0")
    _add_oval(slide, "annotation_dot_1_1")

    findings = check_annotation_contract(slide, 1, payload)
    assert len(findings) == 1
    assert findings[0]["severity"] == "error"
    assert "annotation_leader_" in findings[0]["description"]
    assert "found 1" in findings[0]["description"]


def test_an01_flags_extra_label_shape():
    _, slide = _make_presentation()
    payload = _make_payload(["Rudder", "Mainsail"], casing_width_pt=0)
    _add_label(slide, "annotation_label_1_0", "Rudder")
    _add_label(slide, "annotation_label_1_1", "Mainsail")
    _add_label(slide, "annotation_label_1_2", "Extra")  # N+1
    _add_connector(slide, "annotation_leader_1_0")
    _add_connector(slide, "annotation_leader_1_1")
    _add_oval(slide, "annotation_dot_1_0")
    _add_oval(slide, "annotation_dot_1_1")

    findings = check_annotation_contract(slide, 1, payload)
    assert len(findings) == 1
    assert "annotation_label_" in findings[0]["description"]
    assert "found 3" in findings[0]["description"]


def test_an01_passes_when_counts_exact():
    _, slide = _make_presentation()
    payload = _make_payload(["Rudder", "Mainsail"], casing_width_pt=3.5)
    _build_full_shape_set(slide, 2, casing=True)

    findings = check_annotation_contract(slide, 1, payload)
    assert findings == []


# --- AN-02: check_label_text_verbatim ------------------------------------


def test_an02_flags_text_mismatch():
    _, slide = _make_presentation()
    payload = _make_payload(["Rudder"])
    _add_label(slide, "annotation_label_1_0", "Ruddar")  # typo

    findings = check_label_text_verbatim(slide, 1, payload)
    assert len(findings) == 1
    assert findings[0]["severity"] == "error"
    assert findings[0]["category"] == "annotation"
    assert "Rudder" in findings[0]["description"]
    assert "Ruddar" in findings[0]["description"]


def test_an02_passes_on_exact_match():
    _, slide = _make_presentation()
    payload = _make_payload(["Rudder", "Mainsail"])
    _add_label(slide, "annotation_label_1_0", "Rudder")
    _add_label(slide, "annotation_label_1_1", "Mainsail")

    findings = check_label_text_verbatim(slide, 1, payload)
    assert findings == []


# --- AN-03: check_labels_within_bounds ------------------------------------


def test_an03_flags_box_off_slide():
    prs, slide = _make_presentation()
    # Label pushed past the right edge of the slide.
    _add_label(slide, "annotation_label_1_0", "Rudder",
               left=SLIDE_W - Emu(10000), top=Inches(1), width=Inches(2), height=Inches(0.4))

    findings = check_labels_within_bounds(slide, 1, prs)
    assert len(findings) == 1
    assert findings[0]["severity"] == "warning"
    assert findings[0]["category"] == "annotation"


def test_an03_passes_when_all_in_bounds():
    prs, slide = _make_presentation()
    _add_label(slide, "annotation_label_1_0", "Rudder",
               left=Inches(1), top=Inches(1), width=Inches(1.5), height=Inches(0.4))

    findings = check_labels_within_bounds(slide, 1, prs)
    assert findings == []


# --- run_qa integration ---------------------------------------------------


def _write_strategy_map(deck_dir, slides):
    (deck_dir / "strategy-map.json").write_text(json.dumps({
        "approval_mode": "review",
        "slides": slides,
    }))


def test_annotation_checks_skipped_for_raster_and_none(tmp_path):
    """Non-native slides (raster / no annotation_mode key) emit no AN findings."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "full_bleed", "annotation_mode": "raster",
         "rationale": "x", "render_funnel": ["ollama"]},
        {"slide_number": 2, "strategy": "composed",
         "rationale": "x", "render_funnel": ["ollama"]},
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    annotation_findings = [f for f in report["findings"] if f["category"] == "annotation"]
    assert annotation_findings == []


def test_run_qa_native_branch_exempts_annotation_shapes_from_structural(tmp_path):
    """A sub-floor annotation label does NOT trip AP-02, while a sub-floor
    NON-annotation textbox on the same slide still does (F3a exemption)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "full_bleed", "annotation_mode": "native",
         "rationale": "x", "render_funnel": ["ollama"]},
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Deliberately sub-floor (10pt < 18pt AP-02 floor) annotation label — exempt.
    _add_label(slide, "annotation_label_1_0", "Rudder", font_pt=10)
    # Deliberately sub-floor NON-annotation textbox — should still be flagged.
    _add_label(slide, "plain_body_text", "Some body copy", font_pt=10)
    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    font_findings = [f for f in report["findings"] if f["category"] == "consistency"]
    affected = {f["affected_element"] for f in font_findings}

    assert "annotation_label_1_0" not in affected
    assert "plain_body_text" in affected


def test_run_qa_retains_full_strategy_entries(tmp_path):
    """The refactored loader exposes annotation_mode (and the rest of the
    entry) per slide, not just the resolved strategy string (F3b)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "full_bleed", "annotation_mode": "native",
         "rationale": "test rationale", "render_funnel": ["ollama"]},
    ])

    entries = _load_slide_strategy_entries(str(deck_dir))
    assert entries[1]["annotation_mode"] == "native"
    assert entries[1]["strategy"] == "full_bleed"
    assert entries[1]["rationale"] == "test rationale"


def test_load_slide_strategy_entries_empty_when_no_strategy_map(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    assert _load_slide_strategy_entries(str(deck_dir)) == {}


@pytest.mark.parametrize("mode", ["none", None])
def test_run_qa_ignores_annotation_mode_none_or_absent(tmp_path, mode):
    """annotation_mode 'none'/absent must not route to the native branch."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    entry = {"slide_number": 1, "strategy": "composed",
             "rationale": "x", "render_funnel": ["ollama"]}
    if mode is not None:
        entry["annotation_mode"] = mode
    _write_strategy_map(deck_dir, [entry])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    assert [f for f in report["findings"] if f["category"] == "annotation"] == []


# --- v2.1 T8: composed-strategy native slides (F-08, §5.4) -----------------


def test_composed_native_runs_an_checks(tmp_path):
    """AN-01/02/03 fire for a composed native slide exactly as for a
    full-slide one -- zone-agnostic, no per-zone variant needed."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "composed", "annotation_mode": "native",
         "rationale": "x", "render_funnel": ["ollama"]},
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.slides.add_slide(prs.slide_layouts[6])
    # No annotation shapes at all -- AN-01 fires (absent payload, since no
    # image-manifest/annotations file exists in deck_dir either).
    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    annotation_findings = [f for f in report["findings"] if f["category"] == "annotation"]
    assert len(annotation_findings) == 1
    assert annotation_findings[0]["severity"] == "error"


def test_composed_native_chrome_still_structurally_checked(tmp_path):
    """A composed native slide's REAL headline/body text (not
    annotation_-prefixed) is NOT exempted from structural checks -- a
    genuinely deficient body still flags (mirrors the full-slide exemption
    pin, but proves the exemption doesn't over-reach on composed chrome)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "composed", "annotation_mode": "native",
         "rationale": "x", "render_funnel": ["ollama"]},
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_label(slide, "annotation_label_1_0", "Rudder", font_pt=10)  # exempt
    _add_label(slide, "composed_body_text", "Some real body copy", font_pt=10)  # NOT exempt
    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    font_findings = [f for f in report["findings"] if f["category"] == "consistency"]
    affected = {f["affected_element"] for f in font_findings}

    assert "annotation_label_1_0" not in affected
    assert "composed_body_text" in affected


def test_composed_native_runs_visual_checks(tmp_path):
    """F-08: a composed-strategy native slide gets VISUAL_CHECKS findings
    (a planted low-contrast defect flags); a full-slide native slide does
    NOT run VISUAL_CHECKS at all (v2 parity)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_strategy_map(deck_dir, [
        {"slide_number": 1, "strategy": "composed", "annotation_mode": "native",
         "rationale": "x", "render_funnel": ["ollama"]},
        {"slide_number": 2, "strategy": "full_bleed", "annotation_mode": "native",
         "rationale": "x", "render_funnel": ["ollama"]},
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    composed_slide = prs.slides.add_slide(prs.slide_layouts[6])
    # White-on-white text -- a planted AP-07 low-contrast defect. Default
    # slide background is white, so this trips check_contrast.
    white_tb = composed_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    white_tb.name = "composed_body_text"
    white_tb.text_frame.text = "Invisible text"
    white_tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    full_slide_slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_white_tb = full_slide_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    full_white_tb.name = "annotation_label_2_0"
    full_white_tb.text_frame.text = "Invisible text"
    full_white_tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    report = run_qa(str(pptx_path), str(deck_dir))
    contrast_findings = [f for f in report["findings"] if f["category"] == "contrast"]
    composed_contrast = [f for f in contrast_findings if f["slide_number"] == 1]
    full_slide_contrast = [f for f in contrast_findings if f["slide_number"] == 2]

    assert composed_contrast, "composed-strategy native slide must run VISUAL_CHECKS"
    assert full_slide_contrast == [], "full-slide native slide must stay VISUAL_CHECKS-free (v2 parity)"
