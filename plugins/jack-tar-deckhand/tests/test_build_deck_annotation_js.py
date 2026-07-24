"""Tests for the annotate-figure v2 native-annotation builder in the
PptxGenJS assembler path (issue #142 v2, T6).

Design: docs/superpowers/plans/2026-07-17-annotate-figure-v2.md §3, §4,
§6.2, §8.3.

Covers (all via subprocess: `node build_deck.js --deck-dir ...` then OOXML
inspection with python-pptx — never pixels):

- Exact per-prefix shape counts (F10) + verbatim label text (the whole
  point of AN-02).
- `objectName` survival into the OOXML `name` attribute for every
  annotation shape kind (guards the Spike 1 `name`-drop regression).
- Pure figure (F2): native AND raster annotated slides drop headline,
  body_points, and the footer logo.
- `flipV` correctness (§4.1's bbox+flip leader-line rule) for 2+ anchor/
  label quadrant relationships.
- Contain-fit picture placement (§3.1/§3.2) for both `native` and `raster`
  modes, distinct from the existing cover-crop/stretch full-slide paths.
- The F5-adjacent JS-path contract: a missing or unparseable annotations
  payload never crashes the assembler and never silently draws a "pretty"
  slide — it logs a `console.warn` and skips the overlay.
- Backward compatibility: a slide with no `annotation_mode` key is
  untouched by the new routing branch (mirrors `test_full_bleed_scale.py`).

Do not `Read` PNG files directly — every image fixture here is either the
existing 16:9 hero fixture or a synthetic solid-colour PIL image built on
disk in `tmp_path`; only OOXML geometry and text are asserted.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
PROJECT_ROOT = PLUGIN_ROOT.parent.parent
sys.path.insert(0, str(PLUGIN_ROOT))

from src.process_image import compute_content_hash  # noqa: E402

PLUGIN_JS_ASSEMBLER = PLUGIN_ROOT / "src" / "assembler" / "build_deck.js"
HERO_IMAGE_FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "minimal_deck" / "images" / "slide-01-hero.png"

# 16:9 EMU-equivalent slide dims (matches build_deck.js's SLIDE_W/SLIDE_H
# defaults of 13.333in x 7.5in, in inches here for python-side maths).
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
EMU_PER_IN = 914400
EMU_TOLERANCE = 5000  # ~0.005in — comfortably above rounding noise

FORBIDDEN_HEADLINE = "This headline must not appear"
FORBIDDEN_BODY = "this body must not appear"


def _have_node_with_pptxgenjs() -> bool:
    """True only if `node` is on PATH and pptxgenjs resolves from the plugin dir."""
    if shutil.which("node") is None:
        return False
    check = subprocess.run(
        ["node", "-e", "require('pptxgenjs')"],
        cwd=str(PLUGIN_ROOT),
        capture_output=True,
        text=True,
    )
    return check.returncode == 0


pytestmark = pytest.mark.skipif(
    not _have_node_with_pptxgenjs(), reason="node + pptxgenjs not available"
)


# --- fixture helpers ---------------------------------------------------------


def _make_image(path: Path, size, color=(90, 120, 140)):
    Image.new("RGB", size, color).save(path)
    return str(path)


def _style_guide():
    return {
        "palette": {"primary": "1B3A4B", "accent": "C0392B", "text_primary": "111111",
                    "text_secondary": "555555", "background": "FFFFFF"},
        "typography": {"heading_font": "Inter", "body_font": "Inter",
                       "heading_sizes": {"slide_heading": 32}, "body_sizes": {"body": 18}},
        "layout": {"templates": {}},
        "image_style_tokens": {},
        "slide_palette": {},
    }


def _annotation_style():
    return {
        "leader_width_pt": 1.5, "casing_width_pt": 3.5, "casing_color": "FFFFFF",
        "leader_color": "141414", "dot_radius_pt": 3.0, "box_fill": "FFFFFF",
        "box_border": "141414", "box_border_width_pt": 1.0, "text_color": "141414",
        "font_face": "Calibri", "font_size_pt": 18,
    }


# Anchor/label_pos pairs deliberately chosen so the leader endpoints span
# BOTH flipV cases (§4.1: flipV = (ax<=bx) !== (ay<=by)) — verified against
# the real assembler output, not just predicted:
#   Alpha:    anchor and label both increase in x and y  -> flipV False
#   Beta One: anchor's x is greater but y is smaller      -> flipV True
#   Gamma-3:  anchor's x and y are both greater            -> flipV False
LABELS = [
    {"text": "Alpha", "anchor": [0.1, 0.1], "label_pos": [0.4, 0.4]},
    {"text": "Beta One", "anchor": [0.7, 0.2], "label_pos": [0.3, 0.6]},
    {"text": "Gamma-3", "anchor": [0.85, 0.85], "label_pos": [0.6, 0.55]},
]


def _write_common_contracts(deck_dir: Path, outline_slides):
    (deck_dir / "outline.json").write_text(json.dumps({
        "narrative_arc": "test", "estimated_duration_minutes": 4,
        "total_slides": len(outline_slides), "slides": outline_slides,
    }))
    (deck_dir / "style-guide.json").write_text(json.dumps(_style_guide()))
    (deck_dir / "chart-manifest.json").write_text(json.dumps({
        "generated_at": "2026-07-17T00:00:00Z", "charts": [],
        "summary": {"total_charts": 0, "rendered_count": 0, "failed_count": 0},
    }))
    (deck_dir / "speaker-notes.json").write_text(json.dumps({
        "generated_at": "2026-07-17T00:00:00Z", "talk_duration_minutes": 4, "notes": [],
    }))


def _write_annotations_payload(
        deck_dir: Path, slide_number: int, image_path: str,
        image_dims, labels, placement_zone="annotated_full_slide"):
    payload = {
        "slide_number": slide_number,
        "source": "generated",
        "base_image_path": image_path,
        "base_image_hash": compute_content_hash(image_path),
        "image_dimensions": {"width": image_dims[0], "height": image_dims[1]},
        "placement_zone": placement_zone,
        "fit": "contain",
        "labels": labels,
        "style": _annotation_style(),
    }
    ann_dir = deck_dir / "annotations"
    ann_dir.mkdir(exist_ok=True)
    rel_path = f"annotations/slide-{slide_number:02d}-annotations.json"
    (deck_dir / rel_path).write_text(json.dumps(payload))
    return rel_path


def _run_assembler(deck_dir: Path):
    (deck_dir / "output").mkdir(exist_ok=True)
    result = subprocess.run(
        ["node", str(PLUGIN_JS_ASSEMBLER), "--deck-dir", str(deck_dir)],
        capture_output=True, text=True, timeout=60,
    )
    assert result.returncode == 0, f"JS assembler failed: {result.stderr}"
    return deck_dir / "output" / "presentation.pptx", result.stderr


def _shapes_named(slide, prefix):
    return [s for s in slide.shapes if s.name.startswith(prefix)]


def _flipv(shape):
    xfrm = shape._element.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
    )
    if xfrm is None:
        return None
    val = xfrm.get("flipV")
    return val == "1"


def _build_deck_dir(
        deck_dir: Path, *, image_path: str, image_dims, annotation_mode,
        labels=None, annotations_path=None, placement_zone="annotated_full_slide",
        strategy="full_bleed", show_headline=False, include_manifest_dimensions=True):
    """Build the full contract set for a 2-slide deck: slide 1 stays plain
    `composed` (backward-compat control), slide 2 carries the annotation
    contract under test."""
    outline_slides = [
        {"slide_number": 1, "slide_type": "title", "headline": "Composed Title",
         "body_points": [], "visual_type": "hero_image", "layout_template": "title"},
        {"slide_number": 2, "slide_type": "content", "headline": FORBIDDEN_HEADLINE,
         "body_points": [FORBIDDEN_BODY], "visual_type": "hero_image", "layout_template": "content"},
    ]
    _write_common_contracts(deck_dir, outline_slides)

    images_dir = deck_dir / "images"
    images_dir.mkdir(exist_ok=True)
    dest = images_dir / Path(image_path).name
    if str(dest) != image_path:
        shutil.copy(image_path, dest)
    rel_image_path = f"./images/{dest.name}"

    image_entry = {
        "image_id": "slide-02-figure", "slide_number": 2, "file_path": rel_image_path,
        "placement_zone": placement_zone,
        "source_prompt": "test", "model_used": "test", "alt_text": "figure",
        "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0,
    }
    if include_manifest_dimensions:
        image_entry["dimensions"] = {"width": image_dims[0], "height": image_dims[1]}
    if annotations_path is not None:
        image_entry["annotations_path"] = annotations_path
    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [image_entry],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    (deck_dir / "image-manifest.json").write_text(json.dumps(image_manifest))

    annotation_block = {"labels": [
        {"text": lbl["text"], "target": lbl["text"]}
        for lbl in (labels or [{"text": "x"}])
    ]}
    if show_headline:
        annotation_block["show_headline"] = True

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "title",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": strategy, "rationale": "annotated figure",
             "render_funnel": ["ollama"], "speaker_override": None,
             "annotation_mode": annotation_mode,
             "annotation": annotation_block},
        ],
    }
    (deck_dir / "strategy-map.json").write_text(json.dumps(strategy_map))
    return dest


# --- exact shape counts + verbatim text (F10, AN-02's whole point) ---------


def test_js_assembler_native_annotation_emits_exact_shape_counts_and_verbatim_text(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    # Wire the annotations_path into the already-written manifest.
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    n = len(LABELS)
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_",
                   "annotation_casing_", "annotation_dotring_"):
        shapes = _shapes_named(slide, prefix)
        assert len(shapes) == n, f"{prefix}: expected {n}, got {len(shapes)}"

    label_texts = {
        s.text_frame.text for s in _shapes_named(slide, "annotation_label_")
        if getattr(s, "has_text_frame", False)
    }
    assert label_texts == {lbl["text"] for lbl in LABELS}


# --- objectName survival (Spike 1 regression guard) -------------------------


def test_js_assembler_native_annotation_objectname_survives(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_",
                   "annotation_casing_", "annotation_dotring_"):
        names = [s.name for s in slide.shapes if s.name.startswith(prefix)]
        assert len(names) > 0, (
            f"No shapes named with prefix {prefix!r} survived into the OOXML "
            "-- objectName may have been dropped (Spike 1 regression)."
        )


# --- pure figure (F2) --------------------------------------------------------


def test_js_assembler_native_annotation_drops_headline_and_body(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    for s in slide.shapes:
        if getattr(s, "has_text_frame", False):
            text = s.text_frame.text or ""
            assert FORBIDDEN_HEADLINE not in text
            assert FORBIDDEN_BODY not in text

    # No footer logo image beyond the base figure picture.
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_js_assembler_raster_annotation_drops_headline_and_body(tmp_path):
    """Pure figure applies to raster mode too (§6.2: "same builder... pure
    figure" with a no-op overlay) -- no headline/body even without a payload."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="raster", labels=None, annotations_path=None,
    )

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    for s in slide.shapes:
        if getattr(s, "has_text_frame", False):
            text = s.text_frame.text or ""
            assert FORBIDDEN_HEADLINE not in text
            assert FORBIDDEN_BODY not in text


# --- flipV correctness (§4.1) ------------------------------------------------


def test_js_assembler_native_annotation_flipv_matches_quadrant_geometry(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    leaders = {s.name: s for s in _shapes_named(slide, "annotation_leader_")}
    # Index 0 (Alpha): anchor and label both increase in x and y -> flipV False.
    alpha = next(s for name, s in leaders.items() if name.endswith("_0"))
    assert _flipv(alpha) in (None, False)
    # Index 1 (Beta One): anchor's x decreases toward the label while y
    # increases -> the two axes disagree -> flipV True.
    beta = next(s for name, s in leaders.items() if name.endswith("_1"))
    assert _flipv(beta) is True
    # Both flipV values are represented -- at least 2 quadrant cases covered.
    assert {_flipv(alpha), _flipv(beta)} != {_flipv(alpha)}


# --- contain-fit picture placement (§3.1/§3.2) ------------------------------


def test_js_assembler_native_annotation_contain_fits_off_aspect_image(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    one_label = [{"text": "Centre", "anchor": [0.5, 0.5], "label_pos": [0.2, 0.2]}]
    dest = _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="native", labels=one_label,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1000, 1000), one_label)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]

    expected_h = SLIDE_H_IN * EMU_PER_IN
    expected_w = expected_h  # square image, fit-to-height
    expected_x = (SLIDE_W_IN * EMU_PER_IN - expected_w) / 2

    assert abs(pic.height - expected_h) < EMU_TOLERANCE
    assert abs(pic.width - expected_w) < EMU_TOLERANCE
    assert abs(pic.left - expected_x) < EMU_TOLERANCE
    # Contain-fit, not cover: the picture must NOT fill the full slide width.
    assert pic.width < SLIDE_W_IN * EMU_PER_IN - EMU_TOLERANCE


def test_js_assembler_raster_annotation_contains_not_covers(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="raster", labels=None, annotations_path=None,
    )

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]

    expected_h = SLIDE_H_IN * EMU_PER_IN
    assert abs(pic.height - expected_h) < EMU_TOLERANCE
    assert pic.width < SLIDE_W_IN * EMU_PER_IN - EMU_TOLERANCE

    # Raster mode has no payload -- no overlay shapes at all.
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_",
                   "annotation_casing_", "annotation_dotring_"):
        assert _shapes_named(slide, prefix) == []


# --- missing / unparseable payload: warn-and-skip, never crash --------------


def test_js_assembler_native_annotation_missing_payload_warns_and_skips_overlay(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
        annotations_path="annotations/slide-02-annotations.json",  # never written
    )

    output_pptx, stderr = _run_assembler(deck_dir)
    assert "annotation" in stderr.lower() or "skip" in stderr.lower(), (
        f"Expected a console.warn about the missing annotation payload; stderr was: {stderr!r}"
    )

    prs = Presentation(output_pptx)
    slide = prs.slides[1]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1  # base image still placed, contain-fit
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_"):
        assert _shapes_named(slide, prefix) == []


def test_js_assembler_native_annotation_unparseable_payload_warns_and_skips_overlay(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS,
        annotations_path="annotations/slide-02-annotations.json",
    )
    ann_dir = deck_dir / "annotations"
    ann_dir.mkdir(exist_ok=True)
    (ann_dir / "slide-02-annotations.json").write_text("{ not valid json ")

    output_pptx, stderr = _run_assembler(deck_dir)
    assert "annotation" in stderr.lower() or "unparseable" in stderr.lower() or "skip" in stderr.lower(), (
        f"Expected a console.warn about the unparseable annotation payload; stderr was: {stderr!r}"
    )

    prs = Presentation(output_pptx)
    slide = prs.slides[1]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_"):
        assert _shapes_named(slide, prefix) == []


# --- backward compatibility --------------------------------------------------


def test_js_assembler_backward_compat_without_annotation_mode_unchanged(tmp_path):
    """A full_bleed slide with NO annotation_mode key routes through the
    pre-existing cover-fit builder untouched (mirrors test_full_bleed_scale.py)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    (deck_dir / "output").mkdir()

    outline_slides = [
        {"slide_number": 1, "slide_type": "title", "headline": "Composed Title",
         "body_points": [], "visual_type": "hero_image", "layout_template": "title"},
        {"slide_number": 2, "slide_type": "content", "headline": "Full bleed, no annotation",
         "body_points": ["body"], "visual_type": "hero_image", "layout_template": "content"},
    ]
    _write_common_contracts(deck_dir, outline_slides)

    images_dir = deck_dir / "images"
    images_dir.mkdir()
    shutil.copy(HERO_IMAGE_FIXTURE, images_dir / "slide-02-hero.png")
    image_manifest = {
        "generated_at": "2026-07-17T00:00:00Z", "image_backend": "ollama",
        "images": [
            {"image_id": "slide-02-hero", "slide_number": 2,
             "file_path": "./images/slide-02-hero.png",
             "placement_zone": "background", "dimensions": {"width": 1920, "height": 1080},
             "source_prompt": "test", "model_used": "test", "alt_text": "hero",
             "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0},
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    (deck_dir / "image-manifest.json").write_text(json.dumps(image_manifest))

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "title",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "infographic register",
             "render_funnel": ["ollama"], "speaker_override": None},
        ],
    }
    (deck_dir / "strategy-map.json").write_text(json.dumps(strategy_map))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    non_pictures = [s for s in slide.shapes if s.shape_type != MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert len(non_pictures) == 0  # full_bleed's zero-chrome contract, unchanged

    pic = pictures[0]
    # Cover-fit (existing behaviour): the picture fills the ENTIRE slide,
    # unlike the new contain-fit annotation path.
    assert pic.left == 0 and pic.top == 0
    assert abs(pic.width - SLIDE_W_IN * EMU_PER_IN) < EMU_TOLERANCE
    assert abs(pic.height - SLIDE_H_IN * EMU_PER_IN) < EMU_TOLERANCE


# --- v2.1 Feature A: composed annotated zone (§2.2, §5.3) --------------------


def test_js_composed_native_emits_chrome_and_overlay(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS, strategy="composed",
    )
    rel_ann_path = _write_annotations_payload(
        deck_dir, 2, str(dest), (1920, 1080), LABELS, placement_zone="annotated_image_zone",
    )
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    texts = [s.text_frame.text for s in slide.shapes if getattr(s, "has_text_frame", False)]
    assert any(FORBIDDEN_HEADLINE in t for t in texts), "composed chrome must retain the headline"
    assert any(FORBIDDEN_BODY in t for t in texts), "composed chrome must retain body_points"

    n = len(LABELS)
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_"):
        assert len(_shapes_named(slide, prefix)) == n


def test_js_composed_raster_no_overlay_shapes(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="raster", labels=None, annotations_path=None, strategy="composed",
    )

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    texts = [s.text_frame.text for s in slide.shapes if getattr(s, "has_text_frame", False)]
    assert any(FORBIDDEN_HEADLINE in t for t in texts)

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    for prefix in ("annotation_label_", "annotation_leader_", "annotation_dot_",
                   "annotation_casing_", "annotation_dotring_"):
        assert _shapes_named(slide, prefix) == []


def test_js_full_slide_ignores_stale_payload_zone(tmp_path):
    """F-02: a payload whose placement_zone says annotated_image_zone on a
    slide whose effective strategy is full_bleed must still contain-fit the
    FULL slide -- the strategy-derived routing zone wins, the payload zone
    is informational only."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    one_label = [{"text": "Centre", "anchor": [0.5, 0.5], "label_pos": [0.2, 0.2]}]
    dest = _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="native", labels=one_label, strategy="full_bleed",
    )
    # Payload claims the composed image-zone placement, but the slide's
    # strategy is full_bleed -- routing must ignore this.
    rel_ann_path = _write_annotations_payload(
        deck_dir, 2, str(dest), (1000, 1000), one_label, placement_zone="annotated_image_zone",
    )
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]

    # Square image, fit-to-height against the FULL slide (not the smaller
    # composed image-zone rect).
    expected_h = SLIDE_H_IN * EMU_PER_IN
    expected_w = expected_h
    assert abs(pic.height - expected_h) < EMU_TOLERANCE
    assert abs(pic.width - expected_w) < EMU_TOLERANCE


def test_js_composed_native_fit_uses_payload_dimensions(tmp_path):
    """F-01: a non-16:9 base image with a manifest entry lacking `dimensions`
    must still fit using payload.image_dimensions (not the 1024x576
    fallback), and the anchor mapping lands accordingly."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    # Portrait image so a 1024x576 fallback would produce a visibly
    # different fit than the real 400x800 aspect.
    portrait_path = _make_image(tmp_path / "portrait.png", (400, 800))
    one_label = [{"text": "Centre", "anchor": [0.5, 0.5], "label_pos": [0.2, 0.2]}]
    dest = _build_deck_dir(
        deck_dir, image_path=portrait_path, image_dims=(400, 800),
        annotation_mode="native", labels=one_label, strategy="composed",
        include_manifest_dimensions=False,
    )
    rel_ann_path = _write_annotations_payload(
        deck_dir, 2, str(dest), (400, 800), one_label, placement_zone="annotated_image_zone",
    )
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    assert "dimensions" not in manifest["images"][0]
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]

    # Default image zone (content_with_image, no style-guide override):
    zone_w = SLIDE_W_IN * 0.428
    zone_h = SLIDE_H_IN * 0.787
    zone_ratio = zone_w / zone_h
    img_ratio = 400 / 800
    assert img_ratio < zone_ratio  # image taller than the zone -> fit to height
    expected_h = zone_h * EMU_PER_IN
    expected_w = expected_h * img_ratio
    assert abs(pic.height - expected_h) < EMU_TOLERANCE
    assert abs(pic.width - expected_w) < EMU_TOLERANCE
    # NOT the 1024x576 fallback ratio, which would be wide, not tall.
    fallback_ratio = 1024 / 576
    assert abs(pic.width / pic.height - fallback_ratio) > 0.1


# --- v2.1 Feature B: headline opt-in (§3, §5.3) ------------------------------


def test_js_native_headline_band_textbox(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS, strategy="full_bleed", show_headline=True,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    headline_shapes = _shapes_named(slide, "native_headline_")
    assert len(headline_shapes) == 1
    assert FORBIDDEN_HEADLINE in headline_shapes[0].text_frame.text

    n = len(LABELS)
    assert len(_shapes_named(slide, "annotation_label_")) == n

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    band_h_emu = SLIDE_H_IN * 0.14 * EMU_PER_IN
    assert pictures[0].top >= band_h_emu - EMU_TOLERANCE


def test_js_native_headline_absent_by_default(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS, strategy="full_bleed", show_headline=False,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    assert _shapes_named(slide, "native_headline_") == []


def test_js_raster_show_headline_no_band(tmp_path):
    """F-04: the headline band is native-only -- raster + show_headline
    renders identically to raster without the flag."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    square_path = _make_image(tmp_path / "square.png", (1000, 1000))
    _build_deck_dir(
        deck_dir, image_path=square_path, image_dims=(1000, 1000),
        annotation_mode="raster", labels=None, annotations_path=None,
        strategy="full_bleed", show_headline=True,
    )

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    assert _shapes_named(slide, "native_headline_") == []


def test_js_headline_mode_fills_full_slide_background(tmp_path):
    """F-05: with show_headline true, a background-fill rect spans the full
    slide (no unfilled band seam above the shrunk zone)."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    dest = _build_deck_dir(
        deck_dir, image_path=str(HERO_IMAGE_FIXTURE), image_dims=(1920, 1080),
        annotation_mode="native", labels=LABELS, strategy="full_bleed", show_headline=True,
    )
    rel_ann_path = _write_annotations_payload(deck_dir, 2, str(dest), (1920, 1080), LABELS)
    manifest_path = deck_dir / "image-manifest.json"
    manifest = json.loads(manifest_path.read_text())
    manifest["images"][0]["annotations_path"] = rel_ann_path
    manifest_path.write_text(json.dumps(manifest))

    output_pptx, _stderr = _run_assembler(deck_dir)
    prs = Presentation(output_pptx)
    slide = prs.slides[1]

    from pptx.enum.shapes import MSO_SHAPE_TYPE as _SHAPE_TYPE
    rects = [
        s for s in slide.shapes
        if s.shape_type == _SHAPE_TYPE.AUTO_SHAPE
        and abs(s.width - SLIDE_W_IN * EMU_PER_IN) < EMU_TOLERANCE
        and abs(s.height - SLIDE_H_IN * EMU_PER_IN) < EMU_TOLERANCE
    ]
    assert len(rects) >= 1, "expected a full-slide background fill rect in showHeadline mode"
