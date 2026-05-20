"""Tests for the full_bleed strategy (issue #88 — v1.4.2).

The full_bleed strategy is "the image IS the slide" — picture fills the
canvas with ZERO chrome (no title, no body, no footer logo). Distinct
from full_render, which keeps a title overlay and footer logo.

Coverage:

- Schema accepts full_bleed in strategy and speaker_override enums.
- Neither auto-classifier (rule-based ``strategy_classifier`` or the
  outline-driven ``classify_slide_strategy``) can emit full_bleed. It
  is operator-opt-in only.
- ``build_strategy_map`` overrides + ``upgrade_slide_strategy`` propagate
  full_bleed to a cloud-rendering funnel.
- python-pptx ``_apply_full_bleed`` strips chrome, adds a picture at
  canvas size, hoists it to the first drawable position in spTree.
- Edge case: full_bleed without an image leaves an empty slide rather
  than misleading placeholder chrome.
- python-pptx ``build_deck`` end-to-end emits a full_bleed slide with
  one PICTURE shape and nothing else.
- Backward compatibility: existing strategies (full_render, composed)
  remain unchanged when the strategy map contains no full_bleed entries.
- JS assembler subprocess emits a full_bleed slide containing only the
  picture (no title text shapes), gated on Node + pptxgenjs being
  available so CI can skip when the JS toolchain is absent.
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import pytest
from jsonschema import ValidationError, validate
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
PROJECT_ROOT = PLUGIN_ROOT.parent.parent
sys.path.insert(0, str(PLUGIN_ROOT))

from src.assembler.build_deck_template import (  # noqa: E402
    _apply_full_bleed,
    build_deck,
)
from src.conductor import init_pipeline, upgrade_slide_strategy  # noqa: E402
from src.slide_prompt_composer import (  # noqa: E402
    build_strategy_map,
    classify_slide_strategy,
)
from src.strategy_classifier import classify_strategy  # noqa: E402

SCHEMA_PATH = PLUGIN_ROOT / "src" / "schemas" / "strategy_map.schema.json"
TEMPLATE_FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "templates" / "metamirror-template.pptx"
HERO_IMAGE_FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "minimal_deck" / "images" / "slide-01-hero.png"
PLUGIN_JS_ASSEMBLER = PLUGIN_ROOT / "src" / "assembler" / "build_deck.js"


# --- Schema -----------------------------------------------------------------


def _load_schema():
    with open(SCHEMA_PATH) as f:
        return json.load(f)


def _valid_strategy_map(strategy: str, speaker_override=None):
    entry = {
        "slide_number": 1,
        "strategy": strategy,
        "rationale": "test",
        "render_funnel": ["ollama", "cloud_low", "cloud_full"],
    }
    if speaker_override is not None:
        entry["speaker_override"] = speaker_override
    return {"approval_mode": "review", "slides": [entry]}


def test_schema_accepts_full_bleed_strategy():
    validate(instance=_valid_strategy_map("full_bleed"), schema=_load_schema())


def test_schema_accepts_full_bleed_speaker_override():
    validate(
        instance=_valid_strategy_map("composed", speaker_override="full_bleed"),
        schema=_load_schema(),
    )


def test_schema_rejects_unknown_strategy():
    with pytest.raises(ValidationError):
        validate(instance=_valid_strategy_map("not_a_strategy"), schema=_load_schema())


# --- Auto-classification safety --------------------------------------------


def test_rule_classifier_never_emits_full_bleed_for_academic_text():
    result = classify_strategy(
        "Figure 3: Transformer architecture [12]. Algorithm 2 (gradient descent)."
    )
    assert result.strategy != "full_bleed"


def test_rule_classifier_never_emits_full_bleed_for_business_text():
    result = classify_strategy("Quarterly revenue grew 12% in Q4.")
    assert result.strategy != "full_bleed"


@pytest.mark.parametrize(
    "slide",
    [
        {"slide_number": 1, "slide_type": "title", "headline": "X", "visual_type": "hero_image"},
        {"slide_number": 2, "slide_type": "content", "headline": "Y", "body_points": ["a"], "visual_type": "hero_image"},
        {"slide_number": 3, "slide_type": "content", "headline": "Z", "body_points": ["a", "b", "c", "d"], "visual_type": "hero_image"},
        {"slide_number": 4, "slide_type": "closing", "headline": "End", "visual_type": "none"},
        {"slide_number": 5, "slide_type": "data_chart", "headline": "Data", "visual_type": "chart"},
    ],
)
def test_outline_classifier_never_emits_full_bleed(slide):
    """full_bleed must be operator-opt-in only — never auto-classified."""
    result = classify_slide_strategy(slide, template_mode=False)
    assert result["strategy"] != "full_bleed"


def test_outline_classifier_in_template_mode_never_emits_full_bleed():
    full_bleed_layout = {
        "placeholders": [
            {"idx": 13, "type": "picture", "name": "Picture", "x": 0.0, "y": 0.0, "w": 13.33, "h": 7.5},
        ]
    }
    result = classify_slide_strategy(
        {"slide_number": 1, "slide_type": "title", "headline": "T", "visual_type": "hero_image"},
        template_mode=True,
        template_layout=full_bleed_layout,
    )
    assert result["strategy"] != "full_bleed"


# --- Override propagation ---------------------------------------------------


def test_build_strategy_map_override_to_full_bleed_assigns_cloud_funnel():
    outline = {
        "narrative_arc": "infographic-narrative",
        "estimated_duration_minutes": 10,
        "slides": [
            {"slide_number": 1, "slide_type": "content", "headline": "H",
             "body_points": ["a", "b"], "visual_type": "hero_image"},
        ],
    }
    strategy_map = build_strategy_map(outline, overrides={1: "full_bleed"})
    entry = strategy_map["slides"][0]
    assert entry["speaker_override"] == "full_bleed"
    assert entry["render_funnel"] == ["ollama", "cloud_low", "cloud_full"]


def test_upgrade_slide_strategy_to_full_bleed(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    init_pipeline(str(deck_dir))
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(
            {
                "approval_mode": "review",
                "slides": [
                    {"slide_number": 1, "strategy": "composed", "rationale": "x",
                     "render_funnel": ["ollama"], "speaker_override": None},
                ],
            },
            f,
        )
    upgrade_slide_strategy(str(deck_dir), 1, "full_bleed")
    with open(deck_dir / "strategy-map.json") as f:
        result = json.load(f)
    entry = result["slides"][0]
    assert entry["speaker_override"] == "full_bleed"
    assert entry["render_funnel"] == ["ollama", "cloud_low", "cloud_full"]


# --- _apply_full_bleed unit tests ------------------------------------------


def _make_template_slide_with_chrome(tmp_dir):
    """Build a minimal one-slide presentation with title + body + a rect."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Pre-existing title chrome"
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "Pre-existing body chrome"
    out = tmp_dir / "fixture.pptx"
    prs.save(out)
    return out, prs.slide_width, prs.slide_height


def test_apply_full_bleed_strips_all_shapes_and_adds_picture(tmp_path):
    fixture_pptx, slide_w, slide_h = _make_template_slide_with_chrome(tmp_path)
    prs = Presentation(fixture_pptx)
    slide = prs.slides[0]
    initial_shape_count = len(slide.shapes)
    assert initial_shape_count >= 2  # title + body textbox

    _apply_full_bleed(slide, str(HERO_IMAGE_FIXTURE), slide_w, slide_h)

    assert len(slide.shapes) == 1
    only = slide.shapes[0]
    assert only.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert only.left == 0 and only.top == 0
    assert only.width == slide_w and only.height == slide_h


def test_apply_full_bleed_without_picture_leaves_empty_slide(tmp_path):
    fixture_pptx, slide_w, slide_h = _make_template_slide_with_chrome(tmp_path)
    prs = Presentation(fixture_pptx)
    slide = prs.slides[0]

    _apply_full_bleed(slide, None, slide_w, slide_h)

    assert len(slide.shapes) == 0


def test_apply_full_bleed_missing_image_file_leaves_empty_slide(tmp_path):
    fixture_pptx, slide_w, slide_h = _make_template_slide_with_chrome(tmp_path)
    prs = Presentation(fixture_pptx)
    slide = prs.slides[0]

    _apply_full_bleed(slide, str(tmp_path / "does-not-exist.png"), slide_w, slide_h)

    assert len(slide.shapes) == 0


def test_apply_full_bleed_picture_is_first_drawable_in_sp_tree(tmp_path):
    fixture_pptx, slide_w, slide_h = _make_template_slide_with_chrome(tmp_path)
    prs = Presentation(fixture_pptx)
    slide = prs.slides[0]
    _apply_full_bleed(slide, str(HERO_IMAGE_FIXTURE), slide_w, slide_h)

    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    sp_tree = slide.shapes._spTree
    drawables = [child for child in sp_tree if child.tag not in (f"{ns}nvGrpSpPr", f"{ns}grpSpPr")]
    assert len(drawables) == 1
    pic_tag = f"{ns}pic"
    assert drawables[0].tag == pic_tag


# --- python-pptx build_deck integration ------------------------------------


def _write_minimal_template_inputs(deck_dir: Path, full_bleed_slide_nums):
    """Build the contracts build_deck() needs, with one optional full_bleed slide."""
    outline = {
        "narrative_arc": "infographic-narrative",
        "estimated_duration_minutes": 6,
        "total_slides": 2,
        "slides": [
            {
                "slide_number": 1,
                "slide_type": "content",
                "headline": "Composed slide retains chrome",
                "body_points": ["alpha", "beta"],
                "visual_type": "none",
                "layout_template": "content",
            },
            {
                "slide_number": 2,
                "slide_type": "content",
                "headline": "Full-bleed slide has no chrome",
                "body_points": ["this should not appear"],
                "visual_type": "hero_image",
                "layout_template": "content",
            },
        ],
    }
    with open(deck_dir / "outline.json", "w") as f:
        json.dump(outline, f)

    image_manifest = {
        "generated_at": "2026-05-20T00:00:00Z",
        "image_backend": "ollama",
        "images": [
            {
                "image_id": "slide-02-hero",
                "slide_number": 2,
                "file_path": str(HERO_IMAGE_FIXTURE),
                "placement_zone": "background",
                "dimensions": {"width": 1920, "height": 1080},
                "source_prompt": "test",
                "model_used": "test",
                "alt_text": "Full bleed hero",
                "status": "generated",
                "retry_count": 0,
                "generation_time_seconds": 1.0,
            },
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0, "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    if full_bleed_slide_nums:
        strategy_map = {
            "approval_mode": "review",
            "slides": [
                {"slide_number": 1, "strategy": "composed", "rationale": "keep chrome",
                 "render_funnel": ["ollama"], "speaker_override": None},
                {"slide_number": 2, "strategy": "full_bleed", "rationale": "infographic register",
                 "render_funnel": ["ollama", "cloud_low", "cloud_full"], "speaker_override": None},
            ],
        }
        with open(deck_dir / "strategy-map.json", "w") as f:
            json.dump(strategy_map, f)


def _stub_template_profile():
    """Minimal TemplateProfile pointing at the fixture template."""
    return {
        "master_index": 0,
        "layout_mapping": {
            "content": {"layout_index": 1, "layout_name": "Title and Content"},
            "title": {"layout_index": 0, "layout_name": "Title Slide"},
        },
        "unmapped_fallback": {"layout_index": 1, "layout_name": "Title and Content"},
        "layouts": [
            {
                "name": "Title and Content",
                "placeholders": [
                    {"idx": 0, "type": "title", "name": "Title", "x": 0.5, "y": 0.3, "w": 12.0, "h": 1.0},
                    {"idx": 1, "type": "content", "name": "Body", "x": 0.5, "y": 1.5, "w": 12.0, "h": 5.5},
                ],
            },
            {
                "name": "Title Slide",
                "placeholders": [
                    {"idx": 0, "type": "title", "name": "Title", "x": 0.5, "y": 3.0, "w": 12.0, "h": 1.5},
                ],
            },
        ],
    }


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_emits_full_bleed_slide_with_only_picture(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_minimal_template_inputs(deck_dir, full_bleed_slide_nums=[2])

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())

    prs = Presentation(output_path)
    assert len(prs.slides) == 2

    full_bleed_slide = prs.slides[1]
    pictures = [s for s in full_bleed_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    non_pictures = [s for s in full_bleed_slide.shapes if s.shape_type != MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert len(non_pictures) == 0
    assert pictures[0].left == 0 and pictures[0].top == 0
    assert pictures[0].width == prs.slide_width
    assert pictures[0].height == prs.slide_height


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_preserves_composed_slide_when_other_is_full_bleed(tmp_path):
    """Backward-compat: composed slide chrome is untouched even when a sibling is full_bleed."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_minimal_template_inputs(deck_dir, full_bleed_slide_nums=[2])

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)

    composed_slide = prs.slides[0]
    has_title_text = any(
        getattr(s, "has_text_frame", False) and "Composed slide retains chrome" in (s.text_frame.text or "")
        for s in composed_slide.shapes
    )
    assert has_title_text, "composed slide should keep its title text"


@pytest.mark.skipif(not TEMPLATE_FIXTURE.exists(), reason="template fixture missing")
def test_build_deck_template_without_strategy_map_unchanged(tmp_path):
    """Backward-compat: no strategy-map.json → no full_bleed behaviour kicks in."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    _write_minimal_template_inputs(deck_dir, full_bleed_slide_nums=None)

    output_path = build_deck(str(deck_dir), str(TEMPLATE_FIXTURE), _stub_template_profile())
    prs = Presentation(output_path)
    # Slide 2 should retain chrome since no strategy map said full_bleed.
    slide_two_shapes = list(prs.slides[1].shapes)
    assert len(slide_two_shapes) > 0


# --- JS assembler subprocess integration -----------------------------------


def _have_node_with_pptxgenjs() -> bool:
    """True only if `node` is on PATH and pptxgenjs resolves from the plugin dir."""
    if shutil.which("node") is None:
        return False
    check = subprocess.run(
        ["node", "-e", "require('pptxgenjs')"],
        cwd=str(PLUGIN_ROOT),
        capture_output=True,
        text=True,
    )
    return check.returncode == 0


def _write_js_assembler_inputs(deck_dir: Path):
    """Contracts the JS assembler reads, with a single full_bleed slide."""
    outline = {
        "narrative_arc": "infographic-narrative",
        "estimated_duration_minutes": 4,
        "total_slides": 2,
        "slides": [
            {"slide_number": 1, "slide_type": "title", "headline": "Composed Title",
             "body_points": [], "visual_type": "hero_image", "layout_template": "title"},
            {"slide_number": 2, "slide_type": "content", "headline": "This headline must not appear",
             "body_points": ["this body must not appear"], "visual_type": "hero_image",
             "layout_template": "content"},
        ],
    }
    with open(deck_dir / "outline.json", "w") as f:
        json.dump(outline, f)

    style_guide = {
        "palette": {"primary": "1B3A4B", "accent": "C0392B", "text_primary": "111111",
                    "text_secondary": "555555", "background": "FFFFFF"},
        "typography": {"heading_font": "Inter", "body_font": "Inter",
                       "heading_sizes": {"slide_heading": 32}, "body_sizes": {"body": 18}},
        "layouts": {},
        "image_style_tokens": {},
        "slide_palette": {},
    }
    with open(deck_dir / "style-guide.json", "w") as f:
        json.dump(style_guide, f)

    images_dir = deck_dir / "images"
    images_dir.mkdir()
    shutil.copy(HERO_IMAGE_FIXTURE, images_dir / "slide-02-hero.png")

    image_manifest = {
        "generated_at": "2026-05-20T00:00:00Z",
        "image_backend": "ollama",
        "images": [
            {"image_id": "slide-02-hero", "slide_number": 2,
             "file_path": "./images/slide-02-hero.png",
             "placement_zone": "background", "dimensions": {"width": 1920, "height": 1080},
             "source_prompt": "test", "model_used": "test", "alt_text": "hero",
             "status": "generated", "retry_count": 0, "generation_time_seconds": 1.0},
        ],
        "summary": {"total_images": 1, "generated_count": 1, "cached_count": 0,
                    "placeholder_count": 0, "failed_count": 0, "total_generation_seconds": 1.0},
    }
    with open(deck_dir / "image-manifest.json", "w") as f:
        json.dump(image_manifest, f)

    with open(deck_dir / "chart-manifest.json", "w") as f:
        json.dump({"generated_at": "2026-05-20T00:00:00Z", "charts": [],
                   "summary": {"total_charts": 0, "rendered_count": 0, "failed_count": 0}}, f)

    with open(deck_dir / "speaker-notes.json", "w") as f:
        json.dump({"generated_at": "2026-05-20T00:00:00Z", "talk_duration_minutes": 4,
                   "notes": []}, f)

    strategy_map = {
        "approval_mode": "review",
        "slides": [
            {"slide_number": 1, "strategy": "composed", "rationale": "title",
             "render_funnel": ["ollama"], "speaker_override": None},
            {"slide_number": 2, "strategy": "full_bleed", "rationale": "infographic register",
             "render_funnel": ["ollama", "cloud_low", "cloud_full"], "speaker_override": None},
        ],
    }
    with open(deck_dir / "strategy-map.json", "w") as f:
        json.dump(strategy_map, f)


@pytest.mark.skipif(not _have_node_with_pptxgenjs(), reason="node + pptxgenjs not available")
def test_js_assembler_emits_full_bleed_slide_with_no_title_text(tmp_path):
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    (deck_dir / "output").mkdir()
    _write_js_assembler_inputs(deck_dir)

    result = subprocess.run(
        ["node", str(PLUGIN_JS_ASSEMBLER), "--deck-dir", str(deck_dir)],
        capture_output=True, text=True, timeout=60,
    )
    assert result.returncode == 0, f"JS assembler failed: {result.stderr}"

    output_pptx = deck_dir / "output" / "presentation.pptx"
    assert output_pptx.is_file()

    prs = Presentation(output_pptx)
    assert len(prs.slides) == 2

    full_bleed_slide = prs.slides[1]
    text_shapes = [
        s for s in full_bleed_slide.shapes
        if getattr(s, "has_text_frame", False) and (s.text_frame.text or "").strip()
    ]
    forbidden = "This headline must not appear"
    forbidden_body = "this body must not appear"
    for s in text_shapes:
        assert forbidden not in s.text_frame.text
        assert forbidden_body not in s.text_frame.text

    pictures = [s for s in full_bleed_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
