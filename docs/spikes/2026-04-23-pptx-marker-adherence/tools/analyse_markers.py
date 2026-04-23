"""Parse a .pptx, find shapes whose names match the marker protocol.

Marker grammar: (IMAGE|SMARTART|BG):<identifier>
where identifier is one or more of [A-Za-z0-9_-].
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation

MARKER_RE = re.compile(r"^(IMAGE|SMARTART|BG):([A-Za-z0-9_-]+)$")


def analyse_pptx(path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    markers: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            name = getattr(shape, "name", "") or ""
            m = MARKER_RE.match(name)
            if not m:
                continue
            markers.append(
                {
                    "slide_index": slide_idx,
                    "shape_name": name,
                    "kind": m.group(1),
                    "identifier": m.group(2),
                    "left_emu": shape.left or 0,
                    "top_emu": shape.top or 0,
                    "width_emu": shape.width or 0,
                    "height_emu": shape.height or 0,
                }
            )
    by_kind: dict[str, int] = {}
    for mk in markers:
        by_kind[mk["kind"]] = by_kind.get(mk["kind"], 0) + 1
    return {
        "source": str(path),
        "markers": markers,
        "totals": {
            "markers": len(markers),
            "slides": len(prs.slides),
            "by_kind": by_kind,
        },
    }


if __name__ == "__main__":
    import json
    import sys

    if len(sys.argv) != 2:
        print("usage: analyse_markers.py <path.pptx>", file=sys.stderr)
        sys.exit(2)
    print(json.dumps(analyse_pptx(sys.argv[1]), indent=2))
