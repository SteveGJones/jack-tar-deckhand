"""Build a minimal .pptx with three named marker shapes for analyser tests."""
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
blank = prs.slide_layouts[6]

# Slide 1: IMAGE marker
slide1 = prs.slides.add_slide(blank)
img_shape = slide1.shapes.add_shape(1, Inches(1), Inches(1), Inches(3), Inches(2))
img_shape.name = "IMAGE:agent-architecture"
img_shape.text_frame.text = "IMAGE:agent-architecture"

# Slide 2: SMARTART marker
slide2 = prs.slides.add_slide(blank)
sa_shape = slide2.shapes.add_shape(1, Inches(1), Inches(1), Inches(5), Inches(3))
sa_shape.name = "SMARTART:flowchart"
sa_shape.text_frame.text = "SMARTART:flowchart"

# Slide 3: BG marker (small corner label)
slide3 = prs.slides.add_slide(blank)
bg_shape = slide3.shapes.add_shape(1, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.3))
bg_shape.name = "BG:dramatic-contrast"
bg_shape.text_frame.text = "BG:dramatic-contrast"

# Slide 4: non-marker shape (distractor — must not be counted)
slide4 = prs.slides.add_slide(blank)
other = slide4.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
other.name = "TitleBox"
other.text_frame.text = "Plain content, no marker"

prs.save("docs/spikes/2026-04-23-pptx-marker-adherence/tests/fixtures/minimal.pptx")
print("Wrote minimal.pptx with 3 markers + 1 distractor")
