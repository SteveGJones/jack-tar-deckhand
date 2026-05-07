"""Construct a non-PptxGenJS control .pptx using python-pptx directly.

Simulates a user handing the bridge a .pptx with NO accompanying build.js —
e.g. a deck exported from Keynote, a corporate template, or a deck the user
edited in PowerPoint. The JS parser cannot work here at all; the OOXML parser
must still return useful results.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent / "control.pptx"

prs = Presentation()
blank = prs.slide_layouts[6]

# Slide 1: title-like slide, no markers
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tb.text_frame.text = "Corporate All-Hands Q2 Update"

# Slide 2: bullets only, unmarked — candidate for SmartArt suggestion
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
tb.text_frame.text = "Strategic priorities:\n- Growth\n- Efficiency\n- Culture"

# Slide 3: deliberately named marker shape (as if a tech-savvy user pre-marked it)
s = prs.slides.add_slide(blank)
shp = s.shapes.add_shape(1, Inches(5), Inches(1), Inches(4), Inches(4))
shp.name = "IMAGE:vision-diagram"
shp.text_frame.text = "IMAGE:vision-diagram"

prs.save(OUT)
print(f"wrote {OUT}")
