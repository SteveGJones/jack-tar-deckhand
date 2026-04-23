from __future__ import annotations
import re
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from slide_facts import Marker, SlideFacts

MARKER_RE = re.compile(r"^(IMAGE|SMARTART|BG):([A-Za-z0-9_-]+)$")


def _shape_element_type(shape) -> str:
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    # text frames and plain shapes get disambiguated
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "text"
    return "shape"


def _background_kind(slide) -> str:
    cSld = slide.element.find(qn("p:cSld"))
    bg = cSld.find(qn("p:bg"))
    if bg is None:
        return "default"
    if bg.find(".//" + qn("a:blipFill")) is not None:
        return "image"
    if bg.find(".//" + qn("a:solidFill")) is not None:
        return "solid"
    return "unknown"


def parse_pptx(path: Union[str, Path]) -> list[SlideFacts]:
    prs = Presentation(str(path))
    results: list[SlideFacts] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        markers: list[Marker] = []
        counts: dict[str, int] = {"text": 0, "shape": 0, "image": 0, "chart": 0, "table": 0}
        for shape in slide.shapes:
            name = getattr(shape, "name", "") or ""
            m = MARKER_RE.match(name)
            if m:
                markers.append(Marker(
                    kind=m.group(1),
                    identifier=m.group(2),
                    left_emu=shape.left or 0,
                    top_emu=shape.top or 0,
                    width_emu=shape.width or 0,
                    height_emu=shape.height or 0,
                ))
                continue
            kind = _shape_element_type(shape)
            counts[kind] = counts.get(kind, 0) + 1
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        facts = SlideFacts(
            slide_index=idx,
            text_content="\n".join(t for t in text_parts if t.strip()),
            markers=markers,
            background_kind=_background_kind(slide),
            element_types=counts,
        )
        results.append(facts)
    return results


if __name__ == "__main__":
    import json
    import sys
    if len(sys.argv) != 2:
        print("usage: pptx_parser.py <path.pptx>", file=sys.stderr)
        sys.exit(2)
    facts = parse_pptx(sys.argv[1])
    print(json.dumps([f.to_dict() for f in facts], indent=2))
