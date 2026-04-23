"""Verify op3 grafts a pptx_native diagram onto the slide that had the SMARTART marker."""
import shutil
import sys
from pathlib import Path
from zipfile import ZipFile

SPIKE_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SPIKE_DIR / "prototypes"))

from op3_inject_smartart import inject_smartart

SEED = SPIKE_DIR / "seed" / "seed.pptx"

# Target chosen from Variant A seed inventory (slide 4, main-content-width three-pillars slide)
TARGET_MARKER = "SMARTART:three-pillars-relationship"
TARGET_SLIDE = 4


def test_inject_adds_diagram_parts(tmp_path):
    work = tmp_path / "op3.pptx"
    shutil.copy(SEED, work)

    spec = {
        "graphic_type": "flowchart",
        "layout_id": "process1",
        "data": {"items": ["Planning", "Memory", "Tool use"]},
    }

    inject_smartart(
        pptx_path=work,
        marker_name=TARGET_MARKER,
        slide_index_1based=TARGET_SLIDE,
        spec=spec,
        carriers_dir=tmp_path / "carriers",
    )

    with ZipFile(work) as z:
        names = z.namelist()
    diagram_parts = [n for n in names if "/diagrams/" in n and n.endswith(".xml")]
    assert len(diagram_parts) >= 3, (
        f"expected at least layout+data+colors+quickStyle diagram parts; got {diagram_parts}"
    )

    # The SMARTART marker shape must be gone (replaced by graphicFrame)
    from pptx import Presentation
    prs = Presentation(str(work))
    names_on_slide = [s.name for s in prs.slides[TARGET_SLIDE - 1].shapes]
    assert TARGET_MARKER not in names_on_slide


def test_file_still_openable_after_injection(tmp_path):
    work = tmp_path / "op3.pptx"
    shutil.copy(SEED, work)
    spec = {
        "graphic_type": "flowchart",
        "layout_id": "process1",
        "data": {"items": ["Planning", "Memory", "Tool use"]},
    }
    inject_smartart(
        pptx_path=work,
        marker_name=TARGET_MARKER,
        slide_index_1based=TARGET_SLIDE,
        spec=spec,
        carriers_dir=tmp_path / "carriers",
    )
    from pptx import Presentation
    prs = Presentation(str(work))
    assert len(prs.slides) == 10
