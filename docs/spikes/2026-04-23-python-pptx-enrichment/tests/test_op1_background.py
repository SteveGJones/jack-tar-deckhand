"""Verify op1 produces a deck where the target slide has a blipFill background pointing at an image part."""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

SPIKE_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SPIKE_DIR / "prototypes"))

from op1_background import apply_background

SEED = SPIKE_DIR / "seed" / "seed.pptx"
IMAGE = SPIKE_DIR / "seed" / "placeholder.png"

# Target chosen from Variant A seed inventory
TARGET_SLIDE = 1
TARGET_MARKER = "BG:title-hero-dark-grid"


def test_applies_blipfill_and_removes_marker(tmp_path):
    work = tmp_path / "op1.pptx"
    shutil.copy(SEED, work)

    apply_background(
        pptx_path=work,
        slide_index_1based=TARGET_SLIDE,
        image_path=IMAGE,
        marker_name=TARGET_MARKER,
    )

    prs = Presentation(str(work))
    slide = prs.slides[TARGET_SLIDE - 1]
    cSld = slide.element.find(qn("p:cSld"))
    bg = cSld.find(qn("p:bg"))
    assert bg is not None, f"slide {TARGET_SLIDE} should have a <p:bg> element after op1"
    blip = bg.find(".//" + qn("a:blip"))
    assert blip is not None, "<p:bg> should contain a <a:blip> image reference"
    rid = blip.get(qn("r:embed"))
    assert rid, "<a:blip> must have r:embed"
    part = slide.part.related_part(rid)
    assert part.content_type.startswith("image/"), (
        f"rId should resolve to an image part, got {part.content_type}"
    )

    marker_names = [s.name for s in slide.shapes]
    assert TARGET_MARKER not in marker_names, (
        "BG marker shape must be removed after background applied"
    )


def test_file_reopens_cleanly_in_python_pptx(tmp_path):
    work = tmp_path / "op1.pptx"
    shutil.copy(SEED, work)
    apply_background(
        pptx_path=work,
        slide_index_1based=TARGET_SLIDE,
        image_path=IMAGE,
        marker_name=TARGET_MARKER,
    )
    prs = Presentation(str(work))
    # seed has 10 slides — must be preserved
    assert len(prs.slides) == 10
