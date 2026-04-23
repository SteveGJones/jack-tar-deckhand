"""Verify op2 removes the IMAGE marker and replaces it with an embedded picture at the same position."""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SPIKE_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SPIKE_DIR / "prototypes"))

from op2_replace_image_shape import replace_image_marker

SEED = SPIKE_DIR / "seed" / "seed.pptx"
IMAGE = SPIKE_DIR / "seed" / "placeholder.png"

# Target chosen from Variant A seed inventory (slide 2, right-side illustration)
TARGET_MARKER = "IMAGE:demo-vs-prod-split"
TARGET_SLIDE_0BASED = 1  # slide 2 → index 1


def test_marker_replaced_with_picture_at_same_position(tmp_path):
    work = tmp_path / "op2.pptx"
    shutil.copy(SEED, work)

    # Capture target geometry before mutation
    before = Presentation(str(work))
    target = next(
        (s for s in before.slides[TARGET_SLIDE_0BASED].shapes if s.name == TARGET_MARKER),
        None,
    )
    assert target is not None, f"seed must have a {TARGET_MARKER} shape on slide {TARGET_SLIDE_0BASED + 1}"
    left, top, width, height = target.left, target.top, target.width, target.height

    replace_image_marker(pptx_path=work, marker_name=TARGET_MARKER, image_path=IMAGE)

    after = Presentation(str(work))
    slide = after.slides[TARGET_SLIDE_0BASED]
    # Marker must be gone
    names = [s.name for s in slide.shapes]
    assert TARGET_MARKER not in names
    # Exactly one new picture shape at the captured geometry
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) >= 1
    match = [
        p for p in pictures
        if p.left == left and p.top == top and p.width == width and p.height == height
    ]
    assert len(match) == 1, (
        f"expected one picture at {(left, top, width, height)}, "
        f"got {[(p.left, p.top, p.width, p.height) for p in pictures]}"
    )


def test_file_reopens_cleanly(tmp_path):
    work = tmp_path / "op2.pptx"
    shutil.copy(SEED, work)
    replace_image_marker(pptx_path=work, marker_name=TARGET_MARKER, image_path=IMAGE)
    prs = Presentation(str(work))
    assert len(prs.slides) == 10
