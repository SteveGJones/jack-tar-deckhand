"""Op1: apply a PNG as a slide background and remove the BG marker shape.

Uses python-pptx's part + relationship API for the image, and hand-constructs the
<p:bg> element via lxml because python-pptx's _Background facade exposes no
picture-fill setter.
"""
from __future__ import annotations

from pathlib import Path
from typing import Union

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.parts.image import Image, ImagePart


def _remove_existing_bg(cSld_elem) -> None:
    existing = cSld_elem.find(qn("p:bg"))
    if existing is not None:
        cSld_elem.remove(existing)


def _build_bg_element(rid: str):
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    xml = (
        f'<p:bg xmlns:p="{nsmap["p"]}" xmlns:a="{nsmap["a"]}" xmlns:r="{nsmap["r"]}">'
        "<p:bgPr>"
        '<a:blipFill dpi="0" rotWithShape="1">'
        f'<a:blip r:embed="{rid}"/>'
        "<a:srcRect/>"
        "<a:stretch><a:fillRect/></a:stretch>"
        "</a:blipFill>"
        "<a:effectLst/>"
        "</p:bgPr>"
        "</p:bg>"
    )
    return etree.fromstring(xml)


def apply_background(
    pptx_path: Union[str, Path],
    slide_index_1based: int,
    image_path: Union[str, Path],
    marker_name: str | None = None,
) -> None:
    pptx_path = Path(pptx_path)
    image_path = Path(image_path)

    prs = Presentation(str(pptx_path))
    if slide_index_1based < 1 or slide_index_1based > len(prs.slides):
        raise IndexError(
            f"slide {slide_index_1based} out of range (1..{len(prs.slides)})"
        )
    slide = prs.slides[slide_index_1based - 1]

    # Add the image as a part and relate it to the slide
    img = Image.from_file(str(image_path))
    image_part = ImagePart.new(slide.part.package, img)
    rid = slide.part.relate_to(image_part, RT.IMAGE)

    # Insert <p:bg> as the first child of <p:cSld>
    cSld = slide.element.find(qn("p:cSld"))
    _remove_existing_bg(cSld)
    bg = _build_bg_element(rid)
    cSld.insert(0, bg)

    # Remove marker shape if requested — matched by <p:cNvPr @name>
    if marker_name:
        spTree = cSld.find(qn("p:spTree"))
        for sp in list(spTree):
            nvSpPr = sp.find(qn("p:nvSpPr"))
            if nvSpPr is None:
                continue
            cNvPr = nvSpPr.find(qn("p:cNvPr"))
            if cNvPr is not None and cNvPr.get("name") == marker_name:
                spTree.remove(sp)

    prs.save(str(pptx_path))


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 5:
        print("usage: op1_background.py <pptx> <slide_1based> <image_png> <marker_name_or_->")
        sys.exit(2)
    mn = sys.argv[4] if sys.argv[4] != "-" else None
    apply_background(sys.argv[1], int(sys.argv[2]), sys.argv[3], mn)
    print("ok")
