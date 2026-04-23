"""Op2: find a shape by name on any slide, replace it with an embedded picture at the same position."""
from __future__ import annotations

from pathlib import Path
from typing import Union

from pptx import Presentation


def _find_shape_on_any_slide(prs, name: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == name:
                return slide, shape
    return None, None


def replace_image_marker(
    pptx_path: Union[str, Path],
    marker_name: str,
    image_path: Union[str, Path],
) -> None:
    pptx_path = Path(pptx_path)
    image_path = Path(image_path)

    prs = Presentation(str(pptx_path))
    slide, shape = _find_shape_on_any_slide(prs, marker_name)
    if shape is None:
        raise LookupError(f"shape {marker_name!r} not found on any slide")
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    # Remove the marker shape via its XML element
    sp_elem = shape._element
    sp_elem.getparent().remove(sp_elem)

    # Add a new picture at the captured geometry
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    prs.save(str(pptx_path))


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 4:
        print("usage: op2_replace_image_shape.py <pptx> <marker_name> <image_path>")
        sys.exit(2)
    replace_image_marker(sys.argv[1], sys.argv[2], sys.argv[3])
    print("ok")
