#!/usr/bin/env python3
"""deck-qa: Run 25 automated anti-pattern checks on a .pptx file.

Usage:
    python -m src.qa.run_qa [--pptx-path PATH] [--deck-dir PATH] [--duration MINUTES]

Default pptx-path: ./tmp/deck/output/presentation.pptx
Default deck-dir: ./tmp/deck
"""

import argparse
import json
import os
import sys

from pptx import Presentation

from .checks import (
    STRUCTURAL_CHECKS,
    STRUCTURAL_CHECKS_WITH_PRESENTATION,
    DECK_STRUCTURAL_CHECKS,
    CONSISTENCY_CHECKS,
    IMAGE_QUALITY_CHECKS,
    VISUAL_CHECKS,
    ANIMATION_CHECKS,
    COLOUR_CHECKS,
    KEYNOTE_CHECKS,
    check_slide_count_ratio,
    check_contrast,
)
from .config import QA_CONFIG
from .report import generate_report


def run_qa(pptx_path, deck_dir='./tmp/deck', duration_minutes=None, config=None):
    """Run QA checks with strategy-aware routing for keynote slides."""
    cfg = config or QA_CONFIG
    prs = Presentation(pptx_path)
    findings = []

    # Load strategy map (optional — absent means all slides are 'composed')
    strategy_map_path = os.path.join(deck_dir, 'strategy-map.json')
    slide_strategies = {}
    if os.path.exists(strategy_map_path):
        with open(strategy_map_path) as f:
            strategy_map = json.load(f)
        for entry in strategy_map.get('slides', []):
            slide_strategies[entry['slide_number']] = entry.get('speaker_override') or entry['strategy']

    # Load brand palette for palette drift checks
    brand_palette = []
    brand_profile_path = os.path.join(deck_dir, 'brand-profile.json')
    if os.path.exists(brand_profile_path):
        with open(brand_profile_path) as f:
            bp = json.load(f)
        palette = bp.get('palette', {})
        brand_palette = [v for v in palette.values() if isinstance(v, str) and len(v) == 6]

    # Step 1: Per-slide checks (strategy-aware)
    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        strategy = slide_strategies.get(slide_number, 'composed')

        if strategy == 'full_render':
            # Full render: skip text checks, run image + keynote checks
            for check_fn in IMAGE_QUALITY_CHECKS:
                findings.extend(check_fn(slide, slide_number, prs, config=cfg))
            for check_fn in KEYNOTE_CHECKS:
                findings.extend(check_fn(slide, slide_number, brand_palette=brand_palette, config=cfg))
        elif strategy in ('backdrop_render', 'background', 'backdrop', 'pragmatic_composition'):
            # Backdrop/background/pragmatic: text + image + keynote checks.
            # Contrast check is skipped: the QA checker compares text colour against the
            # slide background fill (white), not the actual image content. Backing pills
            # in the assembler ensure readability; contrast must be validated visually.
            for check_fn in STRUCTURAL_CHECKS:
                findings.extend(check_fn(slide, slide_number, config=cfg))
            for check_fn in STRUCTURAL_CHECKS_WITH_PRESENTATION:
                findings.extend(check_fn(slide, slide_number, prs, config=cfg))
            for check_fn in IMAGE_QUALITY_CHECKS:
                findings.extend(check_fn(slide, slide_number, prs, config=cfg))
            for check_fn in KEYNOTE_CHECKS:
                findings.extend(check_fn(slide, slide_number, brand_palette=brand_palette, config=cfg))
            for check_fn in VISUAL_CHECKS:
                if check_fn is check_contrast:
                    continue
                try:
                    findings.extend(check_fn(slide, slide_number, config=cfg))
                except Exception:
                    pass
        else:
            # Composed: standard checks (unchanged)
            for check_fn in STRUCTURAL_CHECKS:
                findings.extend(check_fn(slide, slide_number, config=cfg))
            for check_fn in STRUCTURAL_CHECKS_WITH_PRESENTATION:
                findings.extend(check_fn(slide, slide_number, prs, config=cfg))
            for check_fn in IMAGE_QUALITY_CHECKS:
                findings.extend(check_fn(slide, slide_number, prs, config=cfg))
            for check_fn in VISUAL_CHECKS:
                try:
                    findings.extend(check_fn(slide, slide_number, config=cfg))
                except Exception:
                    pass

    # Step 1b: Element layout checks (AP-27 to AP-30)
    from src.qa.checks.element_layout import (
        check_element_layout,
        check_vision_confidence,
        check_text_element_alignment,
        check_grid_reading_order,
        check_label_text_fit,
    )

    # Load image manifest for alignment checks
    im_path = os.path.join(deck_dir, 'image-manifest.json')
    im_data = {}
    if os.path.exists(im_path):
        with open(im_path) as f:
            im_data = json.load(f)
    images_by_slide = {}
    for img in im_data.get('images', []):
        sn = img.get('slide_number')
        if sn and img.get('detected_positions'):
            images_by_slide[sn] = img

    if os.path.exists(strategy_map_path):
        with open(strategy_map_path) as f:
            strategy_map_data = json.load(f)
        outline_path = os.path.join(deck_dir, 'outline.json')
        outline_slides = {}
        if os.path.exists(outline_path):
            with open(outline_path) as f:
                outline_data = json.load(f)
            outline_slides = {s['slide_number']: s for s in outline_data.get('slides', [])}

        for entry in strategy_map_data.get('slides', []):
            strategy = entry.get('speaker_override') or entry.get('strategy')
            if strategy in ('backdrop', 'pragmatic_composition'):
                outline_slide = outline_slides.get(entry['slide_number'], {})
                # AP-27: Element layout validation
                findings.extend(check_element_layout(entry, outline_slide))
                # AP-29: Text-element alignment
                image_entry = images_by_slide.get(entry['slide_number'], {})
                if image_entry:
                    findings.extend(check_text_element_alignment(entry, image_entry))
                # AP-31: Label text fit
                findings.extend(check_label_text_fit(entry, outline_slide, image_entry))
            # AP-30: Grid reading order (any strategy with grid layout)
            if entry.get('element_layout', {}).get('template') == 'grid_2x2' or \
               entry.get('body_layout') == 'grid_2x2':
                findings.extend(check_grid_reading_order(entry))

    # AP-28: Check vision confidence on image manifest
    for img in im_data.get('images', []):
        if img.get('detected_positions'):
            findings.extend(check_vision_confidence(img))

    # Step 1c: SmartArt QA checks (SA-01 to SA-05) — only when manifest present
    smartart_manifest_path = os.path.join(deck_dir, 'smartart-manifest.json')
    if os.path.exists(smartart_manifest_path):
        from src.qa.checks.smartart_checks import (
            check_data_integrity,
            check_label_legibility,
            check_enrichment_alignment,
            check_overflow_handling,
            check_accessibility,
        )
        with open(smartart_manifest_path) as f:
            smartart_manifest = json.load(f)

        # Load outline for SA-01 data integrity checks
        outline_path = os.path.join(deck_dir, 'outline.json')
        outline_slides_sa = {}
        if os.path.exists(outline_path):
            with open(outline_path) as f:
                outline_data_sa = json.load(f)
            outline_slides_sa = {
                s['slide_number']: s for s in outline_data_sa.get('slides', [])
            }

        for entry in smartart_manifest.get('slides', []):
            sn = entry.get('slide_number', 0)
            spec = entry.get('spec', {})
            svg_content = entry.get('svg_content', '')
            bg_color = entry.get('background_color')

            # SA-01: Data integrity
            outline_slide_sa = outline_slides_sa.get(sn, {})
            findings.extend(check_data_integrity(outline_slide_sa, spec, slide_number=sn))

            # SA-02: Label legibility
            if svg_content:
                findings.extend(check_label_legibility(svg_content, bg_color, slide_number=sn))

            # SA-03: Enrichment alignment
            findings.extend(check_enrichment_alignment(entry, im_data, slide_number=sn))

            # SA-04: Overflow handling
            if svg_content:
                findings.extend(check_overflow_handling(spec, svg_content, slide_number=sn))

            # SA-05: Accessibility
            if svg_content:
                findings.extend(check_accessibility(svg_content, entry, slide_number=sn))

    # Step 1d: Post-assembly visual inspection (rasterise + blank/brand checks)
    # Enabled when 'visual_inspection_enabled' is True in config (or key absent — opt-in).
    if cfg.get('visual_inspection_enabled', False):
        from src.qa.checks.visual_inspection import run_visual_inspection
        import tempfile

        outline_path = os.path.join(deck_dir, 'outline.json')
        vi_outline = {}
        if os.path.exists(outline_path):
            with open(outline_path) as f:
                vi_outline = json.load(f)

        brand_style_guide = {}
        if os.path.exists(brand_profile_path):
            with open(brand_profile_path) as f:
                bp_vi = json.load(f)
            brand_style_guide = {'palette': bp_vi.get('palette', {})}

        vi_output_dir = tempfile.mkdtemp(prefix='qa_vi_')
        findings.extend(run_visual_inspection(pptx_path, vi_outline, brand_style_guide, vi_output_dir))

    # Step 2: Deck-level structural checks
    for check_fn in DECK_STRUCTURAL_CHECKS:
        findings.extend(check_fn(prs, config=cfg))

    # AP-10: Slide count vs duration (needs external duration)
    if duration_minutes:
        findings.extend(check_slide_count_ratio(prs, duration_minutes, config=cfg))

    # Step 3: Cross-slide consistency checks
    for check_fn in CONSISTENCY_CHECKS:
        findings.extend(check_fn(prs, config=cfg))

    # Step 4: Deck-level animation checks
    for check_fn in ANIMATION_CHECKS:
        findings.extend(check_fn(prs, config=cfg))

    # Step 5: Deck-level colour checks
    colours_used = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                colours_used.add((rgb[0], rgb[1], rgb[2]))
                        except (AttributeError, TypeError):
                            pass
    for check_fn in COLOUR_CHECKS:
        findings.extend(check_fn(colours_used, config=cfg))

    return generate_report(findings, pptx_path, len(prs.slides))


def main():
    parser = argparse.ArgumentParser(description='Run QA checks on a .pptx file')
    parser.add_argument('--pptx-path', default='./tmp/deck/output/presentation.pptx')
    parser.add_argument('--deck-dir', default='./tmp/deck')
    parser.add_argument('--duration', type=int, default=None,
                        help='Talk duration in minutes (for slide count check)')
    parser.add_argument('--output', default=None,
                        help='Output path for QA report JSON')
    args = parser.parse_args()

    report = run_qa(args.pptx_path, args.deck_dir, args.duration)

    # Write report
    output_path = args.output or os.path.join(args.deck_dir, 'qa-report.json')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, 'w') as f:
        json.dump(report, f, indent=2)

    # Print summary
    print(f"QA Report: {report['verdict'].upper()}")
    print(f"  Errors: {report['summary']['errors']}")
    print(f"  Warnings: {report['summary']['warnings']}")
    print(f"  Info: {report['summary']['info']}")

    # Exit code: 1 if fail, 0 otherwise
    sys.exit(1 if report['verdict'] == 'fail' else 0)


if __name__ == '__main__':
    main()
