"""Template analyser — extract layouts and placeholders from a .pptx template.

Reads a corporate .pptx template via python-pptx and produces a TemplateProfile
JSON structure for template-driven assembly.
"""

import hashlib
import json
import os
import re

from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400


def _emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    return round(emu / EMU_PER_INCH, 2)


# Map OOXML placeholder type strings to our vocabulary
_TYPE_MAP = {
    'TITLE': 'title',
    'CENTER_TITLE': 'title',
    'SUBTITLE': 'subtitle',
    'BODY': 'body',
    'OBJECT': 'content',
    'PICTURE': 'picture',
    'CHART': 'chart',
    'TABLE': 'table',
    'DATE': 'date',
    'FOOTER': 'footer',
    'SLIDE_NUMBER': 'slide_number',
}

# Name patterns that override type-based classification
_NAME_OVERRIDES = [
    (re.compile(r'chapter', re.IGNORECASE), 'chapter_box'),
    (re.compile(r'logo', re.IGNORECASE), 'other'),
]


def classify_placeholder(ph_type_str, ph_name):
    """Classify a placeholder into our type vocabulary.

    Args:
        ph_type_str: String like 'TITLE (1)' or 'BODY (2)' from python-pptx.
        ph_name: The shape name string from the template.

    Returns:
        One of: title, subtitle, body, content, picture, chart, table,
                footer, slide_number, date, chapter_box, other.
    """
    # Check name-based overrides first
    for pattern, override_type in _NAME_OVERRIDES:
        if pattern.search(ph_name):
            return override_type

    # Extract the type keyword before the parenthesized number
    type_key = ph_type_str.split('(')[0].strip().split()[-1] if ph_type_str else ''
    return _TYPE_MAP.get(type_key, 'other')


def extract_layouts(template_path, master_index=0):
    """Extract all slide layouts from a template .pptx file.

    Args:
        template_path: Path to the .pptx template file.
        master_index: Which slide master to use (default 0).

    Returns:
        List of layout dicts with name, index, placeholders, and decorative shape count.
    """
    prs = Presentation(template_path)
    master = prs.slide_masters[master_index]
    results = []

    for i, layout in enumerate(master.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type_str = str(ph.placeholder_format.type) if ph.placeholder_format.type is not None else ''
            placeholders.append({
                'idx': ph.placeholder_format.idx,
                'type': classify_placeholder(ph_type_str, ph.name),
                'name': ph.name,
                'x': _emu_to_inches(ph.left),
                'y': _emu_to_inches(ph.top),
                'w': _emu_to_inches(ph.width),
                'h': _emu_to_inches(ph.height),
            })

        non_placeholder_shapes = [s for s in layout.shapes if not s.is_placeholder]

        results.append({
            'name': layout.name,
            'index': i,
            'placeholder_count': len(placeholders),
            'placeholders': placeholders,
            'decorative_shape_count': len(non_placeholder_shapes),
        })

    return results
