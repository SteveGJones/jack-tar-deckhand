"""Template-driven deck assembly via python-pptx.

Opens a corporate .pptx template, strips example slides, adds new slides
from the outline using the template's mapped layouts, and populates
placeholders with content.
"""

import json
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


def _resolve_layout(prs, profile, slide_type):
    """Find the python-pptx SlideLayout object for a given slide type."""
    mapping = profile.get('layout_mapping', {})
    entry = mapping.get(slide_type) or profile.get('unmapped_fallback')
    if not entry:
        return prs.slide_layouts[0]

    master_index = profile.get('master_index', 0)
    master = prs.slide_masters[master_index]
    layout_index = entry['layout_index']
    return master.slide_layouts[layout_index]


def _strip_existing_slides(prs):
    """Remove all existing slides from the presentation."""
    slide_list = prs.slides._sldIdLst
    for sldId in list(slide_list):
        slide_list.remove(sldId)


def _find_placeholder_by_type(slide, ph_type, profile_layout):
    """Find a slide placeholder matching the given type from profile layout metadata."""
    if not profile_layout:
        return None

    target_indices = set()
    for ph_info in profile_layout.get('placeholders', []):
        if ph_info['type'] == ph_type:
            target_indices.add(ph_info['idx'])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx in target_indices:
            return ph
    return None


def _populate_text(placeholder, text):
    """Set text on a placeholder, preserving template formatting."""
    if placeholder is None or not text:
        return
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text


def _populate_body_points(placeholder, body_points):
    """Set bulleted body points on a content or body placeholder."""
    if placeholder is None or not body_points:
        return
    tf = placeholder.text_frame
    tf.clear()
    for i, point in enumerate(body_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point


def _get_profile_layout(profile, layout_name):
    """Find a layout entry in the profile by name."""
    for layout in profile.get('layouts', []):
        if layout['name'] == layout_name:
            return layout
    return None


def _get_mapped_layout_name(profile, slide_type):
    """Get the layout name for a slide type from the mapping."""
    mapping = profile.get('layout_mapping', {})
    entry = mapping.get(slide_type) or profile.get('unmapped_fallback')
    if entry:
        return entry['layout_name']
    return None


def _emit_smartart_placeholder(slide, slide_number, profile_layout):
    """Add a named rectangle placeholder for pptx_native SmartArt injection.

    Uses the content placeholder bounds from the profile layout.
    The existing assembler_patch.py finds this by name and replaces it.
    """
    content_info = None
    if profile_layout:
        for ph_info in profile_layout.get('placeholders', []):
            if ph_info['type'] == 'content':
                content_info = ph_info
                break

    if not content_info:
        x, y, w, h = Inches(0.6), Inches(2.3), Inches(12.13), Inches(4.57)
    else:
        x = Inches(content_info['x'])
        y = Inches(content_info['y'])
        w = Inches(content_info['w'])
        h = Inches(content_info['h'])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.name = f'pptx_native_placeholder_{slide_number}'
    shape.fill.background()
    shape.line.fill.background()


def build_deck(deck_dir, template_path, template_profile):
    """Assemble a .pptx deck using a template's slide layouts.

    Args:
        deck_dir: Path to the DeckContext directory.
        template_path: Path to the source .pptx template.
        template_profile: dict conforming to TemplateProfile schema.

    Returns:
        str: Path to the output .pptx file.
    """
    with open(os.path.join(deck_dir, 'outline.json')) as f:
        outline = json.load(f)

    with open(os.path.join(deck_dir, 'image-manifest.json')) as f:
        image_manifest = json.load(f)

    notes_path = os.path.join(deck_dir, 'speaker-notes.json')
    speaker_notes = {}
    if os.path.isfile(notes_path):
        with open(notes_path) as f:
            notes_data = json.load(f)
        speaker_notes = {n['slide_number']: n['text'] for n in notes_data.get('notes', [])}

    image_lookup = {}
    for img in image_manifest.get('images', []):
        if img.get('status') in ('generated', 'accepted', 'accepted_with_issues'):
            sn = img['slide_number']
            if sn not in image_lookup:
                image_lookup[sn] = img['file_path']

    # Build a set of slide numbers that need pptx_native SmartArt placeholders
    pptx_native_slides = set()
    sa_manifest_path = os.path.join(deck_dir, 'smartart-manifest.json')
    if os.path.isfile(sa_manifest_path):
        with open(sa_manifest_path) as f:
            sa_manifest = json.load(f)
        for graphic in sa_manifest.get('graphics', []):
            if graphic.get('engine_used') == 'pptx_native':
                pptx_native_slides.add(graphic['slide_number'])

    prs = Presentation(template_path)
    _strip_existing_slides(prs)

    for slide_data in outline.get('slides', []):
        slide_number = slide_data['slide_number']
        slide_type = slide_data.get('slide_type', 'content')
        headline = slide_data.get('headline', '')
        body_points = slide_data.get('body_points', [])

        layout = _resolve_layout(prs, template_profile, slide_type)
        slide = prs.slides.add_slide(layout)

        layout_name = _get_mapped_layout_name(template_profile, slide_type)
        profile_layout = _get_profile_layout(template_profile, layout_name)

        # Populate title
        title_ph = _find_placeholder_by_type(slide, 'title', profile_layout)
        _populate_text(title_ph, headline)

        # Populate body/content
        content_ph = _find_placeholder_by_type(slide, 'content', profile_layout)
        body_ph = _find_placeholder_by_type(slide, 'body', profile_layout)
        if content_ph and body_points:
            _populate_body_points(content_ph, body_points)
        elif body_ph and body_points:
            _populate_body_points(body_ph, body_points)

        # SmartArt placeholder for pptx_native injection
        if slide_number in pptx_native_slides:
            _emit_smartart_placeholder(slide, slide_number, profile_layout)

        # Populate picture placeholder
        if slide_number in image_lookup:
            pic_ph = _find_placeholder_by_type(slide, 'picture', profile_layout)
            if pic_ph:
                image_path = image_lookup[slide_number]
                abs_image_path = os.path.join(deck_dir, image_path) if not os.path.isabs(image_path) else image_path
                if os.path.isfile(abs_image_path):
                    pic_ph.insert_picture(abs_image_path)

        # Speaker notes
        if slide_number in speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes[slide_number]

    output_dir = os.path.join(deck_dir, 'output')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'presentation.pptx')
    prs.save(output_path)

    return output_path
